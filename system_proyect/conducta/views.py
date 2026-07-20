import io, os, json
from accounts.models import CoordPermiso
from django.shortcuts import render, redirect, get_object_or_404
from django.http import JsonResponse, HttpResponse, HttpResponseForbidden
from django.contrib.auth.decorators import login_required
from django.views.decorators.http import require_POST
from django.contrib.auth.models import User
from core.utils_notifications import crear_notificacion
from django.template.loader import render_to_string
from django.contrib import messages
from django.db import connections
from django.db.models import Count, Q
from django.utils import timezone
from django.conf import settings
from django.core.mail import send_mail, EmailMultiAlternatives
import datetime as _dt_module

# Bandeja compartida por área (usada como FROM y TO en notificaciones de email)
# Los correos individuales de coordinadores ya NO se usan como destinatarios de email;
# toda notificación cae a la bandeja compartida de cada área.
SITE_URL = 'https://servicios.ana-hn.org:437'

_EMAIL_COORDI_BL = 'coordinacion_bl@ana-hn.org'

_C1_GRADOS   = ['1ero', '2do', '3ero']
_C2_GRADOS   = ['4to', '5to', '6to', '7mo', '8vo', '9no']
_C3_MATERIAS = [
    'español', 'estudios sociales', 'eess', 'cívica', 'civica',
    'computacion', 'computación', 'penmanship', 'arte',
    'educacion fisica', 'educación física', 'biblia',
    'inteligencia emocional',
]

def _q_grados(grados):
    q = Q()
    for g in grados:
        q |= Q(grado__icontains=g)
    return q

def _q_materias(materias):
    q = Q()
    for m in materias:
        q |= Q(materia__icontains=m)
    return q

def _q_docentes(docentes):
    q = Q()
    for d in docentes:
        q |= Q(docente__icontains=d)
    return q

def _q_for_coord(codigo):
    """Q filter for a coordinator's dashboard: explicit override takes priority over docente match."""
    q_docente = _q_docentes_for_coord(codigo)
    no_override = Q(coord_asignado='') | Q(coord_asignado__isnull=True)
    return Q(coord_asignado=codigo) | (no_override & q_docente)


# Reporte ACADÉMICO → coordinador por GRADO (C1 = grados 1-3, C2 = 4-9).
# Un docente puede dar clases en grados de C1 y C2 (ej. Miss Hernandez), por eso el
# filtro por docente mezclaba; se filtra por grado, respetando el override coord_asignado.
# El campo `grado` se guarda como 'PrimariaBL 2do-_1', 'ColegioBL 7mo-_1', etc.
_GRADO_TOKEN = {1: '1ero', 2: '2do', 3: '3ero', 4: '4to', 5: '5to',
                6: '6to', 7: '7mo', 8: '8vo', 9: '9no'}
_GRADOS_COORD = {'C1': [1, 2, 3], 'C2': [4, 5, 6, 7, 8, 9]}


# ─────────────────────────────────────────────────────────────
# <--- hecho por claude code: ruteo BILINGÜE por mapeo en archivo JSON
# (materia manda → grado → override por maestro). Solo área bilingüe.
# ─────────────────────────────────────────────────────────────
_ROUTING_BL_PATH = os.path.join(settings.MEDIA_ROOT, 'conducta', 'routing_bl.json')

# <--- hecho por claude code: nombres de coordinadores para mostrar en los portales
_COORD_NOMBRES = {'C1': 'Mrs Varela', 'C2': 'Mr Ruiz', 'C3': 'Miss Alcerro', 'C4': 'Mr Martinez'}

# <--- hecho por claude code: correo de cada coordinador (para permisos por coordinador en Ruteo)
_COORD_EMAIL = {'C1': 'cvarela@ana-hn.org', 'C2': 'druiz@ana-hn.org',
                'C3': 'ialcerro@ana-hn.org', 'C4': 'jmartinez@ana-hn.org'}
_EMAIL_COORD = {v: k for k, v in _COORD_EMAIL.items()}
# Compat: set de correos que antes tenían solo-lectura (C1, C2, C4). Se usa para migrar.
_COORD_VER_RUTEO_EMAILS = {'cvarela@ana-hn.org', 'druiz@ana-hn.org', 'jmartinez@ana-hn.org'}
_PERMISOS_RUTEO = ('none', 'lectura', 'edit')   # niveles válidos por coordinador

_ROUTING_BL_DEFAULT = {
    'actualizado': '',
    'alumnos': [],
    'grupos': {},   # <--- hecho por claude code: {grado_str: {nombre, nivel, grado_num, seccion, coordinador, clases[]}}
    'grados_c1': [1, 2, 3],
    'grados_c2': [4, 5, 6, 7, 8, 9],
    'materias_c4': ['matematicas'],
    'materias_c3': ['espanol', 'estudios sociales', 'eess', 'civica', 'computacion',
                    'arte', 'penmanship', 'educacion fisica', 'biblia', 'inteligencia emocional'],
    'maestros_override': {},
    'coord_ver_ruteo': False,   # <--- hecho por claude code: (legado) visibilidad global; migrado a coord_permisos
    'coord_permisos': {},       # <--- hecho por claude code: {C1|C2|C3|C4: 'none'|'lectura'|'edit'} — permiso por coordinador
    'docentes_catalogo': [],    # <--- hecho por claude code: lista [{docente, materias, coord}] — un docente puede tener varias entradas
}


def _coord_permisos(cfg):
    """<--- hecho por claude code: permisos por coordinador, con default 'none'.
    Migra el toggle global viejo (coord_ver_ruteo=True → C1/C2/C4 en 'lectura')."""
    base = {c: 'none' for c in _COORD_EMAIL}
    guardado = cfg.get('coord_permisos') or {}
    if guardado:
        for c in base:
            v = str(guardado.get(c, 'none')).lower()
            base[c] = v if v in _PERMISOS_RUTEO else 'none'
    elif cfg.get('coord_ver_ruteo'):   # legado: C1/C2/C4 tenían solo lectura
        for correo in _COORD_VER_RUTEO_EMAILS:
            base[_EMAIL_COORD[correo]] = 'lectura'
    return base


def _permiso_ruteo_email(email, cfg=None):
    """Nivel de permiso ('none'|'lectura'|'edit') del correo dado en Ruteo."""
    cfg = cfg or _cargar_routing_bl()
    cod = _EMAIL_COORD.get((email or '').lower())
    if not cod:
        return 'none'
    return _coord_permisos(cfg).get(cod, 'none')


def _acceso_ruteo(request):
    """Devuelve (puede_ver, solo_lectura) para la pantalla de Ruteo.
    Superusuario → acceso total. Cada coordinador según su permiso: edit → ver+editar, lectura → solo ver."""
    u = request.user
    if not u.is_authenticated:
        return False, True
    if u.is_superuser:
        return True, False
    permiso = _permiso_ruteo_email(u.email)
    if permiso == 'edit':
        return True, False
    if permiso == 'lectura':
        return True, True
    return False, True


def _puede_editar_ruteo(request):
    """<--- hecho por claude code: True si el usuario puede GUARDAR cambios en Ruteo (superuser o coord con permiso 'edit')."""
    u = getattr(request, 'user', None)
    if not (u and u.is_authenticated):
        return False
    if u.is_superuser:
        return True
    return _permiso_ruteo_email(u.email) == 'edit'


def _default_coord_grado(n):
    """Coordinador por defecto de un grupo según su grado (1-3 → C1, 4-9 → C2)."""
    if n in (1, 2, 3):
        return 'C1'
    if n in (4, 5, 6, 7, 8, 9):
        return 'C2'
    return ''


def _nivel_abbr(nivel):
    """Abreviatura del nivel para distinguir grupos con el mismo número (Preescolar vs Primaria)."""
    return {'PreescolarBL': 'Pre', 'PrimariaBL': 'Pri', 'ColegioBL': 'Col'}.get(nivel or '', nivel or '')


def _nombre_grupo(grado_str, grado_num, seccion):
    """Nombre corto de un grupo: '1.1', '7.2'… con fallback al texto crudo.
    (El nivel se antepone al mostrar, ver _nombre_grupo_disp, para no confundir Pre 1.1 con Pri 1.1)."""
    if grado_num and seccion != '':
        return f"{grado_num}.{seccion}"
    if grado_num:
        return str(grado_num)
    return grado_str


def _nombre_grupo_disp(nivel, nombre, seccion=''):
    """Nombre para mostrar. Preescolar → 'Prepa 1', 'Prepa 2' (evita confundir con Primaria).
    Primaria/Colegio → su número tal cual ('1.1', '7.1', '9.1')."""
    if nivel == 'PreescolarBL':
        return f"Prepa {seccion}".strip()
    return nombre


def _reconstruir_grupos(cfg):
    """Arma cfg['grupos'] a partir de cfg['alumnos'] (uno por grado-sección),
    conservando coordinador y clases ya configurados."""
    prev = cfg.get('grupos') or {}
    nuevos = {}
    for a in cfg.get('alumnos', []):
        gr = a.get('grado', '')
        if not gr:
            continue
        if gr not in nuevos:
            existente = prev.get(gr, {})
            n = a.get('grado_num')
            sec = a.get('seccion', '')
            nivel = gr.split(' ')[0] if ' ' in gr else ''
            nuevos[gr] = {
                'nombre': _nombre_grupo(gr, n, sec),
                'nivel': nivel,
                'grado_num': n,
                'seccion': sec,
                'coordinador': existente.get('coordinador') or _default_coord_grado(n),
                'clases': existente.get('clases', []),
            }
    cfg['grupos'] = nuevos
    return cfg


def _norm_txt(s):
    """Minúsculas + sin tildes, para comparar materias sin importar acentos/mayúsculas."""
    s = (s or '').strip().lower()
    for a, b in (('á', 'a'), ('é', 'e'), ('í', 'i'), ('ó', 'o'), ('ú', 'u'), ('ñ', 'n')):
        s = s.replace(a, b)
    return s


# <--- hecho por claude code: al escribir se separan varias materias con ';' (o ','); se aceptan ambos
def _split_materias(s):
    return [x.strip() for x in (s or '').replace(';', ',').split(',') if x.strip()]


def _join_materias(s):
    """Para mostrar en la ventana del maestro: siempre separadas por coma."""
    return ', '.join(_split_materias(s))


def _catalogo_entries(cfg):
    """<--- hecho por claude code: normaliza docentes_catalogo a lista [{docente, materias, coord}].
    Un mismo docente puede tener VARIAS entradas. Migra el formato viejo dict {docente: materias}."""
    raw = cfg.get('docentes_catalogo')
    if isinstance(raw, list):
        out = []
        for e in raw:
            doc = str((e or {}).get('docente', '')).strip()
            if not doc:
                continue
            coord = str((e or {}).get('coord', '')).strip().upper()
            out.append({'docente': doc,
                        'materias': _join_materias((e or {}).get('materias', '')),
                        'coord': coord if coord in {'C1', 'C2', 'C3', 'C4'} else ''})
        return out
    # formato viejo: dict {docente: materias} + maestros_override
    overs = cfg.get('maestros_override', {}) or {}
    out = []
    for doc, mats in (raw or {}).items():
        doc = str(doc).strip()
        if not doc:
            continue
        cs = _norm_override_list(overs.get(doc))
        if cs:
            for c in cs:
                out.append({'docente': doc, 'materias': _join_materias(mats), 'coord': c})
        else:
            out.append({'docente': doc, 'materias': _join_materias(mats), 'coord': ''})
    return out


def _cargar_routing_bl():
    """Lee el JSON de ruteo bilingüe; si no existe o está corrupto, devuelve el default."""
    try:
        with open(_ROUTING_BL_PATH, encoding='utf-8') as fh:
            data = json.load(fh)
        cfg = dict(_ROUTING_BL_DEFAULT)
        cfg.update(data or {})
        return cfg
    except (FileNotFoundError, ValueError, OSError):
        return dict(_ROUTING_BL_DEFAULT)


def _guardar_routing_bl(data):
    """Escribe el JSON de ruteo bilingüe (crea la carpeta si falta)."""
    os.makedirs(os.path.dirname(_ROUTING_BL_PATH), exist_ok=True)
    with open(_ROUTING_BL_PATH, 'w', encoding='utf-8') as fh:
        json.dump(data, fh, ensure_ascii=False, indent=2)


def _norm_override_list(v):
    """Normaliza un override de maestro a lista de códigos válidos (acepta str o lista)."""
    if not v:
        return []
    vals = [v] if isinstance(v, str) else (list(v) if isinstance(v, (list, tuple)) else [])
    out = []
    for x in vals:
        c = str(x).strip().upper()
        if c in {'C1', 'C2', 'C3', 'C4'} and c not in out:
            out.append(c)
    return out


def _coord_bl(materia, docente, grado, cfg=None):
    """Coordinador destino de un reporte BILINGÜE.
    Regla: override único (fuerza) → materia manda (C3/C4) → grupo del alumno → grado.
    Si el maestro tiene VARIOS overrides (ej. reporta a C1 y C2), se elige por el grado/grupo."""
    cfg = cfg or _cargar_routing_bl()
    overrides = cfg.get('maestros_override', {}) or {}
    ov_list = _norm_override_list(overrides.get(docente)) if docente else []
    # 1) Override ÚNICO → fuerza (sin importar materia ni grado)
    if len(ov_list) == 1:
        return ov_list[0]
    # 2) Materia manda (C4 = Matemáticas, C3 = sus materias) — revisa cada materia si hay varias
    m_list = [_norm_txt(x) for x in _split_materias(materia)] or [_norm_txt(materia)]
    c4 = [_norm_txt(x) for x in cfg.get('materias_c4', [])]
    c3 = [_norm_txt(x) for x in cfg.get('materias_c3', [])]
    if any(x in c4 for x in m_list):
        return 'C4'
    if any(x in c3 for x in m_list):
        return 'C3'
    # 3) Grupo del alumno → coordinador del grupo  [nueva fuente]
    grupo = (cfg.get('grupos') or {}).get(grado)
    grupo_coord = grupo.get('coordinador') if grupo else ''
    # 4) Grado → C1/C2 (fallback)
    n, _sec = _parse_grado_num(grado)
    grado_coord = 'C1' if n in cfg.get('grados_c1', [1, 2, 3]) else (
        'C2' if n in cfg.get('grados_c2', [4, 5, 6, 7, 8, 9]) else '')
    destino = grupo_coord or grado_coord
    # 5) Override MÚLTIPLE (ej. C1 y C2): restringe al set permitido, elige por grado/grupo
    if ov_list:
        if destino in ov_list:
            return destino
        if grado_coord in ov_list:
            return grado_coord
        return ov_list[0]
    return destino


def _q_grado_academico(codigo):
    """Q por GRADO del reporte. Grados 1-3 solo Primaria (excluye Preescolar, que también usa 1ero/2do/3ero)."""
    q = Q(pk__in=[])
    for n in _GRADOS_COORD.get(codigo, []):
        cond = Q(grado__icontains=_GRADO_TOKEN[n] + '-')
        if n <= 3:
            cond &= Q(grado__istartswith='PrimariaBL')
        q |= cond
    return q


def _q_academico_for_coord(codigo):
    if codigo not in _GRADOS_COORD:       # C3/C4 u otros: comportamiento original
        return _q_for_coord(codigo)
    no_override = Q(coord_asignado='') | Q(coord_asignado__isnull=True)
    # Grados 1-9 → por grado. Preescolar (u otros niveles) → por docente, para no perder reportes.
    q_pre = Q(grado__istartswith='PreescolarBL') & _q_docentes_for_coord(codigo)
    return Q(coord_asignado=codigo) | (no_override & (_q_grado_academico(codigo) | q_pre))

def _q_docentes_for_coord(codigo):
    """
    Builds Q() for reports belonging to coordinator `codigo`.
    For teachers shared across coordinators, restricts by the materias
    assigned to this coordinator in MateriaDocenteBilingue so that
    reports from shared teachers only appear in the right dashboard.
    """
    all_records = list(MateriaDocenteBilingue.objects.filter(activo=True))
    docente_all_coords = {}
    for md in all_records:
        for c in [x.strip() for x in md.coordinador.split(',')]:
            docente_all_coords.setdefault(md.docente, set()).add(c)
    q = Q()
    for md in all_records:
        if codigo not in [x.strip() for x in md.coordinador.split(',')]:
            continue
        other_coords = docente_all_coords[md.docente] - {codigo}
        if other_coords:
            q_mat = Q()
            for m in [x.strip() for x in md.materia.split(',') if x.strip()]:
                q_mat |= Q(materia__icontains=m)
            q |= Q(docente__icontains=md.docente) & q_mat
        else:
            q |= Q(docente__icontains=md.docente)
    return q

_TIPO_COLOR = {
    'Reporte Informativo': '#1971c2',
    'Reporte Conductual':  '#f59f00',
    'Progress Report':     '#0ca678',
}
_TIPO_ICON = {
    'Reporte Informativo': '&#8505;',
    'Reporte Conductual':  '&#9888;',
    'Progress Report':     '&#128200;',
}


def _notificar_coordinadores(tipo, maestro, alumno, grado, materia, area, coordinador_bl=None, md_colegio=None, subtipo=None):
    nombre = maestro.get_full_name() or maestro.username
    asunto = f"[TechCare] Nuevo {tipo} registrado – {alumno}"
    color = _TIPO_COLOR.get(tipo, '#1971c2')
    icon  = _TIPO_ICON.get(tipo, '&#128203;')
    anio  = _dt_module.datetime.now().year
    filas_html = ""
    for etiqueta, valor in [("Tipo", tipo), ("Maestro", nombre), ("Alumno", alumno),
                             ("Grado", grado), ("Materia", materia), ("Área", area)]:
        if valor:
            filas_html += (
                f'<tr><td style="padding:8px 0;font-size:12px;color:#6c757d;'
                f'text-transform:uppercase;letter-spacing:.4px;width:90px;">{etiqueta}</td>'
                f'<td style="padding:8px 0;font-size:14px;color:#1a1a2e;font-weight:600;">{valor}</td></tr>'
            )
    html_cuerpo = f"""<!DOCTYPE html>
<html lang="es"><head><meta charset="UTF-8"></head>
<body style="margin:0;padding:0;background:#f0f4f8;font-family:Arial,sans-serif;">
<table width="100%" cellpadding="0" cellspacing="0" style="background:#f0f4f8;padding:32px 0;">
  <tr><td align="center">
    <table width="560" cellpadding="0" cellspacing="0" style="background:#ffffff;border-radius:12px;overflow:hidden;box-shadow:0 4px 24px rgba(0,0,0,.10);">
      <tr>
        <td style="background:linear-gradient(135deg,#1864ab,{color});padding:28px 40px;text-align:center;">
          <p style="margin:0;font-size:24px;font-weight:700;color:#fff;letter-spacing:-.5px;">TechCare</p>
          <p style="margin:4px 0 0;font-size:12px;color:rgba(255,255,255,.75);">Asociación Nuevo Amanecer</p>
        </td>
      </tr>
      <tr>
        <td style="padding:32px 40px 20px;">
          <p style="margin:0 0 4px;font-size:18px;font-weight:700;color:#1a1a2e;">
            {icon} Nuevo reporte registrado
          </p>
          <p style="margin:0 0 24px;font-size:14px;color:#6c757d;">Se ha registrado un nuevo reporte en el sistema.</p>
          <table width="100%" cellpadding="0" cellspacing="0"
                 style="background:#f8f9fa;border-radius:8px;padding:16px 20px;margin-bottom:24px;border-left:4px solid {color};">
            {filas_html}
          </table>
          <table width="100%" cellpadding="0" cellspacing="0">
            <tr><td align="center">
              <a href="{SITE_URL}"
                 style="display:inline-block;background:{color};color:#fff;text-decoration:none;
                        padding:12px 32px;border-radius:8px;font-size:14px;font-weight:600;">
                Ver en el sistema
              </a>
            </td></tr>
          </table>
        </td>
      </tr>
      <tr>
        <td style="background:#f8f9fa;padding:14px 40px;border-top:1px solid #e9ecef;text-align:center;">
          <p style="margin:0;font-size:11px;color:#adb5bd;">
            © {anio} Soporte Técnico – Asociación Nuevo Amanecer
          </p>
        </td>
      </tr>
    </table>
  </td></tr>
</table>
</body></html>"""
    texto_plano = (
        f"Nuevo {tipo} registrado.\n\n"
        f"Maestro : {nombre}\nAlumno  : {alumno}\nGrado   : {grado}\n"
        f"Materia : {materia}\nÁrea    : {area}\n\nRevisa: {SITE_URL}"
    )

    # ── EMAIL: solo a usuarios con la flag activa en DestinatarioEmail ──
    # conducta_bl / conducta_col según área; progress_bl para progress
    _CAMPO_EMAIL = {
        ('progress',              'bilingue'): 'progress_bl',
        ('progress',              'colegio'):  'progress_bl',  # no hay progress_col
        ('conductual',            'bilingue'): 'conducta_bl',
        ('conductual',            'colegio'):  'conducta_col',
        ('informativo_academico', 'bilingue'): 'conducta_bl',
        ('informativo_academico', 'colegio'):  'conducta_col',
        ('informativo_conductual','bilingue'): 'conducta_bl',
        ('informativo_conductual','colegio'):  'conducta_col',
    }
    campo_email = _CAMPO_EMAIL.get((subtipo, area))
    if campo_email:
        try:
            from accounts.models import DestinatarioEmail
            emails = list(
                DestinatarioEmail.objects
                .filter(**{campo_email: True})
                .exclude(user__email='')
                .values_list('user__email', flat=True)
            )
            if emails:
                msg = EmailMultiAlternatives(
                    asunto, texto_plano,
                    settings.DEFAULT_FROM_EMAIL,
                    emails,
                )
                msg.attach_alternative(html_cuerpo, "text/html")
                msg.send(fail_silently=True)
        except Exception as _e:
            print(f"[WARN] email conducta: {_e}")

    # ── CAMPANITA + TOAST ──
    # Progress → DestinatarioEmail.progress_bl (todos los marcados)
    # Resto   → ConfiguracionNotificacion por subtipo (controla quién recibe cada tipo)
    bell_users = []

    if subtipo == 'progress':
        from accounts.models import DestinatarioEmail
        bell_users = [
            d.user for d in
            DestinatarioEmail.objects.filter(progress_bl=True).select_related('user')
            if d.user
        ]
    else:
        _SUBTIPO_CAMPO = {
            'conductual':             'recibe_conductual',
            'informativo_academico':  'recibe_informativo_academico',
            'informativo_conductual': 'recibe_informativo_conductual',
        }
        if subtipo and subtipo in _SUBTIPO_CAMPO:
            campo = _SUBTIPO_CAMPO[subtipo]
            reglas = ConfiguracionNotificacion.objects.filter(
                area=area, activo=True, **{campo: True}
            ).select_related('coordinador__usuario')
            bell_users = [r.coordinador.usuario for r in reglas if r.coordinador.usuario]

    if not bell_users:
        # Fallback: todos los coordinadores activos del área
        _area_cfg = 'bilingue' if area != 'colegio' else 'colegio'
        cfgs = ConfiguracionCoordinador.objects.filter(
            area=_area_cfg, activo=True
        ).select_related('usuario')
        bell_users = [c.usuario for c in cfgs if c.usuario]

    mensaje_noti = f"Nuevo {tipo} — {alumno} ({grado})"
    if materia:
        mensaje_noti += f" · {materia}"
    tipo_noti = 'alerta' if tipo == 'Reporte Conductual' else 'info'
    for usuario in bell_users:
        try:
            crear_notificacion(
                usuario,
                mensaje_noti,
                modulo='conducta',
                tipo=tipo_noti,
                enviar_correo=False,
            )
        except Exception:
            pass

# PDF y PowerPoint
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from reportlab.lib.units import mm
from reportlab.lib import colors
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Pt

from .forms import (
    ReporteInformativoForm,
    ReporteConductualForm,
    ProgressReportForm
)
from .models import (
    IncisoConductual,
    ReporteInformativo,
    ReporteConductual,
    ProgressReport,
    MateriaDocenteBilingue,
    MateriaDocenteColegio,
    EvidenciaReporte,  # <--- hecho por claude code: importar modelo de evidencias para la view subir_evidencia
    ConfiguracionNotificacion,
    ConfiguracionCoordinador,
)
from tickets.models import Ticket


# ---------------------------
# FUNCIONES AUXILIARES
# ---------------------------
def obtener_alumnos_bilingue():
    query = """
    SELECT d.PersonaID, 
           ISNULL(d.Nombre1,'') + 
           CASE WHEN d.Nombre2 IS NOT NULL AND d.Nombre2 <> '' THEN ' ' + d.Nombre2 ELSE '' END + ' ' +
           ISNULL(d.Apellido1,'') + 
           CASE WHEN d.Apellido2 IS NOT NULL AND d.Apellido2 <> '' THEN ' ' + d.Apellido2 ELSE '' END AS NombreCompl,
           da.Descripcion AS AreaDesc, c.CrsoNumero, c.GrupoNumero
      FROM dbo.tblPrsDtosGen AS d
      JOIN dbo.tblPrsTipo           AS t  ON d.PersonaID = t.PersonaID
      JOIN dbo.tblEdcArea           AS a  ON t.IngrEgrID  = a.IngrEgrID
      JOIN dbo.tblEdcEjecCrso       AS ec ON a.AreaID     = ec.AreaID
      JOIN dbo.tblEdcCrso           AS c  ON ec.CrsoID    = c.CrsoID
      JOIN dbo.tblEdcDescripAreaEdc AS da ON a.DescrAreaEdcID = da.DescrAreaEdcID
     WHERE d.Alum = 1
     AND DATEPART(yy, c.FechaInicio) = DATEPART(yyyy, GETDATE())
       AND da.Descripcion IN (N'PrimariaBL', N'ColegioBL')   -- <--- hecho por claude code: Preescolar no manda reportes
       AND ec.Desertor = 0 AND ec.TrasladoPer = 0
     ORDER BY da.Descripcion, c.CrsoNumero, c.GrupoNumero, NombreCompl
    """
    alumnos = []
    try:
        with connections['padres_sqlserver'].cursor() as cursor:
            cursor.execute(query)
            for pid, nombre, area, crso, grupo in cursor.fetchall():
                label = nombre.strip()
                grado = f"{area} {crso}-{grupo}"
                alumnos.append({'id': str(pid), 'label': label, 'grado': grado})
    except Exception as e:
        print(">>> ERROR SQL BILINGUE:", e)
    return alumnos


# <--- hecho por claude code: alumnado BL desde el JSON de config (no BD en vivo).
# Fallback a SQL si el JSON aún no tiene alumnos (antes del primer "Refrescar").
def obtener_alumnos_bilingue_cfg():
    cfg = _cargar_routing_bl()
    alumnos = cfg.get('alumnos') or []
    if not alumnos:
        return obtener_alumnos_bilingue()
    return [{'id': str(a.get('id')),
             'label': a.get('label') or a.get('nombre') or '',
             'grado': a.get('grado', '')} for a in alumnos]


# <--- hecho por claude code: alumnado BL agrupado por grupo (grado-sección) para el <optgroup> del formulario
def obtener_alumnos_bilingue_agrupados():
    from itertools import groupby
    cfg = _cargar_routing_bl()
    grupos = cfg.get('grupos') or {}
    alumnos = obtener_alumnos_bilingue_cfg()

    def orden(a):
        # <--- hecho por claude code: por grado (1.1 primero … 7/8/9 al final)
        g = grupos.get(a['grado'], {})
        return (g.get('grado_num') or 99, str(g.get('seccion', '')), a['label'])

    alumnos.sort(key=orden)
    salida = []
    for grado, items in groupby(alumnos, key=lambda a: a['grado']):
        g = grupos.get(grado, {})
        nombre = g.get('nombre') or grado
        nivel = g.get('nivel', '')
        etiqueta = _nombre_grupo_disp(nivel, nombre, g.get('seccion', ''))   # 'Prepa 1', '1.1', '9.1'
        salida.append({'grupo': etiqueta or grado, 'alumnos': list(items)})
    return salida


# <--- hecho por claude code: opciones materia/docente por grupo (filtradas a las clases del grupo)
# con el coordinador REAL calculado (materia manda → grupo → override). Fallback: todas las materias.
def _md_por_grado_bl():
    cfg = _cargar_routing_bl()
    grupos = cfg.get('grupos') or {}
    mds = list(MateriaDocenteBilingue.objects.filter(activo=True))

    def _opt(materia, docente, grado):
        coord = _coord_bl(materia, docente, grado, cfg)
        cn = _COORD_NOMBRES.get(coord, '')
        return {'value': f"carga:{docente}|{materia}",
                'label': f"{materia} / {docente}" + (f" ({cn})" if cn else "")}

    def opcion_md(md, grado):
        prim = _split_materias(md.materia)
        coord = _coord_bl(prim[0] if prim else md.materia, md.docente, grado, cfg)
        cn = _COORD_NOMBRES.get(coord, '')
        return {'value': str(md.pk),
                'label': f"{_join_materias(md.materia)} / {md.docente}" + (f" ({cn})" if cn else "")}

    salida = {}
    for grado, g in grupos.items():
        clases = g.get('clases') or []
        if clases:
            # <--- hecho por claude code: opciones desde las clases del grupo, agrupando por docente
            # las materias que van al MISMO coordinador → una sola opción (para maestros con varias clases).
            opts, vistos = [], set()
            for cl in clases:
                doc = (cl.get('docente') or '').strip()
                por_coord = {}
                for mat in _split_materias(cl.get('materia', '')):
                    coord = _coord_bl(mat, doc, grado, cfg)
                    por_coord.setdefault(coord, [])
                    if mat not in por_coord[coord]:
                        por_coord[coord].append(mat)
                for coord, mats in por_coord.items():
                    materia_str = ', '.join(mats)
                    key = (doc.lower(), materia_str.lower())
                    if key in vistos:
                        continue
                    vistos.add(key)
                    cn = _COORD_NOMBRES.get(coord, '')
                    opts.append({'value': f"carga:{doc}|{materia_str}",
                                 'label': f"{materia_str} / {doc}" + (f" ({cn})" if cn else "")})
            salida[grado] = opts
        else:
            salida[grado] = [opcion_md(md, grado) for md in mds]  # respaldo: grupo sin clases
    return salida

def obtener_alumnos_colegio():
    query = """
    SELECT TOP (100) PERCENT d.PersonaID,
           ISNULL(d.Nombre1, '') + 
           CASE WHEN d.Nombre2 IS NOT NULL AND d.Nombre2 <> '' THEN ' ' + d.Nombre2 ELSE '' END + ' ' +
           ISNULL(d.Apellido1, '') + 
           CASE WHEN d.Apellido2 IS NOT NULL AND d.Apellido2 <> '' THEN ' ' + d.Apellido2 ELSE '' END AS NombreCompl,
           da.Descripcion AS AreaDesc, c.CrsoNumero, c.GrupoNumero
      FROM dbo.tblPrsDtosGen AS d
     INNER JOIN dbo.tblPrsTipo           AS t  ON d.PersonaID = t.PersonaID
     INNER JOIN dbo.tblEdcArea           AS a  ON t.IngrEgrID  = a.IngrEgrID
     INNER JOIN dbo.tblEdcEjecCrso       AS ec ON a.AreaID     = ec.AreaID
     INNER JOIN dbo.tblEdcCrso           AS c  ON ec.CrsoID    = c.CrsoID
     INNER JOIN dbo.tblEdcDescripAreaEdc AS da ON a.DescrAreaEdcID = da.DescrAreaEdcID
     WHERE (d.Alum = 1) 
       AND (DATEPART(yy, c.FechaInicio) = DATEPART(yyyy, GETDATE()))
       AND (da.Descripcion IN (N'Colegio', N'Bachillerato')) 
       AND (ec.Desertor = 0) AND (ec.TrasladoPer = 0)
     ORDER BY AreaDesc, c.CrsoNumero, c.GrupoNumero, NombreCompl
    """
    alumnos = []
    try:
        with connections['padres_sqlserver'].cursor() as cursor:
            cursor.execute(query)
            for pid, nombre, area, crso, grupo in cursor.fetchall():
                label = nombre.strip()
                grado = f"{area} {crso}-{grupo}"
                alumnos.append({'id': str(pid), 'label': label, 'grado': grado})
    except Exception as e:
        print(">>> ERROR SQL COLEGIO:", e)
    return alumnos


def get_materia_docente_choices(area):
    if area == 'bilingue':
        # <--- hecho por claude code: etiqueta "Materia / Docente (Coordinador)"
        choices = []
        for md in MateriaDocenteBilingue.objects.filter(activo=True):
            codigos = [c.strip() for c in md.coordinador.split(',') if c.strip()]
            nombres = ', '.join(_COORD_NOMBRES.get(c, c) for c in codigos)
            label = f"{_join_materias(md.materia)} / {md.docente}"  # <--- materias con coma en el display
            if nombres:
                label += f" ({nombres})"
            choices.append((str(md.pk), label, md.coordinador))
        return choices
    else:
        choices = []
        for md in MateriaDocenteColegio.objects.filter(activo=True).prefetch_related('coordinadores'):
            coords = md.coordinadores.all()
            label = f"{md.materia} — {md.docente}"
            coord_codes = ','.join(c.codigo for c in coords)
            if coords:
                label += f" ({', '.join(c.nombre for c in coords)})"
            choices.append((str(md.pk), label, coord_codes))
        return choices

# ---------------------------
# DASHBOARDS Y FORMULARIOS
# ---------------------------
_USUARIOS_MULTI_AREA = {'druiz@ana-hn.org', 'admin2@ana-hn.org'}

@login_required
def dashboard_maestro(request):
    user = request.user
    area = request.GET.get('area')
    if user.username in _USUARIOS_MULTI_AREA:
        if area:
            request.session['maestro_area'] = area
        else:
            area = request.session.get('maestro_area')
    if not area:
        if user.groups.filter(name='maestros_bilingue').exists():
            area = 'bilingue'
        elif user.groups.filter(name='maestros_colegio').exists():
            area = 'colegio'
    return render(request, 'conducta/dashboard_maestros.html', {'area': area})

# ------------ REPORTE INFORMATIVO  ------------
@login_required
def reporte_informativo_bilingue(request):
    area = 'bilingue'
    students = obtener_alumnos_bilingue_cfg()  # <--- hecho por claude code: alumnado desde JSON
    materia_docente_choices = get_materia_docente_choices(area)
    fecha = timezone.now().strftime("%Y-%m-%dT%H:%M")

    if request.method == 'POST':
        alumno_id = request.POST.get('alumno')
        grado = request.POST.get('grado')
        materia_docente_id = request.POST.get('materia_docente')
        comentario = request.POST.get('comentario', "")
        tipo_reporte = request.POST.get('tipo_reporte', 'academico')
        alumno_obj = next((a for a in students if a['id'] == alumno_id), None)
        alumno_label = alumno_obj['label'] if alumno_obj else ""
        materia = docente = coord_bl = ""
        mdid = materia_docente_id or ''
        if mdid.startswith('carga:'):
            # <--- hecho por claude code: opción tomada de la carga del grupo (docente|materia)
            docente, _sep, materia = mdid[6:].partition('|')
            coord_bl = _coord_bl(materia, docente, grado)
        elif mdid:
            md_obj = MateriaDocenteBilingue.objects.filter(pk=mdid).first()
            if md_obj:
                materia = md_obj.materia
                docente = md_obj.docente
                coord_bl = md_obj.coordinador
        ReporteInformativo.objects.create(
            usuario=request.user,
            area=area,
            alumno_id=alumno_id,
            alumno_nombre=alumno_label,
            grado=grado,
            materia=materia,
            docente=docente,
            tipo_reporte=tipo_reporte,
            comentario=comentario,
            # <--- hecho por claude code: ruteo por mapeo JSON (materia manda → grado → override)
            coord_asignado=_coord_bl(materia, docente, grado),
        )
        _notificar_coordinadores("Reporte Informativo", request.user, alumno_label, grado, materia, area,
                                 coordinador_bl=coord_bl, subtipo=f"informativo_{tipo_reporte}")
        messages.success(request, "¡Reporte registrado correctamente!")
        return redirect('reporte_informativo_bilingue')

    return render(request, 'conducta/form_informativo.html', {
        'students': students,
        'students_agrupados': obtener_alumnos_bilingue_agrupados(),  # <--- hecho por claude code: alumnos por grupo
        'materia_docente_choices': materia_docente_choices,
        'md_por_grado_json': json.dumps(_md_por_grado_bl()),  # <--- hecho por claude code: materia/docente por grupo
        'fecha': fecha,
        'area': area,
    })

@login_required
def reporte_informativo_colegio(request):
    area = 'colegio'
    students = obtener_alumnos_colegio()
    materia_docente_choices = get_materia_docente_choices(area)
    fecha = timezone.now().strftime("%Y-%m-%dT%H:%M")

    if request.method == 'POST':
        alumno_id = request.POST.get('alumno')
        grado = request.POST.get('grado')
        materia_docente_id = request.POST.get('materia_docente')
        comentario = request.POST.get('comentario', "")
        tipo_reporte = request.POST.get('tipo_reporte', 'academico')
        alumno_obj = next((a for a in students if a['id'] == alumno_id), None)
        alumno_label = alumno_obj['label'] if alumno_obj else ""
        materia = docente = ""
        md_obj = None
        if materia_docente_id:
            md_obj = MateriaDocenteColegio.objects.filter(pk=materia_docente_id).first()
            if md_obj:
                materia = md_obj.materia
                docente = md_obj.docente
        ReporteInformativo.objects.create(
            usuario=request.user,
            area=area,
            alumno_id=alumno_id,
            alumno_nombre=alumno_label,
            grado=grado,
            materia=materia,
            docente=docente,
            tipo_reporte=tipo_reporte,
            comentario=comentario
        )
        _notificar_coordinadores("Reporte Informativo", request.user, alumno_label, grado, materia, area,
                                 md_colegio=md_obj, subtipo=f"informativo_{tipo_reporte}")
        messages.success(request, "¡Reporte registrado correctamente!")
        return redirect('reporte_informativo_colegio')

    return render(request, 'conducta/form_informativo.html', {
        'students': students,
        'materia_docente_choices': materia_docente_choices,
        'fecha': fecha,
        'area': area,
    })

# ------------ REPORTE CONDUCTUAL  ------------
@login_required
def reporte_conductual_bilingue(request):
    area = 'bilingue'
    students = obtener_alumnos_bilingue_cfg()  # <--- hecho por claude code: alumnado desde JSON
    materia_docente_choices = get_materia_docente_choices(area)
    fecha = timezone.now().strftime("%Y-%m-%d")
    # Obtiene incisos por severidad
    incisos_leve = IncisoConductual.objects.filter(activo=True, tipo='leve').order_by('descripcion')
    incisos_grave = IncisoConductual.objects.filter(activo=True, tipo='grave').order_by('descripcion')
    incisos_muygrave = IncisoConductual.objects.filter(activo=True, tipo='muy_grave').order_by('descripcion')

    if request.method == 'POST':
        alumno_id = request.POST.get('alumno')
        grado = request.POST.get('grado')
        materia_docente_id = request.POST.get('materia_docente')
        fecha_val = request.POST.get('fecha')
        comentario = request.POST.get('comentario', "")

        # Obtén las listas seleccionadas si el checkbox está activo
        ids_leve = request.POST.getlist('inciso_leve[]') if request.POST.get('chk_leve') else []
        ids_grave = request.POST.getlist('inciso_grave[]') if request.POST.get('chk_grave') else []
        ids_muygrave = request.POST.getlist('inciso_muygrave[]') if request.POST.get('chk_muygrave') else []

        alumno_obj = next((a for a in students if a['id'] == alumno_id), None)
        alumno_label = alumno_obj['label'] if alumno_obj else ""
        materia = docente = coord_bl = ""
        mdid = materia_docente_id or ''
        if mdid.startswith('carga:'):
            # <--- hecho por claude code: opción tomada de la carga del grupo (docente|materia)
            docente, _sep, materia = mdid[6:].partition('|')
            coord_bl = _coord_bl(materia, docente, grado)
        elif mdid:
            md_obj = MateriaDocenteBilingue.objects.filter(pk=mdid).first()
            if md_obj:
                materia = md_obj.materia
                docente = md_obj.docente
                coord_bl = md_obj.coordinador

        # Crea el reporte sin los ManyToMany
        # Conductual bilingüe → siempre C3 (Isabel Alcerro)
        reporte = ReporteConductual.objects.create(
            usuario=request.user,
            area=area,
            alumno_id=alumno_id,
            alumno_nombre=alumno_label,
            grado=grado,
            materia=materia,
            docente=docente,
            fecha=fecha_val,
            comentario=comentario,
            coord_asignado='C3' if area == 'bilingue' else '',
        )
        if ids_leve:
            reporte.incisos_leve.set(ids_leve)
        if ids_grave:
            reporte.incisos_grave.set(ids_grave)
        if ids_muygrave:
            reporte.incisos_muygrave.set(ids_muygrave)

        _notificar_coordinadores("Reporte Conductual", request.user, alumno_label, grado, materia, area,
                                 coordinador_bl=coord_bl, subtipo='conductual')
        messages.success(request, "¡Reporte conductual registrado correctamente!")
        return redirect('reporte_conductual_bilingue')

    return render(request, 'conducta/form_conductual.html', {
        'students': students,
        'students_agrupados': obtener_alumnos_bilingue_agrupados(),  # <--- hecho por claude code: alumnos por grupo
        'md_por_grado_json': json.dumps(_md_por_grado_bl()),  # <--- hecho por claude code: materia/docente por grupo
        'materia_docente_choices': materia_docente_choices,
        'fecha': fecha,
        'area': area,
        'incisos_leve': incisos_leve,
        'incisos_grave': incisos_grave,
        'incisos_muygrave': incisos_muygrave,
    })

@login_required
def reporte_conductual_colegio(request):
    area = 'colegio'
    students = obtener_alumnos_colegio()
    materia_docente_choices = get_materia_docente_choices(area)
    fecha = timezone.now().strftime("%Y-%m-%d")
    incisos_leve = IncisoConductual.objects.filter(activo=True, tipo='leve').order_by('descripcion')
    incisos_grave = IncisoConductual.objects.filter(activo=True, tipo='grave').order_by('descripcion')
    incisos_muygrave = IncisoConductual.objects.filter(activo=True, tipo='muy_grave').order_by('descripcion')

    if request.method == 'POST':
        alumno_id = request.POST.get('alumno')
        grado = request.POST.get('grado')
        materia_docente_id = request.POST.get('materia_docente')
        fecha_val = request.POST.get('fecha')
        comentario = request.POST.get('comentario', "")

        ids_leve = request.POST.getlist('inciso_leve[]') if request.POST.get('chk_leve') else []
        ids_grave = request.POST.getlist('inciso_grave[]') if request.POST.get('chk_grave') else []
        ids_muygrave = request.POST.getlist('inciso_muygrave[]') if request.POST.get('chk_muygrave') else []

        alumno_obj = next((a for a in students if a['id'] == alumno_id), None)
        alumno_label = alumno_obj['label'] if alumno_obj else ""
        materia = docente = ""
        md_obj = None
        if materia_docente_id:
            md_obj = MateriaDocenteColegio.objects.filter(pk=materia_docente_id).first()
            if md_obj:
                materia = md_obj.materia
                docente = md_obj.docente

        reporte = ReporteConductual.objects.create(
            usuario=request.user,
            area=area,
            alumno_id=alumno_id,
            alumno_nombre=alumno_label,
            grado=grado,
            materia=materia,
            docente=docente,
            fecha=fecha_val,
            comentario=comentario,
            coord_asignado='C3' if area == 'bilingue' else '',
        )
        if ids_leve:
            reporte.incisos_leve.set(ids_leve)
        if ids_grave:
            reporte.incisos_grave.set(ids_grave)
        if ids_muygrave:
            reporte.incisos_muygrave.set(ids_muygrave)

        _notificar_coordinadores("Reporte Conductual", request.user, alumno_label, grado, materia, area,
                                 md_colegio=md_obj, subtipo='conductual')
        messages.success(request, "¡Reporte conductual registrado correctamente!")
        return redirect('reporte_conductual_colegio')

    return render(request, 'conducta/form_conductual.html', {
        'students': students,
        'materia_docente_choices': materia_docente_choices,
        'fecha': fecha,
        'area': area,
        'incisos_leve': incisos_leve,
        'incisos_grave': incisos_grave,
        'incisos_muygrave': incisos_muygrave,
    })

#-------------- PROGRESS REPORT -----------------
@login_required
def progress_report_bilingue(request):
    MATERIAS_PRIMARIA = [
        "Math", "Phonics", "Reading", "Language",
        "Science", "Español", "CCSS", "Asociadas"
    ]
    MATERIAS_COLEGIO = [
        "Math", "Spelling", "Reading", "Language", "Science",
        "Español", "CCSS", "Cívica", "Asociadas"
    ]

    students = obtener_alumnos_bilingue_cfg()  # <--- hecho por claude code: alumnado desde JSON
    alumnos_choices = [(s['id'], s['label']) for s in students]

    materias = MATERIAS_PRIMARIA
    grado = ""
    alumno_id = ""
    if request.method == 'POST':
        form = ProgressReportForm(alumnos_choices=alumnos_choices, data=request.POST)
        alumno_id = request.POST.get('alumno')
        grado = request.POST.get('grado', '')
        if alumno_id:
            alumno_obj = next((s for s in students if s['id'] == alumno_id), None)
            if alumno_obj:
                grado = alumno_obj['grado']
        if grado and ("Primaria" in grado or "Preescolar" in grado):
            materias = MATERIAS_PRIMARIA
        else:
            materias = MATERIAS_COLEGIO

        if form.is_valid():
            # --- PROCESAR MATERIAS ---
            usuario_actual = request.user.get_full_name() or request.user.username
            materias_list = []
            for materia in materias:
                if materia == "Asociadas":
                    asignaciones = request.POST.getlist('asignacion_Asociadas[]')
                    comentarios = request.POST.getlist('comentario_Asociadas[]')
                    # Máximo 5 asociadas
                    for idx, asignacion in enumerate(asignaciones[:5]):
                        comentario = comentarios[idx] if idx < len(comentarios) else ""
                        materias_list.append({
                            'materia': 'Asociadas',
                            'asignacion': asignacion,
                            'comentario': comentario,
                            'docente': usuario_actual if (asignacion.strip() or comentario.strip()) else '',
                        })
                else:
                    asignacion = request.POST.get(f"asignacion_{materia}", "")
                    comentario = request.POST.get(f"comentario_{materia}", "")
                    materias_list.append({
                        'materia': materia,
                        'asignacion': asignacion,
                        'comentario': comentario,
                        'docente': usuario_actual if (asignacion.strip() or comentario.strip()) else '',
                    })
            # --- DATOS GENERALES ---
            semana_inicio = form.cleaned_data.get('semana_inicio')
            semana_fin = form.cleaned_data.get('semana_fin')
            comentario_general = form.cleaned_data.get('comentario_general', "")
            alumno_id = form.cleaned_data.get('alumno')
            alumno_obj = next((s for s in students if s['id'] == alumno_id), None)
            alumno_label = alumno_obj['label'] if alumno_obj else ""
            grado = form.cleaned_data.get('grado')

            # --- CREAR REPORTE ---
            ProgressReport.objects.create(
                usuario=request.user,
                alumno_id=alumno_id,
                alumno_nombre=alumno_label,
                grado=grado,
                semana_inicio=semana_inicio,
                semana_fin=semana_fin,
                comentario_general=comentario_general,
                materias_json=materias_list
            )
            _notificar_coordinadores("Progress Report", request.user, alumno_label, grado, "", "bilingue",
                                     subtipo='progress')
            messages.success(request, "¡Progress report registrado correctamente!")
            return redirect('progress_report_bilingue')
    else:
        form = ProgressReportForm(alumnos_choices=alumnos_choices)
        materias = MATERIAS_PRIMARIA

    return render(request, 'conducta/form_progress.html', {
        'form': form,
        'materias': materias,
        'students': students,
        'students_agrupados': obtener_alumnos_bilingue_agrupados(),  # <--- hecho por claude code: alumnos por grupo
        'grado': grado,
    })


#-------------- HISTORIAL MAESTROS -----------------
def _agendas_historial(usuario, areas):
    """Todas las agendas del área (el maestro ve todas, aunque no las haya creado)."""
    from agendas.models import Agenda
    qs = (Agenda.objects.filter(grado__area__in=areas)
          .select_related('grado', 'usuario').order_by('-semana_inicio', 'grado__nombre'))
    out = []
    for a in qs:
        mats = a.materias_json or []
        llenas = sum(1 for m in mats if any(m.get(d, '').strip()
                     for d in ['lunes', 'martes', 'miercoles', 'jueves', 'viernes', 'nota']))
        docente = (a.usuario.get_full_name() or a.usuario.username) if a.usuario else '—'
        out.append({
            'pk': a.pk,
            'grado': a.grado.nombre,
            'docente': docente,
            'semana_inicio': a.semana_inicio,
            'semana_fin': a.semana_fin,
            'estado': a.get_estado_display(),
            'n_materias': llenas,
        })
    return out


@login_required
def historial_maestro_bilingue(request):
    usuario = request.user
    usuario_actual = usuario.get_full_name()
    request.session['maestro_area'] = 'bilingue'
    reportes_informativo = ReporteInformativo.objects.filter(usuario=usuario, area='bilingue').order_by('-fecha')
    reportes_conductual = ReporteConductual.objects.filter(usuario=usuario, area='bilingue').order_by('-fecha')

    reportes_progress = ProgressReport.objects.all().order_by('-fecha')

    # Tutorías solicitadas por este docente (marcas de la tabla por grado)
    from .models import ConvocatoriaMarca, GRADO_TUTORIA_LABEL
    _tut_map = {}
    for m in (ConvocatoriaMarca.objects.filter(area='bilingue', marcado_por=usuario)
              .order_by('-creado_at')):
        key = (m.alumno_id, m.parcial, m.anio)
        d = _tut_map.get(key)
        if not d:
            lbl = GRADO_TUTORIA_LABEL.get(m.grado_num, m.grado_num)
            d = {'alumno_nombre': m.alumno_nombre,
                 'grado': f"{lbl} {m.grado_num}.{m.seccion}" if m.seccion else str(lbl),
                 'parcial': m.parcial, 'anio': m.anio, 'asigs': [], 'fecha': m.creado_at}
            _tut_map[key] = d
        d['asigs'].append(m.asignatura)
    convocatorias_tut = []
    for d in _tut_map.values():
        d['asigs_txt'] = ', '.join(d['asigs'])
        convocatorias_tut.append(d)

    tickets_usuario = Ticket.objects.filter(email=usuario.email).order_by('-created_at')

    agendas_usuario = _agendas_historial(usuario, ['primaria', 'colegio_bl'])

    es_admin = request.user.groups.filter(name='administracion').exists()
    back_url = 'tickets/submit_ticket/' if es_admin else '/conducta/dashboard/maestro/'

    return render(request, 'conducta/historial_maestro.html', {
        'reportes_informativo': reportes_informativo,
        'reportes_conductual': reportes_conductual,
        'reportes_progress': reportes_progress,
        'convocatorias_tut': convocatorias_tut,
        'agendas_usuario': agendas_usuario,
        'tickets_usuario': tickets_usuario,
        'area': 'bilingue',
        'back_url': back_url,
        'today': timezone.now().strftime('%Y-%m-%d'),
    })


@login_required
def historial_maestro_colegio(request):
    usuario = request.user
    request.session['maestro_area'] = 'colegio'
    reportes_informativo = ReporteInformativo.objects.filter(usuario=usuario, area='colegio').order_by('-fecha')
    reportes_conductual = ReporteConductual.objects.filter(usuario=usuario, area='colegio').order_by('-fecha')

    tickets_usuario = Ticket.objects.filter(email=usuario.email).order_by('-created_at')

    agendas_usuario = _agendas_historial(usuario, ['colegio'])

    es_admin = request.user.groups.filter(name='administracion').exists()
    back_url = 'tickets/submit_ticket/' if es_admin else '/conducta/dashboard/maestro/'

    return render(request, 'conducta/historial_maestro.html', {
        'reportes_informativo': reportes_informativo,
        'reportes_conductual': reportes_conductual,
        'reportes_progress': [],
        'agendas_usuario': agendas_usuario,
        'tickets_usuario': tickets_usuario,
        'area': 'colegio',
        'back_url': back_url,
        'today': timezone.now().strftime('%Y-%m-%d'),
    })

# ────────────────────────────────────────────────────────────
# REENVIAR NOTIFICACIONES A COORDINADORES (solo superuser)
# ────────────────────────────────────────────────────────────
@login_required
def reenviar_reportes_coordinadores(request):
    if not request.user.is_superuser:
        return JsonResponse({'ok': False, 'error': 'Sin permisos'}, status=403)

    fecha_str = request.POST.get('fecha', '').strip()
    if not fecha_str:
        return JsonResponse({'ok': False, 'error': 'Fecha requerida'}, status=400)

    from datetime import date as _date_type
    try:
        fecha = _date_type.fromisoformat(fecha_str)
    except ValueError:
        return JsonResponse({'ok': False, 'error': 'Fecha inválida'}, status=400)

    enviados = 0
    for r in ReporteInformativo.objects.filter(fecha__date=fecha):
        _notificar_coordinadores("Reporte Informativo", r.usuario, r.alumno_nombre, r.grado, r.materia, r.area,
                                 subtipo=f"informativo_{r.tipo_reporte}")
        enviados += 1
    for r in ReporteConductual.objects.filter(fecha=fecha):
        _notificar_coordinadores("Reporte Conductual", r.usuario, r.alumno_nombre, r.grado, r.materia, r.area,
                                 subtipo='conductual')
        enviados += 1
    for r in ProgressReport.objects.filter(semana_inicio__lte=fecha, semana_fin__gte=fecha):
        _notificar_coordinadores("Progress Report", r.usuario, r.alumno_nombre, r.grado, "", "bilingue",
                                 subtipo='progress')
        enviados += 1

    return JsonResponse({'ok': True, 'enviados': enviados})


# ────────────────
# EDITAR CONDUCTUAL
# ────────────────

@login_required
def editar_reporte_conductual(request, pk):
    reporte = get_object_or_404(ReporteConductual, pk=pk)

    # Elige coordinadores según área
    if reporte.area == "bilingue":
        coordinadores = COORDINADORES_BL
    else:
        coordinadores = COORDINADORES_COLEGIO

    # <--- hecho por claude code: se cargan todos los incisos activos para que el
    # coordinador pueda agregar/quitar incisos al editar el reporte conductual.
    incisos_leve     = IncisoConductual.objects.filter(activo=True, tipo='leve').order_by('descripcion')
    incisos_grave    = IncisoConductual.objects.filter(activo=True, tipo='grave').order_by('descripcion')
    incisos_muygrave = IncisoConductual.objects.filter(activo=True, tipo='muy_grave').order_by('descripcion')

    if request.method == "POST":
        reporte.comentario            = request.POST.get("comentario", "")
        reporte.comentario_coordinador = request.POST.get("comentario_coordinador", "")
        reporte.estado                = request.POST.get("estado", "enviado")
        reporte.coordinador_firma     = request.POST.get("coordinador_firma", "")
        reporte.save()
        # <--- hecho por claude code: actualizar incisos desde los checkboxes del coordinador
        ids_leve     = request.POST.getlist("incisos_leve")
        ids_grave    = request.POST.getlist("incisos_grave")
        ids_muygrave = request.POST.getlist("incisos_muygrave")
        reporte.incisos_leve.set(ids_leve)
        reporte.incisos_grave.set(ids_grave)
        reporte.incisos_muygrave.set(ids_muygrave)
        messages.success(request, "Reporte conductual actualizado correctamente.")
        return redirect('dashboard_coordinador', area=reporte.area)
    return render(request, "conducta/editor_conductual.html", {
        "reporte": reporte,
        "coordinadores": coordinadores,
        "incisos_leve": incisos_leve,
        "incisos_grave": incisos_grave,
        "incisos_muygrave": incisos_muygrave,
    })

# ────────────────
# EDITAR INFORMATIVO
# ────────────────

@login_required
def editar_reporte_informativo(request, pk):
    reporte = get_object_or_404(ReporteInformativo, pk=pk)

    # Elige coordinadores según área
    if reporte.area == "bilingue":
        coordinadores = COORDINADORES_BL
    else:
        coordinadores = COORDINADORES_COLEGIO

    if request.method == "POST":
        reporte.comentario            = request.POST.get("comentario", "")
        reporte.comentario_coordinador = request.POST.get("comentario_coordinador", "")
        reporte.estado                = request.POST.get("estado", "enviado")
        reporte.coordinador_firma     = request.POST.get("coordinador_firma", "")
        reporte.tipo_reporte          = request.POST.get("tipo_reporte", "academico")
        reporte.save()
        messages.success(request, "Reporte informativo actualizado correctamente.")
        return redirect('dashboard_coordinador', area=reporte.area)
    return render(request, "conducta/editor_informativo.html", {
        "reporte": reporte,
        "coordinadores": coordinadores,
    })

# ────────────────────────────────────────────────────────────────
# ELIMINAR REPORTES (AJAX) — solo coordinadores/staff
# Notifica a todos los superusuarios al eliminar.
# ────────────────────────────────────────────────────────────────

@login_required
def eliminar_reporte_conductual(request, pk):
    if request.method != 'POST':
        return JsonResponse({'ok': False}, status=405)
    if not (request.user.is_staff or request.user.groups.filter(name__icontains='coordinador').exists()):
        return JsonResponse({'ok': False, 'error': 'Sin permiso'}, status=403)

    reporte = get_object_or_404(ReporteConductual, pk=pk)
    info = {
        'pk': reporte.pk,
        'alumno': reporte.alumno_nombre,
        'materia': reporte.materia,
        'docente': reporte.docente,
        'fecha': reporte.fecha.strftime('%d/%m/%Y %H:%M') if reporte.fecha else '—',
        'eliminado_por': request.user.get_full_name() or request.user.username,
    }
    reporte.delete()

    msg = (f"⚠️ Reporte Conductual #{info['pk']} eliminado por {info['eliminado_por']}. "
           f"Alumno: {info['alumno']} | Materia: {info['materia']} – {info['docente']} | "
           f"Fecha: {info['fecha']}")
    for su in User.objects.filter(is_superuser=True):
        crear_notificacion(su, msg, 'conducta', tipo='alerta')

    return JsonResponse({'ok': True})


@login_required
def eliminar_reporte_informativo(request, pk):
    if request.method != 'POST':
        return JsonResponse({'ok': False}, status=405)
    if not (request.user.is_staff or request.user.groups.filter(name__icontains='coordinador').exists()):
        return JsonResponse({'ok': False, 'error': 'Sin permiso'}, status=403)

    reporte = get_object_or_404(ReporteInformativo, pk=pk)
    info = {
        'pk': reporte.pk,
        'alumno': reporte.alumno_nombre,
        'materia': reporte.materia,
        'docente': reporte.docente,
        'tipo': reporte.get_tipo_reporte_display(),
        'fecha': reporte.fecha.strftime('%d/%m/%Y %H:%M') if reporte.fecha else '—',
        'eliminado_por': request.user.get_full_name() or request.user.username,
    }
    reporte.delete()

    msg = (f"⚠️ Reporte Informativo #{info['pk']} ({info['tipo']}) eliminado por {info['eliminado_por']}. "
           f"Alumno: {info['alumno']} | Materia: {info['materia']} – {info['docente']} | "
           f"Fecha: {info['fecha']}")
    for su in User.objects.filter(is_superuser=True):
        crear_notificacion(su, msg, 'conducta', tipo='alerta')

    return JsonResponse({'ok': True})


@login_required
@require_POST
def bulk_eliminar_reportes(request):
    """
    AJAX: elimina múltiples reportes a la vez.
    Body: { "items": [{"pk": 1, "tipo": "informativo"}, ...] }
    <--- hecho por claude code
    """
    if not (request.user.is_staff or request.user.groups.filter(name__icontains='coordinador').exists()):
        return JsonResponse({'ok': False, 'error': 'Sin permiso'}, status=403)

    try:
        body  = json.loads(request.body or b'{}')
        items = body.get('items', [])
    except Exception as e:
        return JsonResponse({'ok': False, 'error': str(e)})

    eliminador = request.user.get_full_name() or request.user.username
    eliminados = 0
    errores    = []

    for item in items:
        pk   = item.get('pk')
        tipo = item.get('tipo', '')
        try:
            if tipo in ('informativo', 'academico', 'conductual_inf'):
                obj = ReporteInformativo.objects.get(pk=pk)
            elif tipo == 'conductual':
                obj = ReporteConductual.objects.get(pk=pk)
            elif tipo == 'progress':
                obj = ProgressReport.objects.get(pk=pk)
            else:
                errores.append(f'Tipo desconocido: {tipo}')
                continue
            obj.delete()
            eliminados += 1
        except Exception as e:
            errores.append(f'pk={pk}: {e}')

    if eliminados:
        msg = f"⚠️ Eliminación masiva: {eliminados} reporte(s) eliminados por {eliminador}."
        for su in User.objects.filter(is_superuser=True):
            crear_notificacion(su, msg, 'conducta', tipo='alerta')

    return JsonResponse({'ok': True, 'eliminados': eliminados, 'errores': errores})


@login_required
def set_coord_reporte(request):
    """AJAX: superuser cambia el coordinador asignado de un reporte."""
    if not request.user.is_superuser:
        return JsonResponse({'ok': False, 'error': 'Sin permiso'}, status=403)
    if request.method != 'POST':
        return JsonResponse({'ok': False}, status=405)

    pk    = request.POST.get('pk')
    tipo  = request.POST.get('tipo')
    coord = request.POST.get('coord', '').strip().upper()

    if coord not in {'C1', 'C2', 'C3', 'C4', ''}:
        return JsonResponse({'ok': False, 'error': 'Coordinador inválido'}, status=400)

    if tipo == 'conductual':
        obj = get_object_or_404(ReporteConductual, pk=pk)
    elif tipo == 'informativo':
        obj = get_object_or_404(ReporteInformativo, pk=pk)
    else:
        return JsonResponse({'ok': False, 'error': 'Tipo inválido'}, status=400)

    obj.coord_asignado = coord
    obj.save(update_fields=['coord_asignado'])
    return JsonResponse({'ok': True, 'coord': coord or obj.coord_codigo})


# ────────────────────────────────────────────────────────────────
# EDITAR REPORTE AJAX — modal inline del dashboard coordinador
# GET  → devuelve JSON con los campos del reporte
# POST → guarda los cambios y devuelve {'ok': True}
# ────────────────────────────────────────────────────────────────
@login_required
def editar_reporte_ajax(request):
    if not es_coordinador(request.user):
        return JsonResponse({'ok': False, 'error': 'Sin permiso'}, status=403)

    pk   = request.GET.get('pk')   or request.POST.get('pk')
    tipo = request.GET.get('tipo') or request.POST.get('tipo')

    if tipo == 'conductual':
        obj = get_object_or_404(ReporteConductual, pk=pk)
    elif tipo == 'informativo':
        obj = get_object_or_404(ReporteInformativo, pk=pk)
    elif tipo == 'progress':
        obj = get_object_or_404(ProgressReport, pk=pk)
    else:
        return JsonResponse({'ok': False, 'error': 'Tipo inválido'}, status=400)

    if request.method == 'GET':
        data = {
            'ok':    True,
            'pk':    obj.pk,
            'tipo':  tipo,
            'alumno': obj.alumno_nombre,
            'grado':  obj.grado,
            'fecha':  obj.fecha.strftime('%d/%m/%Y %H:%M') if obj.fecha else '',
            'estado':              obj.estado,
            'coordinador_firma':   obj.coordinador_firma or '',
            'comentario_coordinador': obj.comentario_coordinador or '',
            'area': getattr(obj, 'area', 'bilingue'),
        }
        if tipo in ('conductual', 'informativo'):
            data['materia']    = obj.materia
            data['docente']    = obj.docente
            data['comentario'] = obj.comentario or ''
        if tipo == 'informativo':
            data['tipo_reporte'] = obj.tipo_reporte
        return JsonResponse(data)

    if request.method == 'POST':
        obj.estado             = request.POST.get('estado', obj.estado)
        obj.coordinador_firma  = request.POST.get('coordinador_firma', '') or None
        obj.comentario_coordinador = request.POST.get('comentario_coordinador', '') or None
        if tipo in ('conductual', 'informativo'):
            obj.comentario = request.POST.get('comentario', '') or None
        if tipo == 'informativo':
            obj.tipo_reporte = request.POST.get('tipo_reporte', obj.tipo_reporte)
        obj.save()
        return JsonResponse({'ok': True})

    return JsonResponse({'ok': False}, status=405)


# ────────────────
# EDITAR PROGRESS
# ────────────────

@login_required
def editar_progress_report(request, pk):
    # <--- hecho por claude code: view reescrita para corregir 3 bugs críticos:
    # Bug 1: Maestro A reclamaba TODAS las materias al guardar (incluso vacías),
    #         bloqueando a Maestro B. Fix: docente solo se asigna si hay contenido.
    # Bug 2: Indentación rota en el loop de Asociadas → NameError cuando el
    #         coordinador guardaba. Fix: loop correctamente indentado.
    # Bug 3: Solo se guardaba UNA fila de Asociadas (la última). Fix: loop correcto.
    # Extra: json.dumps() en un JSONField causaba doble serialización. Fix: lista directa.

    COORDINADORES_BL = ["Mr. Martinez", "Miss Alcerro", "Mr. Ruiz" , "Mrs. Varela"]

    reporte = get_object_or_404(ProgressReport, pk=pk)
    es_coord = es_coordinador(request.user)
    usuario_actual = request.user.get_full_name() or request.user.username

    # Decodificar materias — JSONField puede devolver str si fue guardado con json.dumps()
    materias = reporte.materias_json
    if isinstance(materias, str):
        try:
            materias = json.loads(materias)
        except Exception:
            materias = []
    if not materias:
        materias = []

    # Marcar editabilidad:
    # - Coordinador: siempre editable
    # - Materia vacía (sin contenido y sin docente): cualquier maestro puede llenarla
    # - Materia del propio maestro: editable
    # - Materia de otro maestro, o con contenido sin docente (legacy): readonly
    for mat in materias:
        docente_guardado  = mat.get("docente", "").strip()
        tiene_contenido   = bool(mat.get("asignacion", "").strip() or mat.get("comentario", "").strip())
        if es_coord:
            mat["editable"] = True
        elif docente_guardado == usuario_actual:
            mat["editable"] = True
        elif not docente_guardado and not tiene_contenido:
            mat["editable"] = True   # materia vacía, disponible para cualquiera
        else:
            mat["editable"] = False  # tiene contenido de otro maestro → protegida

    if request.method == 'POST':
        nuevas_materias = []

        # Separar Asociadas existentes para preservar sus docentes
        asociadas_existentes = [m for m in materias if m.get('materia') == 'Asociadas']

        # ── 1. Procesar materias normales (no Asociadas) ──────────────────────
        for mat in materias:
            if mat.get('materia') == 'Asociadas':
                continue  # se procesan en el bloque siguiente

            nombre = mat.get('materia', '')

            if mat["editable"]:
                asignacion_nueva   = request.POST.get(f'asignacion_{nombre}', '').strip()
                comentario_nuevo   = request.POST.get(f'comentario_{nombre}', '').strip()
                asignacion_previa  = mat.get('asignacion', '').strip()
                comentario_previo  = mat.get('comentario', '').strip()

                mat['asignacion'] = asignacion_nueva
                mat['comentario'] = comentario_nuevo

                if not es_coord:
                    # <--- hecho por claude code: fix crítico
                    # Solo reclamar docencia si el maestro REALMENTE cambió el valor.
                    # Sin esto, abrir el form y guardar sin tocar una materia de otro
                    # maestro la reclama como propia (readonly envía el valor igual).
                    cambio_real = (asignacion_nueva != asignacion_previa) or \
                                  (comentario_nuevo != comentario_previo)

                    if cambio_real:
                        if asignacion_nueva or comentario_nuevo:
                            mat['docente'] = usuario_actual  # el usuario editó → es suya
                        else:
                            mat['docente'] = ''  # borró todo → libera la materia
                    # Si no hubo cambio real, preservar el docente original tal como está

            # Guardar sin la clave 'editable' (es solo para el template, no va a la BD)
            nuevas_materias.append({k: v for k, v in mat.items() if k != 'editable'})

        # ── 2. Procesar TODAS las filas de Asociadas (máximo 5) ──────────────
        asignaciones_asociadas = request.POST.getlist('asignacion_Asociadas[]')
        comentarios_asociadas  = request.POST.getlist('comentario_Asociadas[]')

        for i, asignacion in enumerate(asignaciones_asociadas[:5]):
            comentario = comentarios_asociadas[i] if i < len(comentarios_asociadas) else ''

            # Determinar docente: preservar el original si la fila ya existía y no es editable
            if i < len(asociadas_existentes):
                docente_original  = asociadas_existentes[i].get('docente', '')
                era_editable      = asociadas_existentes[i].get('editable', True)
                if era_editable and not es_coord:
                    docente_valor = usuario_actual if (asignacion.strip() or comentario.strip()) else docente_original
                elif es_coord:
                    docente_valor = docente_original  # coordinador no cambia docente
                else:
                    docente_valor = docente_original  # fila ajena → preservar
            else:
                # Fila nueva agregada con el botón "+"
                docente_valor = usuario_actual if (asignacion.strip() or comentario.strip()) else ''

            nuevas_materias.append({
                'materia':    'Asociadas',
                'asignacion': asignacion,
                'comentario': comentario,
                'docente':    docente_valor,
            })

        # ── 3. Guardar ────────────────────────────────────────────────────────
        # JSONField acepta la lista directamente; json.dumps() causaba doble serialización
        reporte.materias_json     = nuevas_materias
        reporte.comentario_general = request.POST.get('comentario_general', reporte.comentario_general)

        if es_coord:
            reporte.coordinador_firma      = request.POST.get('coordinador_firma', reporte.coordinador_firma)
            reporte.estado                 = request.POST.get('estado', reporte.estado)
            reporte.comentario_coordinador = request.POST.get('comentario_coordinador', reporte.comentario_coordinador)

        reporte.save()
        messages.success(request, "¡Reporte actualizado correctamente!")

        if es_coord:
            return redirect('dashboard_coordinador', area='bilingue')
        return redirect('historial_maestro_bilingue')

    return render(request, 'conducta/editor_progress.html', {
        'reporte':       reporte,
        'materias':      materias,
        'es_coordinador': es_coord,
        'usuario_actual': usuario_actual,
        'coordinadores':  COORDINADORES_BL,
    })

# -----------  DESCARGA EN PDF  -----------
def draw_paragraph(pdf, text, x, y, max_width, font="Helvetica", font_size=10, bold=False, italic=False, leading=13):
    """
    Dibuja un párrafo con saltos de línea y justificación manual en ReportLab.
    """
    from reportlab.pdfbase.pdfmetrics import stringWidth
    fontname = font
    if bold and italic:
        fontname = "Helvetica-BoldOblique"
    elif bold:
        fontname = "Helvetica-Bold"
    elif italic:
        fontname = "Helvetica-Oblique"
    pdf.setFont(fontname, font_size)
    lines = []
    for raw_line in text.split('\n'):
        line = ""
        for word in raw_line.split():
            test_line = f"{line} {word}".strip()
            if stringWidth(test_line, fontname, font_size) < max_width:
                line = test_line
            else:
                lines.append(line)
                line = word
        lines.append(line)
    for l in lines:
        pdf.drawString(x, y, l)
        y -= leading
    return y


@login_required
def descargar_pdf_informativo(request, pk):
    from .models import ReporteInformativo
    reporte = get_object_or_404(ReporteInformativo, pk=pk)
    buf = io.BytesIO()
    w, h = letter
    pdf = canvas.Canvas(buf, pagesize=letter)
    pdf.setTitle("reporte_informativo.pdf")

    # Logo centrado
    width_logo = 35 * mm
    height_logo = 35 * mm
    x_logo = (w - width_logo) / 2
    logo_path = os.path.join(settings.STATIC_ROOT, "conducta/img/ana-transformed.png")
    if os.path.exists(logo_path):
        pdf.drawImage(
            logo_path,
            x=x_logo,
            y=h-40*mm,
            width=width_logo,
            height=height_logo,
            mask='auto'
        )
    # 🔥 Más margen superior
    y_actual = h - 52*mm

    # Título y datos
    pdf.setFont("Helvetica-Bold", 16)
    pdf.drawCentredString(w/2, y_actual, "Nuevo Amanecer School" if reporte.area == "bilingue" else "C.E.M.N.G Nuevo Amanecer")
    y_actual -= 14*mm
    pdf.setFont("Helvetica", 13)
    pdf.drawCentredString(w/2, y_actual, "Reporte informativo")
    y_actual -= 14*mm

    # <--- hecho por claude code: offsets calculados con stringWidth para que el valor
    # quede pegado a la etiqueta sin espacio extra visible.
    def label_val(pdf, x, y, label, valor, lbl_font="Helvetica-Bold", lbl_size=11, val_font="Helvetica", val_size=10, gap=2):
        pdf.setFont(lbl_font, lbl_size)
        pdf.drawString(x, y, label)
        ancho = pdf.stringWidth(label, lbl_font, lbl_size)
        pdf.setFont(val_font, val_size)
        pdf.drawString(x + ancho + gap, y, valor)

    label_val(pdf, 32*mm, y_actual, "Nombre:", reporte.alumno_nombre)
    label_val(pdf, 118*mm, y_actual, "Grado:", reporte.grado)
    y_actual -= 8*mm
    label_val(pdf, 32*mm, y_actual, "Docente:", reporte.docente or '')
    label_val(pdf, 100*mm, y_actual, "Fecha:", reporte.fecha.strftime('%d/%m/%Y') if reporte.fecha else '')
    y_actual -= 12*mm

    # Comentario del docente
    pdf.setFont("Helvetica-Bold", 11)
    pdf.drawString(32*mm, y_actual, "Comentario del Docente:")
    y_actual -= 8*mm
    pdf.setFont("Helvetica-Oblique", 10)
    y_actual = draw_paragraph(
        pdf, reporte.comentario or "-", x=38*mm, y=y_actual, max_width=w-65*mm,
        font="Helvetica", font_size=10, italic=True, leading=12
    )
    y_actual -= 10*mm

    # Comentario del coordinador si existe
    if hasattr(reporte, 'comentario_coordinador') and reporte.comentario_coordinador:
        pdf.setFont("Helvetica-Bold", 11)
        pdf.drawString(32*mm, y_actual, "Comentario del Coordinador:")
        y_actual -= 8*mm
        pdf.setFont("Helvetica-Oblique", 10)
        y_actual = draw_paragraph(
            pdf, reporte.comentario_coordinador, x=38*mm, y=y_actual, max_width=w-65*mm,
            font="Helvetica", font_size=10, italic=True, leading=12
        )
        y_actual -= 8*mm

    # <--- hecho por claude code: firmas fijas hacia abajo de la página (28mm del borde),
    # líneas más largas, y "Firma Pariente:" en lugar de "Firma Padre/Madre:".
    y_firma = min(y_actual - 15*mm, 28*mm)
    largo_firma = 58 * mm
    x_col1 = 15 * mm
    x_col2 = 78 * mm
    x_col3 = 141 * mm

    pdf.setStrokeColor(colors.black)
    pdf.setLineWidth(0.7)

    for x_col, label, nombre in [
        (x_col1, "Firma:", reporte.docente or ""),
        (x_col2, "Firma:", getattr(reporte, 'coordinador_firma', '') or ""),
        (x_col3, "Firma Pariente:", ""),
    ]:
        pdf.setFont("Helvetica-Bold", 10)
        pdf.drawString(x_col, y_firma + 5*mm, label)
        ancho_lbl = pdf.stringWidth(label, "Helvetica-Bold", 10)
        pdf.setFont("Helvetica", 10)
        pdf.drawString(x_col + ancho_lbl + 2, y_firma + 5*mm, nombre)
        pdf.line(x_col, y_firma, x_col + largo_firma, y_firma)

    pdf.save()
    buf.seek(0)
    grado_slug = reporte.grado.replace(' ', '_').replace('/', '-')
    nombre_archivo = f"informativo_{reporte.alumno_nombre.replace(' ', '_')}_{grado_slug}.pdf"
    response = HttpResponse(buf, content_type="application/pdf")
    response['Content-Disposition'] = f'inline; filename="{nombre_archivo}"'
    return response


@login_required
def descargar_pdf_conductual(request, pk):
    from .models import ReporteConductual
    reporte = get_object_or_404(ReporteConductual, pk=pk)
    buf = io.BytesIO()
    w, h = letter
    pdf = canvas.Canvas(buf, pagesize=letter)
    pdf.setTitle("reporte_conductual.pdf")

    # <--- hecho por claude code: logo reducido a 22mm y reposicionado para comprimir espacio superior
    width_logo = 22 * mm
    height_logo = 22 * mm
    x_logo = (w - width_logo) / 2
    logo_path = os.path.join(settings.STATIC_ROOT, "conducta/img/ana-transformed.png")
    if os.path.exists(logo_path):
        pdf.drawImage(
            logo_path,
            x=x_logo,
            y=h - 28*mm,
            width=width_logo,
            height=height_logo,
            mask='auto'
        )
    y_actual = h - 33*mm

    # Título
    pdf.setFont("Helvetica-Bold", 16)
    titulo = "Nuevo Amanecer School" if getattr(reporte, "area", None) == "bilingue" else "C.E.M.N.G Nuevo Amanecer"
    pdf.drawCentredString(w/2, y_actual, titulo)
    y_actual -= 9*mm
    pdf.setFont("Helvetica", 13)
    pdf.drawCentredString(w/2, y_actual, "Reporte conductual")
    y_actual -= 10*mm

    # <--- hecho por claude code: misma función label_val para eliminar espacios extra
    def label_val(pdf, x, y, label, valor, lbl_font="Helvetica-Bold", lbl_size=11, val_font="Helvetica", val_size=11, gap=2):
        pdf.setFont(lbl_font, lbl_size)
        pdf.drawString(x, y, label)
        ancho = pdf.stringWidth(label, lbl_font, lbl_size)
        pdf.setFont(val_font, val_size)
        pdf.drawString(x + ancho + gap, y, valor)

    label_val(pdf, 32*mm, y_actual, "Nombre:", reporte.alumno_nombre)

    y_actual -= 7*mm
    label_val(pdf, 32*mm, y_actual, "Grado:", reporte.grado)
    label_val(pdf, 118*mm, y_actual, "Fecha:", reporte.fecha.strftime('%d/%m/%Y') if reporte.fecha else '')

    y_actual -= 7*mm
    label_val(pdf, 32*mm, y_actual, "Docente:", reporte.docente or "")
    y_actual -= 9*mm

    # ----------- INCISOS -----------
    maxw = w - 65*mm
    for tipo, label in [("incisos_leve", "Leve"), ("incisos_grave", "Grave"), ("incisos_muygrave", "Muy Grave")]:
        incisos = getattr(reporte, tipo).all()
        pdf.setFont("Helvetica-Bold", 11)
        pdf.drawString(32*mm, y_actual, f"Incisos {label}:")
        y_actual -= 6*mm
        if incisos:
            for i in incisos:
                y_actual = draw_paragraph(
                    pdf, i.descripcion, x=38*mm, y=y_actual, max_width=maxw,
                    font="Helvetica", font_size=10, bold=True, italic=True, leading=12
                )
        else:
            pdf.setFont("Helvetica-Oblique", 10)
            pdf.drawString(38*mm, y_actual, "-")
            y_actual -= 6*mm
        y_actual -= 4*mm

    # ----------- COMENTARIOS -----------
    pdf.setFont("Helvetica-Bold", 11)
    pdf.drawString(32*mm, y_actual, "Comentario del Docente:")
    y_actual -= 6*mm
    y_actual = draw_paragraph(
        pdf, reporte.comentario or "-", x=38*mm, y=y_actual, max_width=maxw,
        font="Helvetica", font_size=10, italic=True, leading=12
    )
    y_actual -= 7*mm

    if hasattr(reporte, 'comentario_coordinador') and reporte.comentario_coordinador:
        pdf.setFont("Helvetica-Bold", 11)
        pdf.drawString(32*mm, y_actual, "Comentario del Coordinador:")
        y_actual -= 6*mm
        y_actual = draw_paragraph(
            pdf, reporte.comentario_coordinador, x=38*mm, y=y_actual, max_width=maxw,
            font="Helvetica", font_size=10, italic=True, leading=12
        )
        y_actual -= 6*mm

    # <--- hecho por claude code: firma dinámica a 5mm (~15px) del último contenido,
    # <--- hecho por claude code: firmas fijas hacia abajo (28mm del borde),
    # líneas más largas y "Firma Pariente:" en lugar de "Firma Padre/Madre:".
    y_firma = min(y_actual - 15*mm, 28*mm)
    largo_firma = 58 * mm
    x_col1 = 15 * mm
    x_col2 = 78 * mm
    x_col3 = 141 * mm

    pdf.setStrokeColor(colors.black)
    pdf.setLineWidth(0.7)

    for x_col, label, nombre in [
        (x_col1, "Firma:", reporte.docente or ""),
        (x_col2, "Firma:", getattr(reporte, 'coordinador_firma', '') or ""),
        (x_col3, "Firma Pariente:", ""),
    ]:
        pdf.setFont("Helvetica-Bold", 10)
        pdf.drawString(x_col, y_firma + 5*mm, label)
        ancho_lbl = pdf.stringWidth(label, "Helvetica-Bold", 10)
        pdf.setFont("Helvetica", 10)
        pdf.drawString(x_col + ancho_lbl + 2, y_firma + 5*mm, nombre)
        pdf.line(x_col, y_firma, x_col + largo_firma, y_firma)

    pdf.save()
    buf.seek(0)
    grado_slug = reporte.grado.replace(' ', '_').replace('/', '-')
    nombre_archivo = f"conductual_{reporte.alumno_nombre.replace(' ', '_')}_{grado_slug}.pdf"
    response = HttpResponse(buf, content_type="application/pdf")
    response['Content-Disposition'] = f'inline; filename="{nombre_archivo}"'
    return response


@login_required
def descargar_pdf_progress(request, pk):
    from pptx.oxml.ns import qn
    from lxml import etree
    from pptx.util import Emu

    # ── helpers ──────────────────────────────────────────────────────────────
    def _add_shadow(run):
        rPr = run._r.get_or_add_rPr()
        for ex in rPr.findall(qn('a:effectLst')):
            rPr.remove(ex)
        effectLst = etree.SubElement(rPr, qn('a:effectLst'))
        outer = etree.SubElement(effectLst, qn('a:outerShdw'),
                                 attrib={'blurRad': '40000', 'dist': '23000',
                                         'dir': '5400000', 'rotWithShape': '0'})
        clr = etree.SubElement(outer, qn('a:srgbClr'), val='000000')
        etree.SubElement(clr, qn('a:alpha'), val='40000')

    def _run_titulo(para, texto):
        run = para.add_run()
        run.text = texto
        run.font.name = 'Book Antiqua'
        run.font.size = Pt(40)
        run.font.bold = True
        run.font.italic = True
        run.font.underline = True
        _add_shadow(run)
        return run

    def _run_label(para, texto, size=28):
        run = para.add_run()
        run.text = texto
        run.font.name = 'Book Antiqua'
        run.font.size = Pt(size)
        run.font.bold = True
        run.font.italic = True
        run.font.underline = True
        _add_shadow(run)
        return run

    def _run_value(para, texto, size=28):
        run = para.add_run()
        run.text = texto
        run.font.name = 'Book Antiqua'
        run.font.size = Pt(size)
        run.font.italic = True
        _add_shadow(run)
        return run

    def _clear_table_style(tbl_obj):
        # Elimina el estilo de tabla predeterminado (causa el fondo gris y quita los bordes)
        tblPr = tbl_obj._tbl.find(qn('a:tblPr'))
        if tblPr is None:
            tblPr = etree.SubElement(tbl_obj._tbl, qn('a:tblPr'))
        for attr in ('bandRow', 'bandCol', 'firstRow', 'firstCol', 'lastRow', 'lastCol'):
            tblPr.set(attr, '0')
        for el in tblPr.findall(qn('a:tableStyleId')):
            tblPr.remove(el)
        # GUID "No Style, No Grid" — tabla limpia sin tema
        etree.SubElement(tblPr, qn('a:tableStyleId')).text = '{2D5ABB26-0587-4C30-8999-92F81FD0307C}'

    def _set_borders(cell, w=19050):
        # Los bordes deben ir ANTES del fill en tcPr (orden OOXML)
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        for tag in ('lnL', 'lnR', 'lnT', 'lnB'):
            for ex in tcPr.findall(qn(f'a:{tag}')):
                tcPr.remove(ex)
        for i, tag in enumerate(('lnL', 'lnR', 'lnT', 'lnB')):
            ln = etree.Element(qn(f'a:{tag}'), attrib={'w': str(w), 'cap': 'flat', 'cmpd': 'sng'})
            sf = etree.SubElement(ln, qn('a:solidFill'))
            etree.SubElement(sf, qn('a:srgbClr'), val='000000')
            tcPr.insert(i, ln)   # insertar antes del fill

    def _ordinal(d):
        if 11 <= d <= 13:
            return f"{d}th"
        return f"{d}{['th','st','nd','rd','th'][min(d % 10, 4)]}"

    def _clean_grado(grado):
        import re
        _map = {'1':'1st','2':'2nd','3':'3rd','4':'4th',
                '5':'5th','6':'6th','7':'7th','8':'8th','9':'9th'}
        m = re.search(r'(\d+)', grado or '')
        if m:
            return _map.get(m.group(1), grado)
        if 'kinder' in (grado or '').lower():
            return 'Kinder'
        return grado

    # ── datos ─────────────────────────────────────────────────────────────────
    reporte  = get_object_or_404(ProgressReport, pk=pk)
    materias = reporte.materias_json or []
    fondo_path = '/home/admin2/techcare_project/system_proyect/conducta/static/conducta/img/plantilla.jpg'

    # ── slide widescreen 13.33" × 7.5" ───────────────────────────────────────
    prs = Presentation()
    prs.slide_width  = Emu(12192000)
    prs.slide_height = Emu(6858000)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(fondo_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

    # Constantes de layout
    LEFT_X = Inches(0.2)
    TAB_X  = Inches(4.7)
    LEFT_W = TAB_X - LEFT_X          # nombre envuelve exactamente al borde de la tabla
    TAB_W  = Inches(8.5)             # tabla más ancha; total 4.7+8.5=13.2" ≈ slide

    # ── TÍTULO (encima de la tabla, derecha) ──────────────────────────────────
    tb = slide.shapes.add_textbox(TAB_X, Inches(0.05), TAB_W, Inches(0.9))
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _run_titulo(p, 'Progress Report.')

    # ── WEEKS ─────────────────────────────────────────────────────────────────
    semana_str = (f"{_ordinal(reporte.semana_inicio.day)} to "
                  f"{_ordinal(reporte.semana_fin.day)} {reporte.semana_fin.strftime('%B')}")
    tb = slide.shapes.add_textbox(LEFT_X, Inches(0.1), LEFT_W, Inches(1.5))
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    _run_label(p, 'Weeks:')
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.LEFT
    _run_value(p2, semana_str)

    # ── GRADE ─────────────────────────────────────────────────────────────────
    tb = slide.shapes.add_textbox(LEFT_X, Inches(1.75), LEFT_W, Inches(0.75))
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    _run_value(p, f'Grade: {_clean_grado(reporte.grado)}')

    # ── NAME label + nombre alumno ────────────────────────────────────────────
    tb = slide.shapes.add_textbox(LEFT_X, Inches(2.6), LEFT_W, Inches(1.4))
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    _run_label(p, 'Name:')
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.LEFT
    rn = p2.add_run()
    rn.text = reporte.alumno_nombre
    rn.font.name = 'Book Antiqua'
    rn.font.size = Pt(20)
    rn.font.italic = True

    # ── TABLA ─────────────────────────────────────────────────────────────────
    num_rows = len(materias) + 1
    tbl_top  = Inches(0.95)
    tbl_h    = Inches(6.55)   # ocupa casi toda la altura

    tbl = slide.shapes.add_table(num_rows, 3, TAB_X, tbl_top, TAB_W, tbl_h).table
    _clear_table_style(tbl)

    # anchos de columna: ~24% / ~19% / ~57%
    tbl.columns[0].width = Inches(2.0)
    tbl.columns[1].width = Inches(1.7)
    tbl.columns[2].width = Inches(4.8)

    def _cell_style(cell, texto, bold=False, align=PP_ALIGN.CENTER):
        cell.text = texto
        _set_borders(cell)         # bordes ANTES del fill (orden correcto en OOXML)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
        para = cell.text_frame.paragraphs[0]
        para.font.name  = 'Arial'
        para.font.size  = Pt(18)
        para.font.bold  = bold
        para.font.color.rgb = RGBColor(0, 0, 0)
        para.alignment  = align

    # cabecera
    for idx, hdr in enumerate(['Materia', 'Asignación', 'Comentario/Observación']):
        _cell_style(tbl.cell(0, idx), hdr, bold=True, align=PP_ALIGN.CENTER)

    # filas de materias
    for row, mat in enumerate(materias, start=1):
        _cell_style(tbl.cell(row, 0), mat.get('materia', ''),    bold=True,  align=PP_ALIGN.CENTER)
        _cell_style(tbl.cell(row, 1), mat.get('asignacion', ''), bold=False, align=PP_ALIGN.LEFT)
        _cell_style(tbl.cell(row, 2), mat.get('comentario', ''), bold=False, align=PP_ALIGN.LEFT)

    # ── descarga ──────────────────────────────────────────────────────────────
    response = HttpResponse(
        content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
    filename = f'progress_{reporte.alumno_nombre.replace(" ", "_")}.pptx'
    response['Content-Disposition'] = f'attachment; filename="{filename}"'
    prs.save(response)
    return response


@login_required
def descargar_zip_reportes(request):
    """Genera un ZIP con todos los reportes (PDF + PPTX) para coordinadores y admin."""
    import zipfile
    from reportlab.lib.pagesizes import letter as rl_letter
    from reportlab.pdfgen import canvas as rl_canvas
    from reportlab.lib.units import mm as rl_mm
    from reportlab.lib import colors as rl_colors
    from pptx import Presentation as PPTXPres
    from pptx.util import Inches as PPTXInches, Pt as PPTXPt, Emu as PPTXEmu
    from pptx.enum.text import PP_ALIGN as PPTXAlign
    from pptx.dml.color import RGBColor as PPTXColor
    from pptx.oxml.ns import qn as pptx_qn
    from lxml import etree as pptx_etree

    # Solo coordinadores (BL y Colegio) y superusuario
    from accounts.panel_roles import _GRUPOS_COORD_BL, _GRUPOS_COORD_COL
    es_coord = request.user.groups.filter(
        name__in=(*_GRUPOS_COORD_BL, *_GRUPOS_COORD_COL)
    ).exists()
    if not (request.user.is_superuser or es_coord):
        return HttpResponseForbidden()

    logo_path = os.path.join(settings.STATIC_ROOT, 'conducta/img/ana-transformed.png')
    fondo_path = '/home/admin2/techcare_project/system_proyect/conducta/static/conducta/img/plantilla.jpg'

    # ── Logo cacheado y reducido (se incrusta en CADA PDF; el original pesa ~1.8MB).
    # Se carga una sola vez y se reduce, evitando releer/recodificar por reporte. ──
    from reportlab.lib.utils import ImageReader as _RLImageReader
    _logo_reader = None
    if os.path.exists(logo_path):
        try:
            from PIL import Image as _PILImage
            _li = _PILImage.open(logo_path).convert('RGBA')
            _li.thumbnail((400, 400))
            _lbuf = io.BytesIO(); _li.save(_lbuf, format='PNG'); _lbuf.seek(0)
            _logo_reader = _RLImageReader(_lbuf)
        except Exception:
            _logo_reader = _RLImageReader(logo_path)

    # ── generador PDF informativo ─────────────────────────────────────────────
    def _pdf_informativo(r):
        buf = io.BytesIO()
        w, h = rl_letter
        pdf = rl_canvas.Canvas(buf, pagesize=rl_letter)
        if _logo_reader:
            pdf.drawImage(_logo_reader, x=(w-35*rl_mm)/2, y=h-40*rl_mm,
                          width=35*rl_mm, height=35*rl_mm, mask='auto')
        y = h - 52*rl_mm
        pdf.setFont('Helvetica-Bold', 16)
        pdf.drawCentredString(w/2, y, 'Nuevo Amanecer School' if r.area == 'bilingue' else 'C.E.M.N.G Nuevo Amanecer')
        y -= 14*rl_mm
        pdf.setFont('Helvetica', 13)
        pdf.drawCentredString(w/2, y, 'Reporte informativo')
        y -= 14*rl_mm

        def lv(x, yy, lbl, val, ls=11, vs=10):
            pdf.setFont('Helvetica-Bold', ls); pdf.drawString(x, yy, lbl)
            aw = pdf.stringWidth(lbl, 'Helvetica-Bold', ls)
            pdf.setFont('Helvetica', vs); pdf.drawString(x+aw+2, yy, val)

        lv(32*rl_mm, y, 'Nombre:', r.alumno_nombre)
        lv(118*rl_mm, y, 'Grado:', r.grado); y -= 8*rl_mm
        lv(32*rl_mm, y, 'Docente:', r.docente or '')
        lv(100*rl_mm, y, 'Fecha:', r.fecha.strftime('%d/%m/%Y') if r.fecha else ''); y -= 12*rl_mm
        pdf.setFont('Helvetica-Bold', 11); pdf.drawString(32*rl_mm, y, 'Comentario del Docente:'); y -= 8*rl_mm
        y = draw_paragraph(pdf, r.comentario or '-', x=38*rl_mm, y=y,
                            max_width=w-65*rl_mm, font='Helvetica', font_size=10, italic=True, leading=12)
        y -= 10*rl_mm
        if getattr(r, 'comentario_coordinador', None):
            pdf.setFont('Helvetica-Bold', 11); pdf.drawString(32*rl_mm, y, 'Comentario del Coordinador:'); y -= 8*rl_mm
            y = draw_paragraph(pdf, r.comentario_coordinador, x=38*rl_mm, y=y,
                                max_width=w-65*rl_mm, font='Helvetica', font_size=10, italic=True, leading=12)
            y -= 8*rl_mm
        y_f = min(y-15*rl_mm, 28*rl_mm); lf = 58*rl_mm
        pdf.setStrokeColor(rl_colors.black); pdf.setLineWidth(0.7)
        for xc, lbl, nom in [(15*rl_mm,'Firma:',r.docente or ''),
                              (78*rl_mm,'Firma:',getattr(r,'coordinador_firma','')),
                              (141*rl_mm,'Firma Pariente:','')]:
            pdf.setFont('Helvetica-Bold',10); pdf.drawString(xc, y_f+5*rl_mm, lbl)
            aw = pdf.stringWidth(lbl,'Helvetica-Bold',10)
            pdf.setFont('Helvetica',10); pdf.drawString(xc+aw+2, y_f+5*rl_mm, nom or '')
            pdf.line(xc, y_f, xc+lf, y_f)
        pdf.save(); buf.seek(0); return buf

    # ── generador PDF conductual ──────────────────────────────────────────────
    def _pdf_conductual(r):
        buf = io.BytesIO()
        w, h = rl_letter
        pdf = rl_canvas.Canvas(buf, pagesize=rl_letter)
        if _logo_reader:
            pdf.drawImage(_logo_reader, x=(w-22*rl_mm)/2, y=h-28*rl_mm,
                          width=22*rl_mm, height=22*rl_mm, mask='auto')
        y = h - 33*rl_mm
        pdf.setFont('Helvetica-Bold', 16)
        pdf.drawCentredString(w/2, y, 'Nuevo Amanecer School' if getattr(r,'area',None)=='bilingue' else 'C.E.M.N.G Nuevo Amanecer')
        y -= 9*rl_mm; pdf.setFont('Helvetica',13)
        pdf.drawCentredString(w/2, y, 'Reporte conductual'); y -= 10*rl_mm

        def lv(x, yy, lbl, val, s=11):
            pdf.setFont('Helvetica-Bold',s); pdf.drawString(x,yy,lbl)
            aw=pdf.stringWidth(lbl,'Helvetica-Bold',s)
            pdf.setFont('Helvetica',s); pdf.drawString(x+aw+2,yy,val)

        lv(32*rl_mm,y,'Nombre:',r.alumno_nombre); y-=7*rl_mm
        lv(32*rl_mm,y,'Grado:',r.grado)
        lv(118*rl_mm,y,'Fecha:',r.fecha.strftime('%d/%m/%Y') if r.fecha else ''); y-=7*rl_mm
        lv(32*rl_mm,y,'Docente:',r.docente or ''); y-=9*rl_mm
        mw = w-65*rl_mm
        for tipo, lbl in [('incisos_leve','Leve'),('incisos_grave','Grave'),('incisos_muygrave','Muy Grave')]:
            incs = getattr(r,tipo).all()
            pdf.setFont('Helvetica-Bold',11); pdf.drawString(32*rl_mm,y,f'Incisos {lbl}:'); y-=6*rl_mm
            if incs:
                for i in incs:
                    y=draw_paragraph(pdf,i.descripcion,x=38*rl_mm,y=y,max_width=mw,
                                     font='Helvetica',font_size=10,bold=True,italic=True,leading=12)
            else:
                pdf.setFont('Helvetica-Oblique',10); pdf.drawString(38*rl_mm,y,'-'); y-=6*rl_mm
            y-=4*rl_mm
        pdf.setFont('Helvetica-Bold',11); pdf.drawString(32*rl_mm,y,'Comentario del Docente:'); y-=6*rl_mm
        y=draw_paragraph(pdf,r.comentario or '-',x=38*rl_mm,y=y,max_width=mw,
                          font='Helvetica',font_size=10,italic=True,leading=12); y-=7*rl_mm
        if getattr(r,'comentario_coordinador',None):
            pdf.setFont('Helvetica-Bold',11); pdf.drawString(32*rl_mm,y,'Comentario del Coordinador:'); y-=6*rl_mm
            y=draw_paragraph(pdf,r.comentario_coordinador,x=38*rl_mm,y=y,max_width=mw,
                              font='Helvetica',font_size=10,italic=True,leading=12); y-=6*rl_mm
        y_f=min(y-15*rl_mm,28*rl_mm); lf=58*rl_mm
        pdf.setStrokeColor(rl_colors.black); pdf.setLineWidth(0.7)
        for xc,lbl,nom in [(15*rl_mm,'Firma:',r.docente or ''),
                            (78*rl_mm,'Firma:',getattr(r,'coordinador_firma','') or ''),
                            (141*rl_mm,'Firma Pariente:','')]:
            pdf.setFont('Helvetica-Bold',10); pdf.drawString(xc,y_f+5*rl_mm,lbl)
            aw=pdf.stringWidth(lbl,'Helvetica-Bold',10)
            pdf.setFont('Helvetica',10); pdf.drawString(xc+aw+2,y_f+5*rl_mm,nom)
            pdf.line(xc,y_f,xc+lf,y_f)
        pdf.save(); buf.seek(0); return buf

    # ── generador PPTX progress ───────────────────────────────────────────────
    def _pptx_progress(r):
        def _shadow(run):
            rPr=run._r.get_or_add_rPr()
            for ex in rPr.findall(pptx_qn('a:effectLst')): rPr.remove(ex)
            ef=pptx_etree.SubElement(rPr,pptx_qn('a:effectLst'))
            os_=pptx_etree.SubElement(ef,pptx_qn('a:outerShdw'),attrib={'blurRad':'40000','dist':'23000','dir':'5400000','rotWithShape':'0'})
            cl=pptx_etree.SubElement(os_,pptx_qn('a:srgbClr'),val='000000')
            pptx_etree.SubElement(cl,pptx_qn('a:alpha'),val='40000')
        def _run(para,txt,size,bold=False,italic=False,underline=False,shadow=False):
            run=para.add_run(); run.text=txt
            run.font.name='Book Antiqua'; run.font.size=PPTXPt(size)
            run.font.bold=bold; run.font.italic=italic; run.font.underline=underline
            if shadow: _shadow(run)
            return run
        def _borders(cell,w=19050):
            tc=cell._tc; tcPr=tc.get_or_add_tcPr()
            for tag in ('lnL','lnR','lnT','lnB'):
                for ex in tcPr.findall(pptx_qn(f'a:{tag}')): tcPr.remove(ex)
            for i,tag in enumerate(('lnL','lnR','lnT','lnB')):
                ln=pptx_etree.Element(pptx_qn(f'a:{tag}'),attrib={'w':str(w),'cap':'flat','cmpd':'sng'})
                sf=pptx_etree.SubElement(ln,pptx_qn('a:solidFill'))
                pptx_etree.SubElement(sf,pptx_qn('a:srgbClr'),val='000000')
                tcPr.insert(i,ln)
        def _ordinal(d):
            if 11<=d<=13: return f'{d}th'
            return f'{d}{["th","st","nd","rd","th"][min(d%10,4)]}'
        def _clean_g(g):
            import re; m=re.search(r'(\d+)',g or '')
            mp={'1':'1st','2':'2nd','3':'3rd','4':'4th','5':'5th','6':'6th','7':'7th','8':'8th','9':'9th'}
            if m: return mp.get(m.group(1),g)
            return 'Kinder' if 'kinder' in (g or '').lower() else g

        prs=PPTXPres(); prs.slide_width=PPTXEmu(12192000); prs.slide_height=PPTXEmu(6858000)
        slide=prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(fondo_path,0,0,width=prs.slide_width,height=prs.slide_height)
        LX=PPTXInches(0.2); TX=PPTXInches(4.7); LW=TX-LX; TW=PPTXInches(8.5)
        # título
        tb=slide.shapes.add_textbox(TX,PPTXInches(0.05),TW,PPTXInches(0.9))
        tf=tb.text_frame; tf.clear(); p=tf.paragraphs[0]; p.alignment=PPTXAlign.CENTER
        _run(p,'Progress Report.',40,bold=True,italic=True,underline=True,shadow=True)
        # weeks
        sem=f"{_ordinal(r.semana_inicio.day)} to {_ordinal(r.semana_fin.day)} {r.semana_fin.strftime('%B')}"
        tb=slide.shapes.add_textbox(LX,PPTXInches(0.1),LW,PPTXInches(1.5))
        tf=tb.text_frame; tf.clear(); tf.word_wrap=True
        p=tf.paragraphs[0]; p.alignment=PPTXAlign.LEFT
        _run(p,'Weeks:',28,bold=True,italic=True,underline=True,shadow=True)
        p2=tf.add_paragraph(); p2.alignment=PPTXAlign.LEFT
        _run(p2,sem,28,italic=True,shadow=True)
        # grade
        tb=slide.shapes.add_textbox(LX,PPTXInches(1.75),LW,PPTXInches(0.75))
        tf=tb.text_frame; tf.clear(); tf.word_wrap=True
        p=tf.paragraphs[0]; p.alignment=PPTXAlign.LEFT
        _run(p,f'Grade: {_clean_g(r.grado)}',28,italic=True,shadow=True)
        # name
        tb=slide.shapes.add_textbox(LX,PPTXInches(2.6),LW,PPTXInches(1.4))
        tf=tb.text_frame; tf.clear(); tf.word_wrap=True
        p=tf.paragraphs[0]; p.alignment=PPTXAlign.LEFT
        _run(p,'Name:',28,bold=True,italic=True,underline=True,shadow=True)
        p2=tf.add_paragraph(); p2.alignment=PPTXAlign.LEFT
        _run(p2,r.alumno_nombre,20,italic=True)
        # tabla
        mats=r.materias_json or []
        tbl=slide.shapes.add_table(len(mats)+1,3,TX,PPTXInches(0.95),TW,PPTXInches(6.55)).table
        tblPr=tbl._tbl.find(pptx_qn('a:tblPr'))
        if tblPr is None: tblPr=pptx_etree.SubElement(tbl._tbl,pptx_qn('a:tblPr'))
        for at in ('bandRow','bandCol','firstRow','firstCol','lastRow','lastCol'): tblPr.set(at,'0')
        for el in tblPr.findall(pptx_qn('a:tableStyleId')): tblPr.remove(el)
        pptx_etree.SubElement(tblPr,pptx_qn('a:tableStyleId')).text='{2D5ABB26-0587-4C30-8999-92F81FD0307C}'
        tbl.columns[0].width=PPTXInches(2.0); tbl.columns[1].width=PPTXInches(1.7); tbl.columns[2].width=PPTXInches(4.8)
        def _cs(cell,txt,bold=False,align=PPTXAlign.CENTER):
            cell.text=txt; _borders(cell)
            cell.fill.solid(); cell.fill.fore_color.rgb=PPTXColor(255,255,255)
            p=cell.text_frame.paragraphs[0]; p.font.name='Arial'
            p.font.size=PPTXPt(18); p.font.bold=bold
            p.font.color.rgb=PPTXColor(0,0,0); p.alignment=align
        for i,h_ in enumerate(['Materia','Asignación','Comentario/Observación']):
            _cs(tbl.cell(0,i),h_,bold=True)
        for row,mat in enumerate(mats,start=1):
            _cs(tbl.cell(row,0),mat.get('materia',''),bold=True,align=PPTXAlign.CENTER)
            _cs(tbl.cell(row,1),mat.get('asignacion',''),bold=False,align=PPTXAlign.LEFT)
            _cs(tbl.cell(row,2),mat.get('comentario',''),bold=False,align=PPTXAlign.LEFT)
        buf=io.BytesIO(); prs.save(buf); buf.seek(0); return buf

    # ── determinar scope según coordinador ───────────────────────────────────
    email = request.user.email
    if request.user.is_superuser:
        qs_info = ReporteInformativo.objects.all()
        qs_cond = ReporteConductual.objects.all()
        qs_prog = ProgressReport.objects.all()
    elif email == 'cvarela@ana-hn.org':
        _q = _q_for_coord('C1')
        qs_info = ReporteInformativo.objects.filter(area='bilingue', tipo_reporte='academico').filter(_q)
        qs_cond = ReporteConductual.objects.filter(area='bilingue').filter(_q)
        qs_prog = ProgressReport.objects.none()
    elif email == 'druiz@ana-hn.org':
        _q = _q_for_coord('C2')
        qs_info = ReporteInformativo.objects.filter(area='bilingue', tipo_reporte='academico').filter(_q)
        qs_cond = ReporteConductual.objects.filter(area='bilingue').filter(_q)
        qs_prog = ProgressReport.objects.none()
    elif email == 'ialcerro@ana-hn.org':
        _q = _q_for_coord('C3')
        qs_info = (
            ReporteInformativo.objects.filter(area='bilingue', tipo_reporte='conductual') |
            ReporteInformativo.objects.filter(area='bilingue', tipo_reporte='academico').filter(_q)
        ).distinct().order_by('-fecha')
        qs_cond = ReporteConductual.objects.filter(area='bilingue').filter(_q)
        qs_prog = ProgressReport.objects.none()
    elif email == 'jmartinez@ana-hn.org':
        _q = _q_for_coord('C4')
        qs_info = ReporteInformativo.objects.filter(area='bilingue', tipo_reporte='academico').filter(_q)
        qs_cond = ReporteConductual.objects.filter(area='bilingue').filter(_q)
        qs_prog = ProgressReport.objects.all()
    elif email == 'coordinacion_bl@ana-hn.org':
        qs_info = ReporteInformativo.objects.none()
        qs_cond = ReporteConductual.objects.none()
        qs_prog = ProgressReport.objects.all()
    elif request.user.groups.filter(name__in=_GRUPOS_COORD_COL).exists():
        qs_info = ReporteInformativo.objects.filter(area='colegio')
        qs_cond = ReporteConductual.objects.filter(area='colegio')
        qs_prog = ProgressReport.objects.none()
    else:
        qs_info = ReporteInformativo.objects.filter(area='bilingue')
        qs_cond = ReporteConductual.objects.filter(area='bilingue')
        qs_prog = ProgressReport.objects.all()

    # ── filtros opcionales (docente para todos · coordinador solo superuser) ──
    f_docente = (request.GET.get('docente') or '').strip()
    if f_docente:
        qs_info = qs_info.filter(docente__icontains=f_docente)
        qs_cond = qs_cond.filter(docente__icontains=f_docente)
    f_coord = (request.GET.get('coord') or '').strip().upper()
    if f_coord in {'C1', 'C2', 'C3', 'C4'} and request.user.is_superuser:
        _qc = _q_for_coord(f_coord)
        qs_info = qs_info.filter(_qc)
        qs_cond = qs_cond.filter(_qc)

    # ── construir ZIP ─────────────────────────────────────────────────────────
    mem = io.BytesIO()
    fecha = timezone.now().strftime('%Y%m%d_%H%M')
    with zipfile.ZipFile(mem, 'w', zipfile.ZIP_DEFLATED) as zf:
        for r in qs_info.order_by('area', 'alumno_nombre'):
            try:
                nombre = r.alumno_nombre.replace(' ', '_')
                grado  = r.grado.replace(' ', '_').replace('/', '-')
                fname  = f'informativo_{nombre}_{grado}.pdf'
                zf.writestr(f'Informativos/{r.area.capitalize()}/{fname}',
                            _pdf_informativo(r).getvalue())
            except Exception:
                pass
        for r in qs_cond.prefetch_related(
                'incisos_leve', 'incisos_grave', 'incisos_muygrave'
        ).order_by('area', 'alumno_nombre'):
            try:
                nombre = r.alumno_nombre.replace(' ', '_')
                grado  = r.grado.replace(' ', '_').replace('/', '-')
                fname  = f'conductual_{nombre}_{grado}.pdf'
                zf.writestr(f'Conductuales/{r.area.capitalize()}/{fname}',
                            _pdf_conductual(r).getvalue())
            except Exception:
                pass
        for r in qs_prog.order_by('alumno_nombre'):
            try:
                nombre = r.alumno_nombre.replace(' ', '_')
                grado  = r.grado.replace(' ', '_').replace('/', '-')
                fname  = f'progress_{nombre}_{grado}.pptx'
                zf.writestr(f'ProgressReports/{fname}',
                            _pptx_progress(r).getvalue())
            except Exception:
                pass

    mem.seek(0)
    response = HttpResponse(mem.getvalue(), content_type='application/zip')
    response['Content-Disposition'] = f'attachment; filename="reportes_{fecha}.zip"'
    return response


@login_required
def descargar_pdf_conductual_3_strikes(request, pk):
    from .models import ReporteConductual
    reporte = get_object_or_404(ReporteConductual, pk=pk)

    # <--- hecho por claude code: regla corregida — el reporte 3 strikes solo se genera
    # cuando el alumno acumula 3 reportes del MISMO tipo (3 leve, 3 grave o 3 muy grave).
    # Antes tomaba cualquier 3 reportes mezclados, lo cual era incorrecto.
    todos = list(ReporteConductual.objects.filter(
        area=reporte.area,
        alumno_id=reporte.alumno_id
    ).prefetch_related('incisos_leve', 'incisos_grave', 'incisos_muygrave').order_by('fecha'))

    leves     = [r for r in todos if r.incisos_leve.exists() and not r.incisos_grave.exists() and not r.incisos_muygrave.exists()]
    graves    = [r for r in todos if r.incisos_grave.exists() and not r.incisos_muygrave.exists()]
    muy_graves = [r for r in todos if r.incisos_muygrave.exists()]

    reportes = None
    if len(leves) >= 3:
        reportes = leves[:3]
    elif len(graves) >= 3:
        reportes = graves[:3]
    elif len(muy_graves) >= 3:
        reportes = muy_graves[:3]

    if reportes is None:
        # <--- hecho por claude code: mensaje de error como alerta Bootstrap en lugar
        # de texto plano (que tenía encoding roto y se veía mal).
        html = """<!DOCTYPE html>
<html lang="es">
<head>
  <meta charset="UTF-8">
  <title>Sin reporte 3 strikes</title>
  <link href="https://cdn.jsdelivr.net/npm/bootstrap@5.3.0/dist/css/bootstrap.min.css" rel="stylesheet">
</head>
<body class="bg-light d-flex align-items-center justify-content-center" style="min-height:100vh;">
  <div class="card shadow p-4" style="max-width:480px;width:100%;">
    <div class="alert alert-warning mb-3">
      <strong>&#9888; Sin reporte 3 strikes</strong><br>
      Este alumno aún no acumula 3 reportes conductuales del mismo tipo
      (3 leve, 3 grave o 3 muy grave).
    </div>
    <a href="javascript:history.back()" class="btn btn-secondary w-100">&larr; Volver</a>
  </div>
</body>
</html>"""
        return HttpResponse(html, status=400)

    buf = io.BytesIO()
    w, h = letter
    pdf = canvas.Canvas(buf, pagesize=letter)
    pdf.setTitle("reporte_conductual_3_strikes.pdf")

    # Logo centrado
    width_logo = 35 * mm
    height_logo = 35 * mm
    x_logo = (w - width_logo) / 2
    logo_path = os.path.join(settings.STATIC_ROOT, "conducta/img/ana-transformed.png")
    if os.path.exists(logo_path):
        pdf.drawImage(
            logo_path,
            x=x_logo,
            y=h-40*mm,
            width=width_logo,
            height=height_logo,
            mask='auto'
        )
    y_actual = h - 47 * mm

    # Título
    pdf.setFont("Helvetica-Bold", 16)
    titulo = "Nuevo Amanecer School" if reporte.area == "bilingue" else "C.E.M.N.G Nuevo Amanecer"
    pdf.drawCentredString(w/2, y_actual, titulo)
    y_actual -= 12*mm
    pdf.setFont("Helvetica", 13)
    pdf.drawCentredString(w/2, y_actual, "Reporte conductual – 3 Strikes")
    y_actual -= 10*mm

    # --- CADA REPORTE ---
    for idx, rep in enumerate(reportes, 1):
        pdf.setFont("Helvetica-Bold", 12)
        pdf.drawString(32*mm, y_actual, f"Reporte #{idx} ({', '.join(['Leve' if rep.incisos_leve.exists() else '', 'Grave' if rep.incisos_grave.exists() else '', 'Muy Grave' if rep.incisos_muygrave.exists() else '']).replace(',,',',').strip(', ')})")
        y_actual -= 8*mm

        # Fila: Nombre/Grado/Docente/Fecha
        pdf.setFont("Helvetica-Bold", 11)
        pdf.drawString(32*mm, y_actual, "Nombre:")
        pdf.setFont("Helvetica", 11)
        pdf.drawString(60*mm, y_actual, rep.alumno_nombre)

        pdf.setFont("Helvetica-Bold", 11)
        pdf.drawString(120*mm, y_actual, "Grado:")
        pdf.setFont("Helvetica", 11)
        pdf.drawString(140*mm, y_actual, rep.grado)
        y_actual -= 6*mm

        pdf.setFont("Helvetica-Bold", 11)
        pdf.drawString(32*mm, y_actual, "Docente/Materia:")
        pdf.setFont("Helvetica", 11)
        pdf.drawString(75*mm, y_actual, rep.docente)
        pdf.setFont("Helvetica-Bold", 11)
        pdf.drawString(120*mm, y_actual, "Fecha:")
        pdf.setFont("Helvetica", 11)
        pdf.drawString(140*mm, y_actual, rep.fecha.strftime('%d/%m/%Y'))
        y_actual -= 7*mm

        # Incisos (párrafo)
        pdf.setFont("Helvetica-Bold", 10)
        pdf.drawString(32*mm, y_actual, "Inciso:")
        y_actual -= 6*mm

        incisos = [i.descripcion for i in rep.incisos_leve.all()] + [i.descripcion for i in rep.incisos_grave.all()] + [i.descripcion for i in rep.incisos_muygrave.all()]
        if incisos:
            text_obj = pdf.beginText(38*mm, y_actual)
            text_obj.setFont("Helvetica-BoldOblique", 10)
            wrap_len = 95  # Ajusta a gusto (número de caracteres por línea)
            for inciso in incisos:
                lines = [inciso[i:i+wrap_len] for i in range(0, len(inciso), wrap_len)]
                for ln in lines:
                    text_obj.textLine(ln)
                    y_actual -= 5*mm
            pdf.drawText(text_obj)
            y_actual = text_obj.getY() - 1*mm
        else:
            pdf.setFont("Helvetica", 10)
            pdf.drawString(38*mm, y_actual, "-")
            y_actual -= 5*mm

        # Descripción (Comentario)
        pdf.setFont("Helvetica-Bold", 10)
        pdf.drawString(32*mm, y_actual, "Descripción:")
        y_actual -= 5*mm
        if rep.comentario:
            text_obj = pdf.beginText(38*mm, y_actual)
            text_obj.setFont("Helvetica-Oblique", 10)
            for line in rep.comentario.split("\n"):
                text_obj.textLine(line)
                y_actual -= 5*mm
            pdf.drawText(text_obj)
            y_actual = text_obj.getY() - 1*mm
        else:
            pdf.setFont("Helvetica", 10)
            pdf.drawString(38*mm, y_actual, "-")
            y_actual -= 5*mm

        # Espacio entre reportes
        y_actual -= 4*mm
        pdf.setStrokeColor(colors.grey)
        pdf.line(30*mm, y_actual, w-30*mm, y_actual)
        y_actual -= 8*mm

    # FIRMAS PARA LLENAR EN FÍSICO
    y_firmas = 20*mm
    x_firma = [34*mm, 91*mm, 146*mm]
    if reporte.area == "colegio":
        etiquetas = ["Firma Padre de Familia", "Firma del Docente/Orientación", "Firma de Consejería"]
    else:
        etiquetas = ["Firma Padre de Familia", "Firma del Docente", "Firma del Coordinador"]

    for i in range(3):
        pdf.setStrokeColor(colors.black)
        pdf.line(x_firma[i], y_firmas, x_firma[i]+42*mm, y_firmas)
        pdf.setFont("Helvetica", 10)
        pdf.drawCentredString(x_firma[i]+21*mm, y_firmas-5*mm, etiquetas[i])

    pdf.save()
    buf.seek(0)
    grado_slug = reporte.grado.replace(' ', '_').replace('/', '-')
    nombre_archivo = f"3strikes_{reporte.alumno_nombre.replace(' ', '_')}_{grado_slug}.pdf"
    response = HttpResponse(buf, content_type="application/pdf")
    response['Content-Disposition'] = f'inline; filename="{nombre_archivo}"'
    return response

_REDIRECT_COORD_BL = {
    'cvarela@ana-hn.org':         'dashboard_c1',
    'druiz@ana-hn.org':           'dashboard_c2',
    'ialcerro@ana-hn.org':        'dashboard_c3',
    'jmartinez@ana-hn.org':       'dashboard_c4',
    'coordinacion_bl@ana-hn.org': 'dashboard_coordi_bl',
}

#--------------  DASHBOARD COORDINADOR -----------------
@login_required
def dashboard_coordinador(request, area):
    if request.user.groups.filter(name='solo_progress').exists():
        return redirect('progress_report_bilingue')

    if area == 'bilingue':
        destino = _REDIRECT_COORD_BL.get(request.user.email)
        if destino:
            return redirect(destino)

        # coordi_bl: solo ve Progress Reports
        if request.user.email == _EMAIL_COORDI_BL:
            contexto = {
                'area': area,
                'reportes_informativo': ReporteInformativo.objects.none(),
                'reportes_conductual':  ReporteConductual.objects.none(),
                'reportes_progress':    _solo_periodo_activo(ProgressReport.objects.all(), campo='semana_inicio', area='bilingue').annotate(num_evid=Count('evidenciareporte')),
                'strikes':              {},
                'today':                timezone.now().strftime('%Y-%m-%d'),
            }
            return render(request, 'conducta/dashboard_coordinador.html', contexto)

        reportes_informativo = ReporteInformativo.objects.filter(area='bilingue')
        reportes_conductual  = ReporteConductual.objects.filter(area='bilingue')
        reportes_progress    = ProgressReport.objects.all()

    elif area == 'colegio':
        reportes_informativo = ReporteInformativo.objects.filter(area='colegio')
        reportes_conductual  = ReporteConductual.objects.filter(area='colegio')
        reportes_progress    = []
    else:
        return redirect('menu')

    # Período activo + anotación de evidencias (rendimiento)
    reportes_informativo, reportes_conductual, reportes_progress = _dash_querysets(
        reportes_informativo, reportes_conductual, reportes_progress, area)

    strikes = {}
    qs = reportes_conductual.values('alumno_id').annotate(total=Count('id')).filter(total__gte=3)
    for row in qs:
        strikes[row['alumno_id']] = row['total']

    coord_nombres = {
        cfg.codigo: cfg.nombre
        for cfg in ConfiguracionCoordinador.objects.filter(area='bilingue', activo=True)
    }

    # ── Permisos editar/eliminar para coordinadores (CoordPermiso) ────────────
    u = request.user
    if u.is_superuser:
        can_edit_dash   = True
        can_delete_dash = True
    else:
        try:
            _cp = u.coord_permiso
            can_edit_dash   = _cp.dashboard_editar
            can_delete_dash = _cp.dashboard_eliminar
        except Exception:
            can_edit_dash   = False
            can_delete_dash = False

    contexto = {
        'area': area,
        'reportes_informativo': reportes_informativo,
        'reportes_conductual':  reportes_conductual,
        'reportes_progress':    reportes_progress,
        'strikes':              strikes,
        'today':                timezone.now().strftime('%Y-%m-%d'),
        'coord_nombres':        coord_nombres,
        'can_edit_dash':        can_edit_dash,
        'can_delete_dash':      can_delete_dash,
        'docentes_filtro':      _docentes_filtro(area),
        'periodo_activo':       _periodo_activo(area),
        'periodos_escolares':   _periodos_escolares(area),
    }
    return render(request, 'conducta/dashboard_coordinador.html', contexto)

def _docentes_de_coord(codigo):
    docentes = set()
    for md in MateriaDocenteBilingue.objects.filter(activo=True):
        codigos = [c.strip() for c in md.coordinador.split(',')]
        if codigo in codigos:
            docentes.add(md.docente)
    return list(docentes)

def _docentes_filtro(area='bilingue'):
    """Catálogo completo de docentes (activos) para el filtro del dashboard.
    Incluye también a los que aún no tienen reportes registrados."""
    modelo = MateriaDocenteColegio if area == 'colegio' else MateriaDocenteBilingue
    return sorted({md.docente.strip() for md in modelo.objects.filter(activo=True) if md.docente})


# ── Períodos escolares (catálogo PROPIO de conducta) ─────────────────────────
def _periodo_activo(area=None):
    """Período escolar activo (define qué reportes se ven en el dashboard)."""
    from .models import PeriodoEscolarConducta
    qs = PeriodoEscolarConducta.objects.filter(activo=True)
    if area:
        qs = qs.filter(area__in=[area, 'ambas'])
    return qs.order_by('-anio', '-parcial').first()


def _periodos_escolares(area=None):
    """Catálogo completo de períodos (para el tab Período Escolar)."""
    from .models import PeriodoEscolarConducta
    qs = PeriodoEscolarConducta.objects.all()
    if area:
        qs = qs.filter(area__in=[area, 'ambas'])
    return list(qs.order_by('-anio', '-parcial'))


def _solo_periodo_activo(qs, campo='fecha__date', area=None):
    """Limita un queryset al rango del período activo. Sin período → sin filtro."""
    p = _periodo_activo(area)
    if not p:
        return qs
    return qs.filter(**{f'{campo}__gte': p.fecha_inicio, f'{campo}__lte': p.fecha_fin})


def _dash_querysets(qs_info, qs_cond, qs_prog, area=None):
    """Optimiza los querysets del dashboard: filtra al período activo y anota
    el conteo de evidencias en UNA consulta (elimina el N+1 por fila)."""
    qs_info = (_solo_periodo_activo(qs_info, area=area)
               .annotate(num_evid=Count('evidenciareporte'))
               .prefetch_related('evidenciareporte_set'))
    qs_cond = (_solo_periodo_activo(qs_cond, area=area)
               .annotate(num_evid=Count('evidenciareporte'))
               .prefetch_related('evidenciareporte_set', 'incisos_leve', 'incisos_grave', 'incisos_muygrave'))
    if hasattr(qs_prog, 'filter'):
        qs_prog = (_solo_periodo_activo(qs_prog, campo='semana_inicio', area=area)
                   .annotate(num_evid=Count('evidenciareporte'))
                   .prefetch_related('evidenciareporte_set'))
    return qs_info, qs_cond, qs_prog


def _coord_dash_perms(user):
    """Retorna (can_edit_dash, can_delete_dash) según CoordPermiso del usuario."""
    if user.is_superuser:
        return True, True
    try:
        cp = CoordPermiso.objects.get(user_id=user.pk)
        return cp.dashboard_editar, cp.dashboard_eliminar
    except CoordPermiso.DoesNotExist:
        return False, False


def dashboard_c1(request):
    q = _q_for_coord('C1')
    qs_info = ReporteInformativo.objects.filter(area='bilingue', tipo_reporte='academico').filter(_q_academico_for_coord('C1'))
    qs_cond = ReporteConductual.objects.filter(area='bilingue').filter(q)
    qs_prog = ProgressReport.objects.all()
    qs_info, qs_cond, qs_prog = _dash_querysets(qs_info, qs_cond, qs_prog, 'bilingue')
    strikes = {r['alumno_id']: r['total'] for r in qs_cond.values('alumno_id').annotate(total=Count('id')).filter(total__gte=3)}
    reportes_nuevos_bl = (
        ReporteInformativo.objects.filter(area='bilingue', estado='enviado').count() +
        ReporteConductual.objects.filter(area='bilingue', estado='enviado').count() +
        ProgressReport.objects.filter(estado='enviado').count()
    )
    can_edit_dash, can_delete_dash = _coord_dash_perms(request.user)
    return render(request, 'conducta/dashboard_coordinador.html', {
        'area': 'bilingue', 'reportes_informativo': qs_info,
        'reportes_conductual': qs_cond, 'reportes_progress': qs_prog,
        'strikes': strikes, 'today': timezone.now().strftime('%Y-%m-%d'),
        'coord_codigo': 'C1',
        'mostrar_reportes_nuevos': True, 'reportes_nuevos_bl': reportes_nuevos_bl,
        'can_edit_dash': can_edit_dash, 'can_delete_dash': can_delete_dash,
        'docentes_filtro': _docentes_filtro(),
        'periodo_activo': _periodo_activo('bilingue'),
        'periodos_escolares': _periodos_escolares('bilingue'),
    })

@login_required
def dashboard_c2(request):
    q = _q_for_coord('C2')
    qs_info = ReporteInformativo.objects.filter(area='bilingue', tipo_reporte='academico').filter(_q_academico_for_coord('C2'))
    qs_cond = ReporteConductual.objects.filter(area='bilingue').filter(q)
    qs_prog = ProgressReport.objects.all()
    qs_info, qs_cond, qs_prog = _dash_querysets(qs_info, qs_cond, qs_prog, 'bilingue')
    strikes = {r['alumno_id']: r['total'] for r in qs_cond.values('alumno_id').annotate(total=Count('id')).filter(total__gte=3)}
    can_edit_dash, can_delete_dash = _coord_dash_perms(request.user)
    return render(request, 'conducta/dashboard_coordinador.html', {
        'area': 'bilingue', 'reportes_informativo': qs_info,
        'reportes_conductual': qs_cond, 'reportes_progress': qs_prog,
        'strikes': strikes, 'today': timezone.now().strftime('%Y-%m-%d'),
        'coord_codigo': 'C2',
        'can_edit_dash': can_edit_dash, 'can_delete_dash': can_delete_dash,
        'docentes_filtro': _docentes_filtro(),
        'periodo_activo': _periodo_activo('bilingue'),
        'periodos_escolares': _periodos_escolares('bilingue'),
    })

@login_required
def dashboard_c3(request):
    q = _q_for_coord('C3')
    qs_info = (
        ReporteInformativo.objects.filter(area='bilingue', tipo_reporte='conductual') |
        ReporteInformativo.objects.filter(area='bilingue', tipo_reporte='academico').filter(q)
    ).distinct().order_by('-fecha')
    qs_cond = ReporteConductual.objects.filter(area='bilingue').filter(q)
    qs_prog = ProgressReport.objects.all()
    qs_info, qs_cond, qs_prog = _dash_querysets(qs_info, qs_cond, qs_prog, 'bilingue')
    strikes = {r['alumno_id']: r['total'] for r in qs_cond.values('alumno_id').annotate(total=Count('id')).filter(total__gte=3)}
    can_edit_dash, can_delete_dash = _coord_dash_perms(request.user)
    return render(request, 'conducta/dashboard_coordinador.html', {
        'area': 'bilingue', 'reportes_informativo': qs_info,
        'reportes_conductual': qs_cond, 'reportes_progress': qs_prog,
        'strikes': strikes, 'today': timezone.now().strftime('%Y-%m-%d'),
        'coord_codigo': 'C3',
        'can_edit_dash': can_edit_dash, 'can_delete_dash': can_delete_dash,
        'docentes_filtro': _docentes_filtro(),
        'periodo_activo': _periodo_activo('bilingue'),
        'periodos_escolares': _periodos_escolares('bilingue'),
    })

@login_required
def dashboard_c4(request):
    q = _q_for_coord('C4')
    qs_info = ReporteInformativo.objects.filter(area='bilingue', tipo_reporte='academico').filter(q)
    qs_cond = ReporteConductual.objects.filter(area='bilingue').filter(q)
    qs_prog = ProgressReport.objects.all()
    qs_info, qs_cond, qs_prog = _dash_querysets(qs_info, qs_cond, qs_prog, 'bilingue')
    strikes = {r['alumno_id']: r['total'] for r in qs_cond.values('alumno_id').annotate(total=Count('id')).filter(total__gte=3)}
    can_edit_dash, can_delete_dash = _coord_dash_perms(request.user)
    return render(request, 'conducta/dashboard_coordinador.html', {
        'area': 'bilingue', 'reportes_informativo': qs_info,
        'reportes_conductual': qs_cond, 'reportes_progress': qs_prog,
        'strikes': strikes, 'today': timezone.now().strftime('%Y-%m-%d'),
        'coord_codigo': 'C4',
        'can_edit_dash': can_edit_dash, 'can_delete_dash': can_delete_dash,
        'docentes_filtro': _docentes_filtro(),
        'periodo_activo': _periodo_activo('bilingue'),
        'periodos_escolares': _periodos_escolares('bilingue'),
    })

@login_required
def dashboard_coordi_bl(request):
    reportes_nuevos_bl = (
        ReporteInformativo.objects.filter(area='bilingue', estado='enviado').count() +
        ReporteConductual.objects.filter(area='bilingue', estado='enviado').count() +
        ProgressReport.objects.filter(estado='enviado').count()
    )
    can_edit_dash, can_delete_dash = _coord_dash_perms(request.user)
    return render(request, 'conducta/dashboard_coordinador.html', {
        'area': 'bilingue',
        'reportes_informativo': ReporteInformativo.objects.none(),
        'reportes_conductual':  ReporteConductual.objects.none(),
        'reportes_progress':    _solo_periodo_activo(ProgressReport.objects.all(), campo='semana_inicio', area='bilingue').annotate(num_evid=Count('evidenciareporte')),
        'strikes': {}, 'today': timezone.now().strftime('%Y-%m-%d'),
        'mostrar_reportes_nuevos': True, 'reportes_nuevos_bl': reportes_nuevos_bl,
        'can_edit_dash': can_edit_dash, 'can_delete_dash': can_delete_dash,
        'docentes_filtro': _docentes_filtro(),
        'periodo_activo': _periodo_activo('bilingue'),
        'periodos_escolares': _periodos_escolares('bilingue'),
    })

@login_required
def periodo_conducta_save(request):
    """Crear/editar un período escolar de conducta (solo superuser)."""
    if not request.user.is_superuser or request.method != 'POST':
        return redirect('menu')
    from .models import PeriodoEscolarConducta
    pk = request.POST.get('pk') or None
    try:
        obj = PeriodoEscolarConducta.objects.get(pk=pk) if pk else PeriodoEscolarConducta()
        obj.nombre       = request.POST.get('nombre', '').strip()[:60]
        obj.parcial      = int(request.POST.get('parcial') or 1)
        obj.anio         = int(request.POST.get('anio') or timezone.now().year)
        obj.area         = request.POST.get('area') if request.POST.get('area') in ('bilingue', 'colegio', 'ambas') else 'ambas'
        obj.fecha_inicio = request.POST.get('fecha_inicio')
        obj.fecha_fin    = request.POST.get('fecha_fin')
        obj.activo       = request.POST.get('activo') == 'on'
        obj.save()
        if obj.activo:
            # Un solo período activo por área
            areas = ['bilingue', 'colegio', 'ambas'] if obj.area == 'ambas' else [obj.area, 'ambas']
            PeriodoEscolarConducta.objects.exclude(pk=obj.pk).filter(area__in=areas).update(activo=False)
        messages.success(request, f'Período "{obj.nombre}" guardado.')
    except Exception as e:
        messages.error(request, f'No se pudo guardar el período: {e}')
    return redirect(request.META.get('HTTP_REFERER', '/'))


@login_required
def periodo_conducta_delete(request, pk):
    if not request.user.is_superuser or request.method != 'POST':
        return redirect('menu')
    from .models import PeriodoEscolarConducta
    PeriodoEscolarConducta.objects.filter(pk=pk).delete()
    messages.success(request, 'Período eliminado.')
    return redirect(request.META.get('HTTP_REFERER', '/'))


@login_required
def historial_alumnado(request, area):
    """Historial de reportes agrupado por parcial → grado (carga perezosa AJAX).
    Regla: el reporte cae en el parcial cuyo rango cubre su fecha; sin match → Parcial 2."""
    if area not in ('bilingue', 'colegio'):
        from django.http import Http404
        raise Http404
    periodos = _periodos_escolares(area)

    def _parcial_de(f):
        for p in periodos:
            if p.fecha_inicio <= f <= p.fecha_fin:
                return p.parcial
        return 2  # acumulado sin período definido → Parcial 2

    reportes = []
    qs_info = ReporteInformativo.objects.filter(area=area).annotate(
        num_evid=Count('evidenciareporte'))
    for r in qs_info:
        reportes.append({'pk': r.pk, 'tipo': 'Informativo / ' + ('Conductual' if r.tipo_reporte == 'conductual' else 'Académico'),
                         'tipo_key': 'informativo', 'color': 'blue',
                         'alumno': r.alumno_nombre, 'alumno_id': r.alumno_id, 'grado': r.grado or '—',
                         'docente': r.docente, 'materia': r.materia, 'fecha': r.fecha, 'estado': r.estado,
                         'num_evid': r.num_evid, 'parcial': _parcial_de(r.fecha.date())})
    qs_cond = ReporteConductual.objects.filter(area=area).annotate(
        num_evid=Count('evidenciareporte'))
    for r in qs_cond:
        reportes.append({'pk': r.pk, 'tipo': 'Conductual', 'tipo_key': 'conductual', 'color': 'orange',
                         'alumno': r.alumno_nombre, 'alumno_id': r.alumno_id, 'grado': r.grado or '—',
                         'docente': r.docente, 'materia': r.materia, 'fecha': r.fecha, 'estado': r.estado,
                         'num_evid': r.num_evid, 'parcial': _parcial_de(r.fecha.date())})
    if area == 'bilingue':
        for r in ProgressReport.objects.annotate(num_evid=Count('evidenciareporte')):
            reportes.append({'pk': r.pk, 'tipo': 'Progress', 'tipo_key': 'progress', 'color': 'teal',
                             'alumno': r.alumno_nombre, 'alumno_id': r.alumno_id, 'grado': r.grado or '—',
                             'docente': '', 'materia': '', 'fecha': r.semana_inicio, 'estado': r.estado,
                             'num_evid': r.num_evid, 'parcial': _parcial_de(r.semana_inicio)})

    # Estructura: parciales 1..3 → grados → reportes
    parciales = []
    for n in (1, 2, 3):
        del_parcial = [r for r in reportes if r['parcial'] == n]
        grados = {}
        for r in del_parcial:
            grados.setdefault(r['grado'], []).append(r)
        grados_list = [
            {'nombre': g, 'reportes': sorted(rs, key=lambda x: str(x['fecha']), reverse=True), 'total': len(rs)}
            for g, rs in sorted(grados.items(), key=lambda kv: kv[0])
        ]
        parciales.append({'num': n, 'total': len(del_parcial), 'grados': grados_list})

    can_edit_dash, can_delete_dash = _coord_dash_perms(request.user)
    return render(request, 'conducta/_historial_alumnado.html', {
        'parciales': parciales, 'area': area, 'total': len(reportes),
        'can_edit_dash': can_edit_dash, 'can_delete_dash': can_delete_dash,
    })


@login_required
def materias_docentes_bl(request):
    """Página independiente: tabla Materia-Docente Bilingüe. Solo superuser."""
    if not request.user.is_superuser:
        from django.core.exceptions import PermissionDenied
        raise PermissionDenied
    from .models import MateriaDocenteBilingue
    q = request.GET.get('q', '').strip()
    materias = MateriaDocenteBilingue.objects.order_by('materia', 'docente')
    if q:
        materias = materias.filter(materia__icontains=q) | materias.filter(docente__icontains=q)
        materias = materias.order_by('materia', 'docente')
    configs = ConfiguracionCoordinador.objects.filter(area='bilingue', activo=True)
    color_map = {'C1': '#c92a2a', 'C2': '#1971c2', 'C3': '#2f9e44', 'C4': '#e67700'}
    # Construir mapa código → nombre para usar en template
    coord_names = {c.codigo: c.nombre for c in configs}
    return render(request, 'conducta/materias_docentes_bl.html', {
        'materias':    materias,
        'configs':     configs,
        'coord_names': coord_names,
        'color_map':   color_map,
        'q':           q,
        'total':       materias.count(),
        'nav_home_url': '/',
    })


# ── CRUD AJAX: MateriaDocenteBilingue ──────────────────────────────
# <--- hecho por claude code: crear, editar y eliminar materias desde la UI sin admin

def _superuser_required(request):
    """Shortcut: devuelve error JSON si no es superuser."""
    if not request.user.is_authenticated or not request.user.is_superuser:
        return JsonResponse({'ok': False, 'error': 'Sin permiso'}, status=403)
    return None


@login_required
@require_POST
def materia_bl_create(request):
    """AJAX POST: crea una nueva MateriaDocenteBilingue."""
    err = _superuser_required(request)
    if err:
        return err
    try:
        body     = json.loads(request.body or b'{}')
        materia  = body.get('materia', '').strip()[:100]
        docente  = body.get('docente', '').strip()[:100]
        coord    = ','.join(c.strip() for c in body.get('coordinadores', []) if c.strip())
        activo   = bool(body.get('activo', True))
        if not materia or not docente:
            return JsonResponse({'ok': False, 'error': 'Materia y docente son obligatorios.'})
    except Exception as e:
        return JsonResponse({'ok': False, 'error': str(e)})

    obj = MateriaDocenteBilingue.objects.create(
        materia=materia, docente=docente,
        coordinador=coord or 'C2', activo=activo
    )
    return JsonResponse({'ok': True, 'pk': obj.pk, 'materia': obj.materia, 'docente': obj.docente,
                         'coordinador': obj.coordinador, 'activo': obj.activo})


@login_required
@require_POST
def materia_bl_update(request, pk):
    """AJAX POST: actualiza un MateriaDocenteBilingue existente."""
    err = _superuser_required(request)
    if err:
        return err
    obj = get_object_or_404(MateriaDocenteBilingue, pk=pk)
    try:
        body    = json.loads(request.body or b'{}')
        materia = body.get('materia', '').strip()[:100]
        docente = body.get('docente', '').strip()[:100]
        coord   = ','.join(c.strip() for c in body.get('coordinadores', []) if c.strip())
        activo  = bool(body.get('activo', True))
        if not materia or not docente:
            return JsonResponse({'ok': False, 'error': 'Materia y docente son obligatorios.'})
    except Exception as e:
        return JsonResponse({'ok': False, 'error': str(e)})

    obj.materia     = materia
    obj.docente     = docente
    obj.coordinador = coord or 'C2'
    obj.activo      = activo
    obj.save()
    return JsonResponse({'ok': True, 'pk': obj.pk, 'materia': obj.materia, 'docente': obj.docente,
                         'coordinador': obj.coordinador, 'activo': obj.activo})


@login_required
@require_POST
def materia_bl_delete(request, pk):
    """AJAX POST: elimina un MateriaDocenteBilingue."""
    err = _superuser_required(request)
    if err:
        return err
    obj = get_object_or_404(MateriaDocenteBilingue, pk=pk)
    obj.delete()
    return JsonResponse({'ok': True})


@login_required
def historial_alumno_coordinador(request, alumno_id):
    """
    Devuelve el historial de reportes (informativo, conductual, progress)
    para un alumno específico. Se usa en el modal del dashboard coordinador.
    """
    informativos = ReporteInformativo.objects.filter(alumno_id=alumno_id).order_by('-fecha')
    conductuales = ReporteConductual.objects.filter(alumno_id=alumno_id).order_by('-fecha')
    progress = ProgressReport.objects.filter(alumno_id=alumno_id).order_by('-fecha')
    context = {
        'informativos': informativos,
        'conductuales': conductuales,
        'progress': progress,
    }
    html = render_to_string('conducta/lista_reportes.html', context)
    return HttpResponse(html)

# ----------- EDITOR DE REPORTES ( SOLO COORDINADOR) -----------
def es_coordinador(user):
    return user.is_staff or user.groups.filter(name__icontains="coordinador").exists()


# Cambia estos según tu lógica o roles
COORDINADORES_BL = [
     "Mr. Martinez",
     "Miss Alcerro",
     "Mr. Ruiz",
    "Mrs. Varela"

]
COORDINADORES_COLEGIO = [
    "Profe. Licona", 
    "Profe. Felipe", 
    "Profe. Gabriela"
    ]


# ─────────────────────────────────────────────────────────────
# subir_evidencia  <--- hecho por claude code
# View que recibe el POST del modal de evidencias del dashboard
# coordinador. Solo coordinadores pueden subir evidencias (chequeo
# con es_coordinador). Según el 'tipo' que llega en el POST,
# asigna la FK correcta del EvidenciaReporte y redirige al
# dashboard del área correspondiente.
# ─────────────────────────────────────────────────────────────
@login_required
def subir_evidencia(request):
    # Solo coordinadores pueden subir evidencias
    if not es_coordinador(request.user):
        return HttpResponseForbidden("Solo coordinadores pueden subir evidencias.")

    if request.method == 'POST':
        tipo = request.POST.get('tipo')
        reporte_id = request.POST.get('reporte_id')
        imagen = request.FILES.get('imagen')
        comentario = request.POST.get('comentario', '')
        # <--- hecho por claude code: se recibe el titulo para identificar la evidencia
        titulo = request.POST.get('titulo', '')

        # Validar que lleguen todos los datos mínimos
        if not imagen or not tipo or not reporte_id:
            messages.error(request, "Faltan datos para subir la evidencia.")
            return redirect(request.META.get('HTTP_REFERER', 'dashboard_maestro'))

        # Asignar la FK correcta según el tipo de reporte
        if tipo == 'conductual':
            reporte = get_object_or_404(ReporteConductual, pk=reporte_id)
            total_evidencias = EvidenciaReporte.objects.filter(reporte_conductual=reporte).count()
            area = reporte.area
        elif tipo == 'informativo':
            reporte = get_object_or_404(ReporteInformativo, pk=reporte_id)
            total_evidencias = EvidenciaReporte.objects.filter(reporte_informativo=reporte).count()
            area = reporte.area
        elif tipo == 'progress':
            reporte = get_object_or_404(ProgressReport, pk=reporte_id)
            total_evidencias = EvidenciaReporte.objects.filter(reporte_progress=reporte).count()
            area = 'bilingue'
        else:
            messages.error(request, "Tipo de reporte inválido.")
            return redirect('dashboard_maestro')

        # Máximo 2 evidencias por reporte
        if total_evidencias >= 2:
            messages.error(request, "Este reporte ya tiene el máximo de 2 evidencias permitidas.")
            return redirect('dashboard_coordinador', area=area)

        # Crear y guardar la evidencia
        ev = EvidenciaReporte(
            tipo=tipo,
            imagen=imagen,
            titulo=titulo,
            comentario=comentario,
            subido_por=request.user
        )

        if tipo == 'conductual':
            ev.reporte_conductual = reporte
        elif tipo == 'informativo':
            ev.reporte_informativo = reporte
        elif tipo == 'progress':
            ev.reporte_progress = reporte

        ev.save()
        messages.success(request, "Evidencia subida correctamente.")
        return redirect('dashboard_coordinador', area=area)

    # Si alguien accede por GET, redirigir al inicio
    return redirect('dashboard_maestro')


@login_required
def eliminar_evidencia(request, pk):
    if not es_coordinador(request.user):
        return JsonResponse({'ok': False, 'error': 'Sin permiso.'}, status=403)
    if request.method != 'POST':
        return JsonResponse({'ok': False, 'error': 'Método no permitido.'}, status=405)
    ev = get_object_or_404(EvidenciaReporte, pk=pk)
    try:
        ev.imagen.delete(save=False)
    except Exception:
        pass
    ev.delete()
    return JsonResponse({'ok': True})


# ================= DIRECTORIO DE TELÉFONOS =================
@login_required
def directorio_telefonos(request):
    user = request.user
    is_admin       = user.is_superuser or user.is_staff
    is_coord_bl    = user.groups.filter(name='coordinador_bilingue').exists()
    is_coord_col   = user.groups.filter(name__in=['coordinadores_colegio', 'coordinador_colegio', 'coordinadores']).exists()
    is_enfermeria  = user.groups.filter(name='enfermeria').exists()

    if is_admin or is_enfermeria:
        areas  = "N'PrimariaBL', N'ColegioBL', N'PreescolarBL', N'Colegio', N'Bachillerato'"
        titulo = 'Todos los Alumnos'
        area_dashboard = 'bilingue'
    elif is_coord_bl:
        areas  = "N'PrimariaBL', N'ColegioBL', N'PreescolarBL'"
        titulo = 'Alumnos Bilingüe'
        area_dashboard = 'bilingue'
    elif is_coord_col:
        areas  = "N'Colegio', N'Bachillerato'"
        titulo = 'Alumnos Colegio'
        area_dashboard = 'colegio'
    else:
        areas  = "N''"
        titulo = ''
        area_dashboard = 'bilingue'

    sql = f"""
      SELECT
        ISNULL(d.Apellido1,'')
          + CASE WHEN d.Apellido2 IS NOT NULL AND d.Apellido2 <> '' THEN ' ' + d.Apellido2 ELSE '' END
          + ', '
          + ISNULL(d.Nombre1,'')
          + CASE WHEN d.Nombre2 IS NOT NULL AND d.Nombre2 <> '' THEN ' ' + d.Nombre2 ELSE '' END AS Nombre,
        da.Descripcion AS Area,
        c.CrsoNumero,
        c.GrupoNumero,
        ISNULL(d.Tel1, '') AS Tel1,
        ISNULL(d.Tel2, '') AS Tel2
      FROM dbo.tblPrsDtosGen AS d
      JOIN dbo.tblPrsTipo           AS t  ON d.PersonaID = t.PersonaID
      JOIN dbo.tblEdcArea           AS a  ON t.IngrEgrID  = a.IngrEgrID
      JOIN dbo.tblEdcEjecCrso       AS ec ON a.AreaID     = ec.AreaID
      JOIN dbo.tblEdcCrso           AS c  ON ec.CrsoID    = c.CrsoID
      JOIN dbo.tblEdcDescripAreaEdc AS da ON a.DescrAreaEdcID = da.DescrAreaEdcID
      WHERE d.Alum = 1
        AND DATEPART(yy, c.FechaInicio) = DATEPART(yyyy, GETDATE())
        AND da.Descripcion IN ({areas})
        AND ec.Desertor    = 0
        AND ec.TrasladoPer = 0
      ORDER BY da.Descripcion, c.CrsoNumero, c.GrupoNumero, Nombre
    """

    import re as _re

    def _moviles(campo1, campo2):
        # Extrae (numero_8_digitos, label_pariente) de un campo de texto.
        # Soporta formatos como "9876-5217 mamá", "8963-1182 papá", "9876/8765".
        def _extraer(texto):
            resultado = []
            patron = r'(?:504[-\s]?)?(\d{4})[-\s./]?(\d{4})\s*([a-zA-ZÀ-ɏ]*)'
            for m in _re.finditer(patron, texto or ''):
                num = m.group(1) + m.group(2)
                label = m.group(3).strip().lower()
                if len(num) == 8 and not num.startswith('2'):
                    resultado.append((num, label))
            return resultado

        vistos = set()
        resultado = []
        for num, label in _extraer(campo1) + _extraer(campo2):
            if num not in vistos:
                vistos.add(num)
                resultado.append({
                    'num':   f'{num[:4]}-{num[4:]}',
                    'wa':    f'https://wa.me/504{num}',
                    'label': label,
                })
        return resultado

    alumnos = []
    try:
        with connections['padres_sqlserver'].cursor() as cursor:
            cursor.execute(sql)
            for nombre, area, crso, grupo, tel1, tel2 in cursor.fetchall():
                moviles = _moviles(tel1.strip(), tel2.strip())
                alumnos.append({
                    'nombre':  nombre.strip(),
                    'grado':   f'{area} {crso}-{grupo}',
                    'moviles': moviles,
                })
    except Exception:
        alumnos = []

    return render(request, 'conducta/directorio_telefonos.html', {
        'alumnos':         alumnos,
        'total':           len(alumnos),
        'titulo':          titulo,
        'area_dashboard':  area_dashboard,
    })



# ════════════════════════════════════════════════════════════════════
# CONVOCATORIA DE TUTORÍAS (Bilingüe)   <--- hecho por claude code
# ════════════════════════════════════════════════════════════════════
import re as _re_tut

ASIGNATURAS_TUTORIA = ['Lang. Arts', 'Lang. Arts/Spelling', 'Math', 'Español',
                       'Language', 'Reading', 'Science', 'Spelling', 'Phonics']


_ORDINAL_GRADO = {
    '1ero': 1, '1mo': 1, 'primero': 1,
    '2do': 2, 'segundo': 2,
    '3ero': 3, '3ro': 3, 'tercero': 3,
    '4to': 4, 'cuarto': 4,
    '5to': 5, 'quinto': 5,
    '6to': 6, 'sexto': 6,
    '7mo': 7, 'septimo': 7, 'séptimo': 7,
    '8vo': 8, 'octavo': 8,
    '9no': 9, 'noveno': 9,
}


def _parse_grado_num(grado_str):
    """'PrimariaBL 1ero-_2' -> (1, '2'). Maneja ordinales (1ero, 2do…) y números."""
    s = grado_str or ''
    m = _re_tut.search(r'([0-9A-Za-zªºáéíóú]+)\s*-\s*([0-9A-Za-z_]+)\s*$', s)
    if not m:
        return None, ''
    crso = m.group(1).strip().lower()
    seccion = m.group(2).lstrip('_').strip()
    n = _ORDINAL_GRADO.get(crso)
    if n is None:
        dm = _re_tut.match(r'(\d+)', crso)
        n = int(dm.group(1)) if dm else None
    if n and 1 <= n <= 9:
        return n, seccion
    return None, seccion


def _horario_grado(grado_num, parcial, anio):
    """Asignaturas de tutoría del grado con los días en que se imparten (1-5).
    Los días vienen del horario y quedan fijos (no los elige el maestro)."""
    from .models import TutoriaHorario
    data, orden = {}, {}
    for r in TutoriaHorario.objects.filter(grado=grado_num, parcial=parcial, anio=anio).order_by('dia'):
        d = data.setdefault(r.asignatura, {'dias': [], 'color': r.color, 'docente': r.docente})
        if r.dia not in d['dias']:
            d['dias'].append(r.dia)
        orden.setdefault(r.asignatura, r.dia)
    out = [{'asignatura': a, 'dias': sorted(v['dias']), 'color': v['color'], 'docente': v['docente']}
           for a, v in data.items()]
    out.sort(key=lambda x: orden.get(x['asignatura'], 99))
    return out


@login_required
def convocatoria_horario_ajax(request):
    """Devuelve las asignaturas+días de tutoría de un grado para un parcial/año."""
    try:
        grado_num = int(request.GET.get('grado_num'))
        parcial = int(request.GET.get('parcial'))
        anio = int(request.GET.get('anio'))
    except (TypeError, ValueError):
        return JsonResponse({'ok': False, 'items': []})
    return JsonResponse({'ok': True, 'items': _horario_grado(grado_num, parcial, anio)})


@login_required
def convocatoria_nueva(request, pk=None):
    """Pantalla del maestro bilingüe: crear o editar convocatoria de tutorías."""
    from .models import (ConvocatoriaTutoria, ConvocatoriaAsignatura,
                         PeriodoEscolarConducta, GRADO_TUTORIA_LABEL)
    area = 'bilingue'
    conv_edit = get_object_or_404(ConvocatoriaTutoria, pk=pk) if pk else None
    students = obtener_alumnos_bilingue()
    periodos = list(PeriodoEscolarConducta.objects.filter(area__in=['bilingue', 'ambas'])
                    .order_by('-anio', '-parcial'))
    activo = _periodo_activo('bilingue')

    if request.method == 'POST':
        alumno_id = (request.POST.get('alumno') or '').strip()
        grado_txt = (request.POST.get('grado') or '').strip()
        try:
            parcial = int(request.POST.get('parcial'))
            anio = int(request.POST.get('anio'))
        except (TypeError, ValueError):
            messages.error(request, "Selecciona un parcial válido.")
            return redirect('convocatoria_nueva')
        # Docente guía y coordinador ya NO se capturan: firman a mano en la carta.
        docente_guia = ''
        coordinador = ''
        fecha_str = (request.POST.get('fecha') or '').strip()
        try:
            fecha = _dt_module.date.fromisoformat(fecha_str) if fecha_str else _dt_module.date.today()
        except ValueError:
            fecha = _dt_module.date.today()

        alumno_obj = next((a for a in students if a['id'] == alumno_id), None)
        alumno_nombre = alumno_obj['label'] if alumno_obj else (request.POST.get('alumno_nombre') or '')
        grado_num, seccion = _parse_grado_num(grado_txt)

        try:
            asigns = json.loads(request.POST.get('asignaturas_json') or '[]')
        except json.JSONDecodeError:
            asigns = []
        asigns = [a for a in asigns if a.get('asignatura') and a.get('dias')]
        redir = 'convocatoria_editar' if conv_edit else 'convocatoria_nueva'
        if not alumno_id or not asigns:
            messages.error(request, "Selecciona el alumno y al menos una asignatura con días.")
            return redirect(redir, pk=conv_edit.pk) if conv_edit else redirect('convocatoria_nueva')

        from .models import grupo_coord_de_grado, ConfiguracionCoordinador
        # grupo solo para ruteo/notificación al coordinador (no aparece en la carta)
        grupo = grupo_coord_de_grado(grado_num)

        if conv_edit:
            conv = conv_edit
            conv.alumno_id, conv.alumno_nombre = alumno_id, alumno_nombre
            conv.grado, conv.grado_num, conv.seccion = grado_txt, grado_num, seccion
            conv.parcial, conv.anio = parcial, anio
            conv.docente_guia, conv.coordinador, conv.coord_grupo = docente_guia, coordinador, grupo
            conv.fecha = fecha
            conv.save()
            conv.asignaturas.all().delete()
        else:
            conv = ConvocatoriaTutoria.objects.create(
                usuario=request.user, area=area, alumno_id=alumno_id, alumno_nombre=alumno_nombre,
                grado=grado_txt, grado_num=grado_num, seccion=seccion,
                parcial=parcial, anio=anio, docente_guia=docente_guia, coordinador=coordinador,
                coord_grupo=grupo, fecha=fecha,
            )
        for a in asigns:
            dias = ','.join(str(int(d)) for d in a['dias'] if str(d).isdigit())
            ConvocatoriaAsignatura.objects.create(
                convocatoria=conv, asignatura=a['asignatura'][:60], dias=dias)

        if not conv_edit and grupo:
            cc = ConfiguracionCoordinador.objects.filter(area='bilingue', codigo=grupo).first()
            if cc and cc.usuario_id:
                try:
                    crear_notificacion(
                        usuario=cc.usuario,
                        mensaje=f"Nueva convocatoria de tutoría: {alumno_nombre} ({grado_txt}) por {request.user.get_full_name() or request.user.username}.",
                        modulo="conducta", tipo="info",
                    )
                except Exception:
                    pass

        if conv_edit:
            messages.success(request, "Convocatoria actualizada.")
            # El coordinador edita → vuelve al listado; el maestro → a su historial
            es_coord = (request.user.is_superuser or
                        ConfiguracionCoordinador.objects.filter(area='bilingue', usuario=request.user).exists())
            return redirect('convocatorias_coordinador') if es_coord else redirect('historial_maestro_bilingue')
        messages.success(request, "¡Convocatoria registrada! Puedes verla en tu Historial.")
        return redirect('historial_maestro_bilingue')

    # GET
    from .models import ConfiguracionCoordinador
    conv_asigs = list(conv_edit.asignaturas.values_list('asignatura', flat=True)) if conv_edit else []
    es_coord = (request.user.is_superuser or
                ConfiguracionCoordinador.objects.filter(area='bilingue', usuario=request.user).exists())
    volver_coord = bool(conv_edit and es_coord)
    return render(request, 'conducta/form_convocatoria.html', {
        'volver_coord': volver_coord,
        'area': area,
        'students': students,
        'periodos': periodos,
        'periodo_activo': activo,
        'anio_actual': (activo.anio if activo else timezone.now().year),
        'parcial_actual': (activo.parcial if activo else 1),
        'fecha_hoy': _dt_module.date.today().isoformat(),
        'grado_labels': GRADO_TUTORIA_LABEL,
        'conv_edit': conv_edit,
        'conv_asigs_json': json.dumps(conv_asigs),
    })


@login_required
@require_POST
def convocatoria_eliminar(request, pk):
    """Elimina una convocatoria (coordinador/superuser o su autor)."""
    from .models import ConvocatoriaTutoria, ConfiguracionCoordinador
    conv = get_object_or_404(ConvocatoriaTutoria, pk=pk)
    es_coord = (request.user.is_superuser or
                ConfiguracionCoordinador.objects.filter(area='bilingue', usuario=request.user).exists())
    if not (es_coord or conv.usuario_id == request.user.id):
        return JsonResponse({'ok': False, 'error': 'Sin permiso'}, status=403)
    conv.delete()
    messages.success(request, "Convocatoria eliminada.")
    return redirect(request.META.get('HTTP_REFERER') or 'convocatorias_coordinador')


def _convocatoria_ctx(conv):
    """Contexto de render (carta) de una convocatoria."""
    from .models import DIA_ABREV, GRADO_TUTORIA_LABEL
    filas = []
    for i, a in enumerate(conv.asignaturas.all(), start=1):
        dl = a.dias_lista
        filas.append({
            'n': i, 'asignatura': a.asignatura,
            'dias_bool': {d: (d in dl) for d in (1, 2, 3, 4, 5)},
        })
    _MESES = ['enero', 'febrero', 'marzo', 'abril', 'mayo', 'junio', 'julio',
              'agosto', 'septiembre', 'octubre', 'noviembre', 'diciembre']
    _PARCIAL_ORD = {1: 'I', 2: 'II', 3: 'III', 4: 'IV'}
    return {
        'conv': conv,
        'filas': filas,
        'grado_label': GRADO_TUTORIA_LABEL.get(conv.grado_num, conv.grado),
        'parcial_rom': _PARCIAL_ORD.get(conv.parcial, conv.parcial),
        'fecha_dia': conv.fecha.day,
        'fecha_mes': _MESES[conv.fecha.month - 1],
        'fecha_anio': conv.fecha.year,
    }


@login_required
def convocatoria_pdf(request, pk):
    """Genera la carta de convocatoria (WeasyPrint)."""
    from weasyprint import HTML as _WHTML
    from .models import ConvocatoriaTutoria
    conv = get_object_or_404(ConvocatoriaTutoria, pk=pk)
    ctx = _convocatoria_ctx(conv)
    # Logo como ruta de archivo (WeasyPrint no siempre resuelve {% static %} por HTTP)
    for _base in (settings.STATIC_ROOT, *[d for d in getattr(settings, 'STATICFILES_DIRS', [])]):
        if not _base:
            continue
        _p = os.path.join(str(_base), 'conducta', 'img', 'encabezado.jpg')
        if os.path.exists(_p):
            ctx['logo_path'] = 'file://' + _p
            break
    html = render_to_string('conducta/pdf/convocatoria_pdf.html', ctx, request=request)
    pdf_bytes = _WHTML(string=html, base_url=request.build_absolute_uri('/')).write_pdf()
    resp = HttpResponse(pdf_bytes, content_type='application/pdf')
    nombre = conv.alumno_nombre.replace(' ', '_')
    resp['Content-Disposition'] = f'inline; filename="convocatoria_{nombre}_P{conv.parcial}.pdf"'
    resp['X-Frame-Options'] = 'SAMEORIGIN'
    return resp


@login_required
def convocatorias_coordinador(request):
    """Solicitudes de tutorías (tab del coordinador) — agrupadas por grado/sección."""
    from .models import (ConvocatoriaMarca, PeriodoEscolarConducta, GRADO_TUTORIA_LABEL,
                         ConfiguracionCoordinador, grupo_coord_de_grado)
    qs = ConvocatoriaMarca.objects.filter(area='bilingue').select_related('marcado_por')

    # Auto-restricción: un coordinador C1/C2 (no superuser) solo ve su grupo
    mi_grupo = ''
    if not request.user.is_superuser:
        cc = ConfiguracionCoordinador.objects.filter(area='bilingue', usuario=request.user,
                                                     codigo__in=['C1', 'C2']).first()
        if cc:
            mi_grupo = cc.codigo

    parcial = request.GET.get('parcial')
    grado = request.GET.get('grado')
    f_grupo = request.GET.get('grupo')
    if parcial:
        qs = qs.filter(parcial=parcial)
    if grado:
        qs = qs.filter(grado_num=grado)

    marcas = list(qs.order_by('grado_num', 'seccion', 'alumno_nombre', 'asignatura'))

    # Agrupar por (grado_num, seccion) -> grupo coord, alumnos
    grupos = {}
    for m in marcas:
        gc = grupo_coord_de_grado(m.grado_num)
        if mi_grupo and gc != mi_grupo:
            continue
        if f_grupo and not mi_grupo and gc != f_grupo:
            continue
        key = (m.grado_num, m.seccion, m.parcial, m.anio)
        g = grupos.get(key)
        if not g:
            lbl = GRADO_TUTORIA_LABEL.get(m.grado_num, m.grado_num)
            g = {
                'grado_num': m.grado_num, 'seccion': m.seccion, 'parcial': m.parcial,
                'anio': m.anio, 'coord_grupo': gc,
                'grado_label': f"{lbl} {m.grado_num}.{m.seccion}" if m.seccion else str(lbl),
                'alumnos': {}, 'n_marcas': 0,
            }
            grupos[key] = g
        a = g['alumnos'].get(m.alumno_id)
        if not a:
            a = {'nombre': m.alumno_nombre, 'asigs': [], 'docentes': set(), 'fecha': m.creado_at}
            g['alumnos'][m.alumno_id] = a
        nombre_doc = (m.marcado_por.get_full_name() or m.marcado_por.username) if m.marcado_por else '—'
        a['asigs'].append({'asignatura': m.asignatura, 'por': nombre_doc})
        a['docentes'].add(nombre_doc)
        g['n_marcas'] += 1

    grupos_list = []
    for g in grupos.values():
        alumnos = list(g['alumnos'].values())
        for a in alumnos:
            a['docentes'] = ', '.join(sorted(a['docentes']))
            a['asigs_txt'] = ', '.join(x['asignatura'] for x in a['asigs'])
        g['alumnos'] = alumnos
        g['n_alumnos'] = len(alumnos)
        grupos_list.append(g)
    grupos_list.sort(key=lambda g: (g['grado_num'], g['seccion']))

    # Firmas (coordinador + docente guía) por grado, con fallback automático a la fija del grado
    for g in grupos_list:
        coord, doc = _firma_convocatoria(g['grado_num'], g['seccion'], g['parcial'], g['anio'])
        g['firma'] = coord
        g['docente_firma'] = doc

    from .models import MateriaDocenteBilingue

    es_coord = request.user.is_staff or request.user.groups.filter(name__icontains='coordinador').exists()
    maestros = sorted(set(MateriaDocenteBilingue.objects.filter(activo=True)
                          .values_list('docente', flat=True)))

    periodos = list(PeriodoEscolarConducta.objects.filter(area__in=['bilingue', 'ambas'])
                    .order_by('-anio', '-parcial'))
    return render(request, 'conducta/convocatorias_coordinador.html', {
        'grupos': grupos_list,
        'periodos': periodos,
        'grados': [(n, GRADO_TUTORIA_LABEL[n]) for n in range(1, 10)],
        'f_parcial': parcial or '',
        'f_grado': grado or '',
        'f_grupo': f_grupo or '',
        'mi_grupo': mi_grupo,
        'es_coord': es_coord,
        'coordinadores': COORDINADORES_BL,
        'maestros': maestros,
    })


@login_required
def tutoria_horario_config(request):
    """Pantalla del coordinador/superuser: editar la matriz Grado × Día × Asignatura."""
    from .models import (TutoriaHorario, PeriodoEscolarConducta, GRADOS_TUTORIA,
                         DIAS_TUTORIA, color_asignatura)
    activo = _periodo_activo('bilingue')
    try:
        parcial = int(request.GET.get('parcial') or (activo.parcial if activo else 1))
        anio = int(request.GET.get('anio') or (activo.anio if activo else timezone.now().year))
    except (TypeError, ValueError):
        parcial = activo.parcial if activo else 1
        anio = activo.anio if activo else timezone.now().year

    if request.method == 'POST':
        try:
            parcial = int(request.POST.get('parcial'))
            anio = int(request.POST.get('anio'))
        except (TypeError, ValueError):
            messages.error(request, "Parcial/año inválido.")
            return redirect('tutoria_horario_config')
        from .models import grupo_coord_de_grado
        TutoriaHorario.objects.filter(parcial=parcial, anio=anio).delete()
        nuevos = []
        for g, _gl in GRADOS_TUTORIA:
            for d, _dl in DIAS_TUTORIA:
                raw = (request.POST.get(f'cell_{g}_{d}') or '').strip()
                vistos = set()
                for asig in [a.strip() for a in raw.split(',') if a.strip()]:
                    key = asig.lower()
                    if key in vistos:
                        continue
                    vistos.add(key)
                    nuevos.append(TutoriaHorario(
                        grado=g, dia=d, asignatura=asig[:60], parcial=parcial, anio=anio,
                        color=color_asignatura(asig)))
        TutoriaHorario.objects.bulk_create(nuevos)
        messages.success(request, f"Horario guardado ({len(nuevos)} asignaturas).")
        return redirect(f"{request.path}?parcial={parcial}&anio={anio}")

    # GET: grilla {grado: {dia: 'A, B'}}
    from .models import grupo_coord_de_grado
    grid = {g: {d: [] for d, _ in DIAS_TUTORIA} for g, _ in GRADOS_TUTORIA}
    for r in TutoriaHorario.objects.filter(parcial=parcial, anio=anio):
        if r.grado in grid and r.dia in grid[r.grado]:
            grid[r.grado][r.dia].append(r.asignatura)
    filas = []
    for g, gl in GRADOS_TUTORIA:
        celdas = [{'dia': d, 'valor': ', '.join(grid[g][d])} for d, _ in DIAS_TUTORIA]
        filas.append({
            'grado': g,
            'label': gl + (" (1 y 2)" if g <= 4 else ""),
            'grupo': grupo_coord_de_grado(g),
            'celdas': celdas,
        })
    periodos = list(PeriodoEscolarConducta.objects.filter(area__in=['bilingue', 'ambas'])
                    .order_by('-anio', '-parcial'))
    return render(request, 'conducta/tutoria_horario_config.html', {
        'filas': filas,
        'dias': DIAS_TUTORIA,
        'parcial': parcial,
        'anio': anio,
        'periodos': periodos,
        'asignaturas_catalogo': ASIGNATURAS_TUTORIA,
    })


# ════════════════════════════════════════════════════════════════════
# CONVOCATORIA — Ventana 1: tabla por grado (alumnos × asignaturas)
# <--- hecho por claude code
# ════════════════════════════════════════════════════════════════════
def _grados_secciones_bilingue(students):
    """Lista ordenada de (grado_num, seccion, label, count) a partir de los alumnos."""
    from .models import GRADO_TUTORIA_LABEL
    agg = {}
    for a in students:
        g, s = _parse_grado_num(a['grado'])
        if not g:
            continue
        agg.setdefault((g, s), 0)
        agg[(g, s)] += 1
    out = []
    for (g, s), n in sorted(agg.items()):
        out.append({'grado_num': g, 'seccion': s,
                    'label': f"{GRADO_TUTORIA_LABEL.get(g, g)} {g}.{s}" if s else f"{GRADO_TUTORIA_LABEL.get(g, g)}",
                    'count': n})
    return out


@login_required
def convocatoria_grado(request):
    """Ventana 1: el docente elige grado-sección y marca alumnos × asignaturas."""
    from .models import (ConvocatoriaMarca, PeriodoEscolarConducta, GRADO_TUTORIA_LABEL,
                         ConfiguracionCoordinador)
    area = 'bilingue'
    activo = _periodo_activo('bilingue')
    periodos = list(PeriodoEscolarConducta.objects.filter(area__in=['bilingue', 'ambas'])
                    .order_by('-anio', '-parcial'))
    try:
        parcial = int(request.GET.get('parcial') or (activo.parcial if activo else 1))
        anio = int(request.GET.get('anio') or (activo.anio if activo else timezone.now().year))
    except (TypeError, ValueError):
        parcial, anio = (activo.parcial if activo else 1), (activo.anio if activo else timezone.now().year)

    students = obtener_alumnos_bilingue()
    grados_secciones = _grados_secciones_bilingue(students)

    sel_g = request.GET.get('g')
    sel_s = request.GET.get('s') or ''
    tabla = None
    if sel_g:
        try:
            gnum = int(sel_g)
        except ValueError:
            gnum = None
        if gnum:
            alumnos = [a for a in students if _parse_grado_num(a['grado']) == (gnum, sel_s)]
            subjects = _horario_grado(gnum, parcial, anio)
            marks = {}
            for mk in ConvocatoriaMarca.objects.filter(
                area=area, parcial=parcial, anio=anio, grado_num=gnum, seccion=sel_s
            ).select_related('marcado_por'):
                nombre = ((mk.marcado_por.get_full_name() or mk.marcado_por.username)
                          if mk.marcado_por_id else '—')
                marks[(mk.alumno_id, mk.asignatura)] = nombre
            filas = []
            for i, a in enumerate(alumnos, start=1):
                cels = []
                for sub in subjects:
                    key = (a['id'], sub['asignatura'])
                    cels.append({'asig': sub['asignatura'], 'color': sub['color'],
                                 'marcado': key in marks, 'por': marks.get(key, '')})
                filas.append({'n': i, 'id': a['id'], 'nombre': a['label'], 'cels': cels})
            tabla = {
                'grado_num': gnum, 'seccion': sel_s,
                'label': f"{GRADO_TUTORIA_LABEL.get(gnum, gnum)} {gnum}.{sel_s}" if sel_s else GRADO_TUTORIA_LABEL.get(gnum, gnum),
                'subjects': subjects, 'filas': filas,
            }

    es_coord = (request.user.is_superuser or
                ConfiguracionCoordinador.objects.filter(area='bilingue', usuario=request.user).exists())
    subjects_json = json.dumps([{'asignatura': s['asignatura'], 'dias': s['dias']}
                               for s in (tabla['subjects'] if tabla else [])])
    return render(request, 'conducta/form_convocatoria_grado.html', {
        'area': area, 'periodos': periodos, 'periodo_activo': activo,
        'parcial': parcial, 'anio': anio,
        'grados_secciones': grados_secciones,
        'sel_g': sel_g or '', 'sel_s': sel_s,
        'tabla': tabla, 'es_coord': es_coord, 'subjects_json': subjects_json,
        'mi_nombre': request.user.get_full_name() or request.user.username,
        'parcial_rom': {1: 'I', 2: 'II', 3: 'III', 4: 'IV'}.get(parcial, parcial),
        'fecha_hoy': _dt_module.date.today().isoformat(),
    })


@login_required
@require_POST
def convocatoria_marca_toggle(request):
    """AJAX: marca/desmarca un alumno en una asignatura (guardado al instante)."""
    from .models import ConvocatoriaMarca
    body = json.loads(request.body or b'{}')
    try:
        parcial = int(body['parcial']); anio = int(body['anio']); grado_num = int(body['grado_num'])
    except (KeyError, ValueError):
        return JsonResponse({'ok': False, 'error': 'Datos inválidos'}, status=400)
    seccion = (body.get('seccion') or '').strip()
    alumno_id = (body.get('alumno_id') or '').strip()
    alumno_nombre = (body.get('alumno_nombre') or '').strip()
    asignatura = (body.get('asignatura') or '').strip()
    checked = bool(body.get('checked'))
    if not alumno_id or not asignatura:
        return JsonResponse({'ok': False, 'error': 'Faltan datos'}, status=400)
    if checked:
        ConvocatoriaMarca.objects.get_or_create(
            area='bilingue', parcial=parcial, anio=anio, alumno_id=alumno_id, asignatura=asignatura,
            defaults={'grado_num': grado_num, 'seccion': seccion,
                      'alumno_nombre': alumno_nombre, 'marcado_por': request.user})
    else:
        ConvocatoriaMarca.objects.filter(
            area='bilingue', parcial=parcial, anio=anio, alumno_id=alumno_id, asignatura=asignatura).delete()
    return JsonResponse({'ok': True})


def _cartas_desde_marcas(grado_num, seccion, parcial, anio):
    """Construye la lista de cartas (una por alumno convocado) desde las marcas."""
    from .models import ConvocatoriaMarca, GRADO_TUTORIA_LABEL
    dias_por_asig = {s['asignatura']: s['dias'] for s in _horario_grado(grado_num, parcial, anio)}
    por_alumno = {}
    for m in ConvocatoriaMarca.objects.filter(
            area='bilingue', parcial=parcial, anio=anio, grado_num=grado_num, seccion=seccion
    ).order_by('alumno_nombre'):
        por_alumno.setdefault((m.alumno_id, m.alumno_nombre), []).append(m.asignatura)
    _MESES = ['enero', 'febrero', 'marzo', 'abril', 'mayo', 'junio', 'julio',
              'agosto', 'septiembre', 'octubre', 'noviembre', 'diciembre']
    cartas = []
    for (aid, nombre), asigs in por_alumno.items():
        filas = []
        for i, asig in enumerate(asigs, start=1):
            dl = dias_por_asig.get(asig, [])
            filas.append({'n': i, 'asignatura': asig, 'dias_bool': {d: (d in dl) for d in (1, 2, 3, 4, 5)}})
        cartas.append({'alumno_nombre': nombre, 'filas': filas})
    return cartas, GRADO_TUTORIA_LABEL.get(grado_num, grado_num)


@login_required
def convocatoria_firma_guardar(request):
    """Guarda el coordinador de área que firma la convocatoria de un grado. Solo coordinadores."""
    if request.method != 'POST':
        return redirect('convocatorias_coordinador')
    if not (request.user.is_staff or request.user.groups.filter(name__icontains='coordinador').exists()):
        messages.error(request, 'No tienes permiso para editar la firma.')
        return redirect('convocatorias_coordinador')
    from .models import ConvocatoriaFirma
    try:
        gnum = int(request.POST['grado']); parcial = int(request.POST['parcial']); anio = int(request.POST['anio'])
    except (KeyError, ValueError):
        messages.error(request, 'Parámetros inválidos.')
        return redirect('convocatorias_coordinador')
    sec = request.POST.get('seccion', '') or ''
    ConvocatoriaFirma.objects.update_or_create(
        area='bilingue', parcial=parcial, anio=anio, grado_num=gnum, seccion=sec,
        defaults={'coordinador_firma': (request.POST.get('coordinador_firma', '') or '').strip(),
                  'docente_firma':     (request.POST.get('docente_firma', '') or '').strip()})
    messages.success(request, 'Firmas de la convocatoria actualizadas.')
    return redirect(request.META.get('HTTP_REFERER') or 'convocatorias_coordinador')


def _firma_convocatoria(gnum, sec, parcial, anio):
    """(coordinador_firma, docente_firma) del grado. Si no hay para este parcial/año,
    usa la firma fija más reciente de ese grado/sección (automático)."""
    from .models import ConvocatoriaFirma
    base = ConvocatoriaFirma.objects.filter(area='bilingue', grado_num=gnum, seccion=sec)
    cf = base.filter(parcial=parcial, anio=anio).first() or base.order_by('-anio', '-parcial').first()
    return (cf.coordinador_firma, cf.docente_firma) if cf else ('', '')


@login_required
def convocatoria_cartas_pdf(request):
    """PDF con las cartas individuales de todos los alumnos convocados del grado."""
    from weasyprint import HTML as _WHTML
    try:
        gnum = int(request.GET.get('g')); parcial = int(request.GET.get('parcial')); anio = int(request.GET.get('anio'))
    except (TypeError, ValueError):
        return HttpResponse('Parámetros inválidos', status=400)
    sec = request.GET.get('s') or ''
    cartas, grado_label = _cartas_desde_marcas(gnum, sec, parcial, anio)
    fecha_str = request.GET.get('fecha') or _dt_module.date.today().isoformat()
    try:
        fecha = _dt_module.date.fromisoformat(fecha_str)
    except ValueError:
        fecha = _dt_module.date.today()
    _MESES = ['enero', 'febrero', 'marzo', 'abril', 'mayo', 'junio', 'julio',
              'agosto', 'septiembre', 'octubre', 'noviembre', 'diciembre']
    ctx = {
        'cartas': cartas, 'grado_label': grado_label, 'seccion': sec,
        'parcial_rom': {1: 'I', 2: 'II', 3: 'III', 4: 'IV'}.get(parcial, parcial),
        'fecha_dia': fecha.day, 'fecha_mes': _MESES[fecha.month - 1], 'fecha_anio': fecha.year,
    }
    ctx['coordinador_firma'], ctx['docente_firma'] = _firma_convocatoria(gnum, sec, parcial, anio)
    for _base in (settings.STATIC_ROOT, *getattr(settings, 'STATICFILES_DIRS', [])):
        p = os.path.join(str(_base), 'conducta', 'img', 'encabezado.jpg') if _base else ''
        if p and os.path.exists(p):
            ctx['logo_path'] = 'file://' + p
            break
    html = render_to_string('conducta/pdf/convocatoria_cartas_pdf.html', ctx, request=request)
    pdf = _WHTML(string=html, base_url=request.build_absolute_uri('/')).write_pdf()
    resp = HttpResponse(pdf, content_type='application/pdf')
    resp['Content-Disposition'] = f'inline; filename="cartas_tutoria_{gnum}-{sec}_P{parcial}.pdf"'
    resp['X-Frame-Options'] = 'SAMEORIGIN'
    return resp


@login_required
def convocatoria_tabla_pdf(request):
    """PDF con la tabla del grado (alumnos × asignaturas con ✔)."""
    from weasyprint import HTML as _WHTML
    from .models import ConvocatoriaMarca, GRADO_TUTORIA_LABEL
    try:
        gnum = int(request.GET.get('g')); parcial = int(request.GET.get('parcial')); anio = int(request.GET.get('anio'))
    except (TypeError, ValueError):
        return HttpResponse('Parámetros inválidos', status=400)
    sec = request.GET.get('s') or ''
    students = obtener_alumnos_bilingue()
    alumnos = [a for a in students if _parse_grado_num(a['grado']) == (gnum, sec)]
    subjects = _horario_grado(gnum, parcial, anio)
    marks = set(ConvocatoriaMarca.objects.filter(
        area='bilingue', parcial=parcial, anio=anio, grado_num=gnum, seccion=sec
    ).values_list('alumno_id', 'asignatura'))
    filas = []
    for i, a in enumerate(alumnos, start=1):
        filas.append({'n': i, 'nombre': a['label'],
                      'cels': [(a['id'], s['asignatura']) in marks for s in subjects]})
    ctx = {
        'grado_label': (f"{GRADO_TUTORIA_LABEL.get(gnum, gnum)} {gnum}.{sec}" if sec else GRADO_TUTORIA_LABEL.get(gnum, gnum)),
        'parcial_rom': {1: 'I', 2: 'II', 3: 'III', 4: 'IV'}.get(parcial, parcial),
        'subjects': [s['asignatura'] for s in subjects], 'filas': filas,
        'coordinador_firma': _firma_convocatoria(gnum, sec, parcial, anio)[0],
    }
    html = render_to_string('conducta/pdf/convocatoria_tabla_pdf.html', ctx, request=request)
    pdf = _WHTML(string=html, base_url=request.build_absolute_uri('/')).write_pdf()
    resp = HttpResponse(pdf, content_type='application/pdf')
    resp['Content-Disposition'] = f'inline; filename="tabla_tutoria_{gnum}-{sec}_P{parcial}.pdf"'
    resp['X-Frame-Options'] = 'SAMEORIGIN'
    return resp


# ═══════════════════════════════════════════════════════════════════
# <--- hecho por claude code: Configuración de ruteo BILINGÜE (solo superusuario)
# Pantalla que administra el JSON (alumnado + mapeo maestro/materia/grado → coordinador).
# ═══════════════════════════════════════════════════════════════════
def _sup_required(request):
    return request.user.is_authenticated and request.user.is_superuser


# <--- hecho por claude code: arma la lista de grupos para el portal (usado por Ruteo y Grupos)
def _grupos_para_portal(cfg):
    alumnos_por_grado = {}
    for a in cfg.get('alumnos', []):
        alumnos_por_grado.setdefault(a.get('grado', ''), []).append(a)
    grupos = []
    for gr, g in (cfg.get('grupos') or {}).items():
        clases_disp = []
        for cl in g.get('clases', []):
            doc = cl.get('docente', '')
            coords = []
            for m in _split_materias(cl.get('materia', '')):
                c = _coord_bl(m, doc, gr, cfg)
                if c and c not in coords:
                    coords.append(c)
            clases_disp.append({
                'docente': doc,
                'materia': _join_materias(cl.get('materia', '')),
                'coords': [{'code': c, 'nombre': _COORD_NOMBRES.get(c, c)} for c in coords],
            })
        grupos.append({
            'key': gr,
            'nombre': _nombre_grupo_disp(g.get('nivel', ''), g.get('nombre') or gr, g.get('seccion', '')),
            'nivel': g.get('nivel', ''),
            'grado_num': g.get('grado_num') or 0,
            'seccion': g.get('seccion', ''),
            'coordinador': g.get('coordinador', ''),
            'clases': clases_disp,
            'alumnos': sorted(alumnos_por_grado.get(gr, []), key=lambda x: x.get('label', '')),
        })
    grupos.sort(key=lambda x: (x['grado_num'], str(x['seccion'])))
    return grupos


@login_required
def routing_bl_config(request):
    puede, solo_lectura = _acceso_ruteo(request)
    if not puede:
        return redirect('menu')
    cfg = _cargar_routing_bl()
    # Maestros/materias conocidos (para overrides, grupos y datalist)
    maestros = sorted({md.docente for md in MateriaDocenteBilingue.objects.filter(activo=True) if md.docente})
    materias = sorted({m.strip() for md in MateriaDocenteBilingue.objects.filter(activo=True)
                       for m in (md.materia or '').split(',') if m.strip()})
    grupos = _grupos_para_portal(cfg)
    return render(request, 'conducta/routing_bl.html', {
        'cfg': cfg,
        'cfg_json': json.dumps(cfg, ensure_ascii=False, indent=2),
        'n_alumnos': len(cfg.get('alumnos', [])),
        'maestros': maestros,
        'materias': materias,
        'coord_nombres': _COORD_NOMBRES,
        'grupos': grupos,
        'n_grupos': len(grupos),
        'actualizado': cfg.get('actualizado', ''),
        'solo_lectura': solo_lectura,   # <--- hecho por claude code: coordinadores = solo lectura
        # <--- hecho por claude code: permisos por coordinador [{code, nombre, correo, permiso}]
        'coord_permisos': [{'code': c, 'nombre': _COORD_NOMBRES.get(c, c),
                            'correo': _COORD_EMAIL.get(c, ''), 'permiso': p}
                           for c, p in _coord_permisos(cfg).items()],
        # <--- hecho por claude code: catálogo de docentes = lista [{docente, materias, coord}]
        'docentes_catalogo': _catalogo_entries(cfg),
        'docentes_catalogo_json': json.dumps(_catalogo_entries(cfg), ensure_ascii=False),
    })


@login_required
@require_POST
def routing_bl_toggle_vis(request):
    """<--- hecho por claude code: fija el permiso de UN coordinador (none|lectura|edit). Solo superuser."""
    if not _sup_required(request):
        return JsonResponse({'ok': False, 'error': 'Sin permiso'}, status=403)
    try:
        body = json.loads(request.body or b'{}')
    except ValueError:
        body = {}
    coord = str(body.get('coord', '')).strip().upper()
    permiso = str(body.get('permiso', '')).strip().lower()
    if coord not in _COORD_EMAIL or permiso not in _PERMISOS_RUTEO:
        return JsonResponse({'ok': False, 'error': 'Datos inválidos'}, status=400)
    cfg = _cargar_routing_bl()
    permisos = _coord_permisos(cfg)   # arranca del estado actual (migrando el legado)
    permisos[coord] = permiso
    cfg['coord_permisos'] = permisos
    cfg.pop('coord_ver_ruteo', None)  # ya migrado al esquema por coordinador
    _guardar_routing_bl(cfg)
    return JsonResponse({'ok': True, 'coord': coord, 'permiso': permiso, 'permisos': permisos})


@login_required
@require_POST
def routing_bl_refrescar(request):
    """Trae el alumnado BL desde SQL Server y lo guarda en el JSON (calcula grado_num)."""
    if not _puede_editar_ruteo(request):
        return JsonResponse({'ok': False, 'error': 'Sin permiso'}, status=403)
    cfg = _cargar_routing_bl()
    alumnos_sql = obtener_alumnos_bilingue()   # golpea SQL Server (único punto)
    alumnos = []
    for a in alumnos_sql:
        n, sec = _parse_grado_num(a['grado'])
        alumnos.append({'id': a['id'], 'label': a['label'], 'grado': a['grado'],
                        'grado_num': n, 'seccion': sec})
    cfg['alumnos'] = alumnos
    _reconstruir_grupos(cfg)   # <--- hecho por claude code: arma/actualiza grupos por grado-sección
    cfg['actualizado'] = timezone.now().strftime('%Y-%m-%d %H:%M')
    _guardar_routing_bl(cfg)
    return JsonResponse({'ok': True, 'n_alumnos': len(alumnos),
                         'n_grupos': len(cfg['grupos']), 'actualizado': cfg['actualizado']})


@login_required
@require_POST
def routing_bl_guardar(request):
    """Guarda el mapeo editado (materias C3/C4, grados C1/C2, overrides). Conserva el alumnado."""
    if not _puede_editar_ruteo(request):
        return JsonResponse({'ok': False, 'error': 'Sin permiso'}, status=403)
    try:
        body = json.loads(request.body or b'{}')
    except ValueError:
        return JsonResponse({'ok': False, 'error': 'JSON inválido'}, status=400)
    cfg = _cargar_routing_bl()

    def _lista_txt(v):
        if isinstance(v, list):
            return [str(x).strip() for x in v if str(x).strip()]
        return [x.strip() for x in str(v or '').split(',') if x.strip()]

    def _lista_int(v, default):
        # <--- hecho por claude code: tolera "1.1" → 1 (toma el número de grado), quita repetidos y sin orden
        items = v if isinstance(v, list) else str(v or '').split(',')
        out = []
        for x in items:
            s = str(x).strip()
            if not s:
                continue
            try:
                n = int(float(s))   # "1.1" → 1 · "5" → 5
            except (TypeError, ValueError):
                continue
            if 1 <= n <= 9 and n not in out:
                out.append(n)
        return out or default

    # <--- hecho por claude code: actualizar SOLO las claves presentes en el body
    # (evita que el autoguardado del mapeo re-agregue overrides de una pestaña vieja)
    if 'materias_c3' in body:
        cfg['materias_c3'] = _lista_txt(body.get('materias_c3'))
    if 'materias_c4' in body:
        cfg['materias_c4'] = _lista_txt(body.get('materias_c4'))
    if 'grados_c1' in body:
        cfg['grados_c1'] = _lista_int(body.get('grados_c1'), [1, 2, 3])
    if 'grados_c2' in body:
        cfg['grados_c2'] = _lista_int(body.get('grados_c2'), [4, 5, 6, 7, 8, 9])
    if 'maestros_override' in body:
        ov = body.get('maestros_override') or {}
        cfg['maestros_override'] = {str(k): _norm_override_list(v) for k, v in ov.items()
                                    if _norm_override_list(v)}
    if 'docentes_catalogo' in body:
        # <--- hecho por claude code: LISTA de entradas [{docente, materias, coord}]
        # un mismo docente puede tener VARIAS entradas (varias cargas / coordinadores).
        # Los overrides (maestros_override) se DERIVAN del catálogo = fuente única.
        raw = body.get('docentes_catalogo') or []
        entradas, overs = [], {}
        if isinstance(raw, list):
            for e in raw:
                doc = str((e or {}).get('docente', '')).strip()
                mats = _join_materias((e or {}).get('materias'))
                coord = str((e or {}).get('coord', '')).strip().upper()
                if coord not in {'C1', 'C2', 'C3', 'C4'}:
                    coord = ''
                if not doc or (not mats and not coord):
                    continue
                entradas.append({'docente': doc, 'materias': mats, 'coord': coord})
                if coord:
                    overs.setdefault(doc, [])
                    if coord not in overs[doc]:
                        overs[doc].append(coord)
        cfg['docentes_catalogo'] = entradas
        cfg['maestros_override'] = overs
    cfg['actualizado'] = timezone.now().strftime('%Y-%m-%d %H:%M')
    _guardar_routing_bl(cfg)
    # <--- hecho por claude code: devolver los valores normalizados para que la UI se sincronice
    return JsonResponse({'ok': True, 'actualizado': cfg['actualizado'], 'norm': {
        'materias_c3': ', '.join(cfg['materias_c3']),
        'materias_c4': ', '.join(cfg['materias_c4']),
        'grados_c1': ', '.join(str(n) for n in cfg['grados_c1']),
        'grados_c2': ', '.join(str(n) for n in cfg['grados_c2']),
    }})


@login_required
def routing_bl_descargar(request):
    """Descarga el JSON completo de ruteo bilingüe."""
    if not _sup_required(request):
        return redirect('menu')
    cfg = _cargar_routing_bl()
    resp = HttpResponse(json.dumps(cfg, ensure_ascii=False, indent=2),
                        content_type='application/json; charset=utf-8')
    resp['Content-Disposition'] = 'attachment; filename="routing_bl.json"'
    return resp


@login_required
@require_POST
def routing_bl_cargar(request):
    """Reemplaza el JSON de ruteo con un archivo subido."""
    if not _sup_required(request):
        return JsonResponse({'ok': False, 'error': 'Sin permiso'}, status=403)
    f = request.FILES.get('archivo')
    if not f:
        return JsonResponse({'ok': False, 'error': 'No se recibió archivo'}, status=400)
    try:
        data = json.loads(f.read().decode('utf-8'))
    except (ValueError, UnicodeDecodeError):
        return JsonResponse({'ok': False, 'error': 'El archivo no es un JSON válido'}, status=400)
    cfg = dict(_ROUTING_BL_DEFAULT)
    cfg.update(data or {})
    cfg['actualizado'] = timezone.now().strftime('%Y-%m-%d %H:%M')
    _guardar_routing_bl(cfg)
    return JsonResponse({'ok': True, 'n_alumnos': len(cfg.get('alumnos', [])), 'actualizado': cfg['actualizado']})


# ═══════════════════════════════════════════════════════════════════
# <--- hecho por claude code: Portal de GRUPOS BILINGÜE (grado-sección → alumnos + docentes/clases + coordinador)
# Los grupos son la fuente de ruteo C1/C2 (materia manda para C3/C4 se resuelve aparte).
# ═══════════════════════════════════════════════════════════════════
@login_required
def grupos_bl_config(request):
    # <--- hecho por claude code: Grupos ahora vive dentro de Ruteo (una sola hoja)
    if not _sup_required(request):
        return redirect('menu')
    return redirect('routing_bl_config')


@login_required
@require_POST
def grupos_bl_guardar(request):
    """Guarda coordinador + clases (docente/materia) de cada grupo. No toca el alumnado."""
    if not _puede_editar_ruteo(request):
        return JsonResponse({'ok': False, 'error': 'Sin permiso'}, status=403)
    try:
        body = json.loads(request.body or b'{}')
    except ValueError:
        return JsonResponse({'ok': False, 'error': 'JSON inválido'}, status=400)
    cfg = _cargar_routing_bl()
    grupos = cfg.get('grupos') or {}
    entrantes = body.get('grupos') or {}
    for key, datos in entrantes.items():
        if key not in grupos:
            continue
        coord = str(datos.get('coordinador', '')).strip().upper()
        if coord in {'C1', 'C2', 'C3', 'C4', ''}:
            grupos[key]['coordinador'] = coord
        clases = []
        for c in (datos.get('clases') or []):
            doc = str(c.get('docente', '')).strip()
            mat = str(c.get('materia', '')).strip()
            if doc or mat:
                clases.append({'docente': doc, 'materia': mat})
        grupos[key]['clases'] = clases
    cfg['grupos'] = grupos
    cfg['actualizado'] = timezone.now().strftime('%Y-%m-%d %H:%M')
    _guardar_routing_bl(cfg)
    # <--- hecho por claude code: devolver el coordinador recalculado por clase (para el autoguardado)
    coords = {}
    for key in entrantes.keys():
        g = grupos.get(key)
        if not g:
            continue
        filas = []
        for cl in g.get('clases', []):
            doc = cl.get('docente', '')
            cs = []
            for m in _split_materias(cl.get('materia', '')):
                c = _coord_bl(m, doc, key, cfg)
                if c and c not in [x['code'] for x in cs]:
                    cs.append({'code': c, 'nombre': _COORD_NOMBRES.get(c, c)})
            filas.append(cs)
        coords[key] = filas
    return JsonResponse({'ok': True, 'actualizado': cfg['actualizado'], 'coords': coords})
