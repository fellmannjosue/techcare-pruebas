import io, os, re
from django.shortcuts import render, redirect, get_object_or_404
from django.http import HttpResponse, HttpResponseForbidden, JsonResponse
from django.contrib.auth.decorators import login_required
from django.contrib import messages
from django.utils import timezone
from django.conf import settings

from django.contrib.auth.models import User
from .models import GradoAgenda, Agenda, ImagenAgenda, AgendaBloqueoConfig


def _agenda_bloqueado(request):
    """Horario semanal automático. Devuelve el mensaje de bloqueo o None.
    Lun: solo permitidos. Mar/Mié: abierto. Jue: hasta la hora límite. Vie: cerrado.
    Coordinadores nunca se bloquean."""
    if _es_coordinador(request.user):
        return None
    cfg = AgendaBloqueoConfig.get()
    if not cfg.activo:
        return None
    from datetime import datetime as _dt
    now = _dt.now()
    wd = now.weekday()  # 0 = Lunes
    if wd == 0:  # Lunes → solo los permitidos pueden llenar
        if not cfg.maestros.filter(pk=request.user.pk).exists():
            return cfg.mensaje
        return None
    if wd == 3:  # Jueves → abierto hasta la hora de cierre
        if now.time() >= cfg.jueves_limite:
            return cfg.mensaje_jueves
        return None
    if wd == 4:  # Viernes → cerrado
        return cfg.mensaje_viernes
    return None  # Martes, Miércoles, fin de semana → abierto


def _agenda_countdown(request):
    """Si hoy es jueves y aún no llega la hora de cierre (y el maestro no es coord),
    devuelve la hora límite ISO para mostrar la cuenta regresiva; None si no aplica."""
    if _es_coordinador(request.user):
        return None
    cfg = AgendaBloqueoConfig.get()
    if not cfg.activo:
        return None
    from datetime import datetime as _dt
    now = _dt.now()
    if now.weekday() == 3 and now.time() < cfg.jueves_limite:
        return _dt.combine(now.date(), cfg.jueves_limite).isoformat()
    return None

# ── Materias por tipo de grado ────────────────────────────────────────────────
MATERIAS_PRIMARIA = [
    "Math", "Phonics", "Reading", "Language",
    "Science", "Español", "CCSS", "Asociadas",
]
MATERIAS_COLEGIO = [
    "Math", "Spelling", "Reading", "Language", "Science",
    "Español", "CCSS", "Cívica", "Asociadas",
]
MATERIAS_COLEGIO_7_9 = [
    "Matemática", "Dibujo Técnico", "Español", "Ciencias Naturales",
    "Cívica", "Estudios Sociales", "Inglés", "Tecnología",
    "Artística", "Computación", "Orientación", "E. Física",
]
MATERIAS_COLEGIO_10 = [
    "Matemática", "Física Elemental", "Robótica", "Español",
    "Química", "Biología", "Sociología", "Psicología",
    "Historia de Honduras", "Inglés Básico", "Inglés Avanzado",
    "Educ. Física", "Informática",
]
MATERIAS_COLEGIO_11 = [
    "Matemática", "Física Elemental", "Dibujo Téc.", "Robótica",
    "Química", "Biología", "Español", "Historia Universal",
    "Economía", "Antropología", "Artística", "Filosofía",
    "Inglés", "Educ. Física",
]
_MATERIAS_MAP = {
    'colegio_bl':  MATERIAS_COLEGIO,
    'colegio_7_9': MATERIAS_COLEGIO_7_9,
    'colegio_10':  MATERIAS_COLEGIO_10,
    'colegio_11':  MATERIAS_COLEGIO_11,
}


def _materias_para_grado(grado_obj):
    return _MATERIAS_MAP.get(grado_obj.tipo_materias, MATERIAS_PRIMARIA)


# ── Clases extra (antes "Asociadas") ──────────────────────────────────────────
# Penmanship: SOLO primaria (1ero–6to). Las demás: primaria + colegio_bl.
CLASE_PENMANSHIP        = "Penmanship"
CLASES_ASOCIADAS_BASE   = ["Arte", "Biblia", "P.E", "Computación", "Speaking", "Spelling"]

# Orden unificado de TODAS las clases por tipo de grado (incluye las extra)
_ORDEN_CLASES = {
    'primaria':   ['Math', 'Language', 'Spelling', 'Phonics', 'Reading', 'Science',
                   'Español', 'CCSS', 'Penmanship', 'Arte', 'Biblia', 'Computación',
                   'Speaking', 'P.E'],
    'colegio_bl': ['Math', 'Language', 'Spelling', 'Reading', 'Science', 'Español',
                   'CCSS', 'Cívica', 'Arte', 'Biblia', 'Computación', 'Speaking', 'P.E'],
}


def _clases_asociadas_para_grado(grado_obj):
    """Clases extra (se guardan solo si tienen contenido)."""
    if not grado_obj or grado_obj.tipo_materias not in ('primaria', 'colegio_bl'):
        return []
    clases = list(CLASES_ASOCIADAS_BASE)
    if grado_obj.tipo_materias == 'primaria':   # 1ero–6to
        clases = [CLASE_PENMANSHIP] + clases
    # Excluir las que ya son materias fijas del grado (p.ej. Spelling en colegio_bl)
    regulares = set(_materias_para_grado(grado_obj))
    return [c for c in clases if c not in regulares]


def _clases_orden(grado_obj):
    """Lista ordenada de TODAS las clases del grado (uniforme, sin separar asociadas)."""
    if grado_obj and grado_obj.tipo_materias in _ORDEN_CLASES:
        return list(_ORDEN_CLASES[grado_obj.tipo_materias])
    # Colegio 7-9/10/11: sin clases extra, se usa la lista fija tal cual
    return [m for m in _materias_para_grado(grado_obj) if m != 'Asociadas'] if grado_obj else []


def _clases_core(grado_obj):
    """Clases que se guardan siempre (aunque vayan vacías)."""
    return {m for m in _materias_para_grado(grado_obj) if m != 'Asociadas'}


def _areas_para_usuario(user, request=None):
    """Devuelve la lista de areas que el usuario puede ver, o None si ve todo."""
    if user.is_superuser:
        return None
    # Cualquier usuario con área seleccionada explícitamente en sesión
    if request and request.session.get("agenda_modo_maestro"):
        area = request.session.get("agenda_area_maestro", "bilingue")
        if area == "colegio":
            return ["colegio"]
        return ["primaria", "colegio_bl"]
    if user.groups.filter(name__in=['coordinador_bilingue', 'coord_progress_bl']).exists():
        return ['primaria', 'colegio_bl']
    if user.groups.filter(name__in=['coordinadores_colegio', 'coordinador_colegio', 'coordinadores']).exists():
        return ['colegio']
    if user.groups.filter(name='maestros_bilingue').exists():
        return ['primaria', 'colegio_bl']
    if user.groups.filter(name='maestros_colegio').exists():
        return ['colegio']
    return None


def _es_coordinador(user):
    return user.is_superuser or user.groups.filter(
        name__in=['coordinador_bilingue', 'coord_progress_bl',
                  'coordinadores_colegio', 'coordinadores']
    ).exists()


# Coordinadores que también son maestros y pueden alternar rol en Agendas
_COORD_MAESTROS = frozenset([
    'cvarela@ana-hn.org',
    'druiz@ana-hn.org',
    'ialcerro@ana-hn.org',
    'jmartinez@ana-hn.org',
])


def _es_coord_maestro(user):
    return user.email.lower() in _COORD_MAESTROS


def _modo_maestro(request):
    """True si el coord-maestro activó el modo maestro en la sesión."""
    return _es_coord_maestro(request.user) and bool(
        request.session.get('agenda_modo_maestro', False)
    )


def _es_coord_efectivo(request):
    """Coordinador efectivo: es coordinador Y NO está en modo maestro."""
    return _es_coordinador(request.user) and not _modo_maestro(request)


def _rol_ctx(request):
    """Contexto de rol que toda vista de agendas debe incluir en render()."""
    return {
        'es_coord':         _es_coord_efectivo(request),
        'es_coord_maestro': _es_coord_maestro(request.user),
        'modo_maestro':     _modo_maestro(request),
    }


# ── TOGGLE MODO MAESTRO / COORDINADOR ─────────────────────────────────────────
@login_required
def toggle_modo_maestro(request):
    if not _es_coord_maestro(request.user):
        return HttpResponseForbidden()
    actual = request.session.get('agenda_modo_maestro', False)
    request.session['agenda_modo_maestro'] = not actual
    if not actual:
        return redirect('agendas:historial_maestro')
    return redirect('agendas:dashboard_coordinador')


# ── FORMULARIO MAESTRO ────────────────────────────────────────────────────────
@login_required
def form_agenda(request):
    if _es_coord_efectivo(request):
        return redirect('agendas:dashboard_coordinador')

    _msg_bloqueo = _agenda_bloqueado(request)
    if _msg_bloqueo:
        return render(request, 'agendas/bloqueado.html',
                      {'mensaje': _msg_bloqueo, **_rol_ctx(request)})

    areas = _areas_para_usuario(request.user, request)
    grados_qs = GradoAgenda.objects.filter(activo=True)
    if areas is not None:
        grados_qs = grados_qs.filter(area__in=areas)
    grados = grados_qs
    usuario_actual = request.user.get_full_name() or request.user.username

    grado_obj  = None
    materias   = []
    semana_ini = ''
    semana_fin = ''

    if request.method == 'POST':
        grado_id   = request.POST.get('grado_id')
        semana_ini = request.POST.get('semana_inicio', '')
        semana_fin = request.POST.get('semana_fin', '')
        grado_obj  = GradoAgenda.objects.filter(pk=grado_id, activo=True).first()

        if not grado_obj:
            messages.error(request, "Selecciona un grado válido.")
            return render(request, 'agendas/form_agenda.html', {'grados': grados, **_rol_ctx(request)})

        materias = _clases_orden(grado_obj)

        # Solo se guarda al presionar "Guardar". El auto-submit por cambio de grado
        # (o fechas faltantes) únicamente recarga la tabla, NO crea la agenda vacía.
        if 'guardar' not in request.POST or not semana_ini or not semana_fin:
            return render(request, 'agendas/form_agenda.html', {
                'grados': grados, 'grado_obj': grado_obj,
                'materias': materias, 'semana_ini': semana_ini, 'semana_fin': semana_fin,
                'clases_asociadas': _clases_asociadas_para_grado(grado_obj),
                **_rol_ctx(request),
            })

        # ── Construir materias_json (lista unificada y ordenada) ──────────────
        core = _clases_core(grado_obj)
        materias_list = []
        for materia in materias:
            lunes     = request.POST.get(f'lunes_{materia}', '')
            martes    = request.POST.get(f'martes_{materia}', '')
            miercoles = request.POST.get(f'miercoles_{materia}', '')
            jueves    = request.POST.get(f'jueves_{materia}', '')
            viernes   = request.POST.get(f'viernes_{materia}', '')
            tiene_contenido = any([lunes, martes, miercoles, jueves, viernes])
            # Clase extra vacía → no se guarda; las core se guardan siempre
            if materia not in core and not tiene_contenido:
                continue
            materias_list.append({
                'materia':   materia,
                'lunes':     lunes,
                'martes':    martes,
                'miercoles': miercoles,
                'jueves':    jueves,
                'viernes':   viernes,
                'nota':      '',
                'docente':   usuario_actual if tiene_contenido else '',
                'es_asociada': materia not in core,
            })

        nota_general = request.POST.get('nota_general', '').strip()[:200]

        # Verificar que no exista ya una agenda para este grado en la misma semana
        from datetime import date as _date
        try:
            ini_d = _date.fromisoformat(str(semana_ini))
            fin_d = _date.fromisoformat(str(semana_fin))
        except ValueError:
            messages.error(request, "Fechas inválidas.")
            return render(request, 'agendas/form_agenda.html', {
                'grados': grados, 'grado_obj': grado_obj,
                'materias': materias, 'semana_ini': semana_ini, 'semana_fin': semana_fin,
                'clases_asociadas': _clases_asociadas_para_grado(grado_obj),
                **_rol_ctx(request),
            })

        # La fecha final no puede ser anterior a la de inicio
        if fin_d < ini_d:
            messages.error(request, "La fecha final no puede ser anterior a la fecha de inicio.")
            return render(request, 'agendas/form_agenda.html', {
                'grados': grados, 'grado_obj': grado_obj,
                'materias': materias, 'semana_ini': semana_ini, 'semana_fin': semana_fin,
                'clases_asociadas': _clases_asociadas_para_grado(grado_obj),
                **_rol_ctx(request),
            })

        # No duplicar agenda del mismo grado: bloquear si las fechas se traslapan
        # O si cae en la misma semana (mismo lunes) que una existente.
        from django.db.models import Q as _Q
        from datetime import timedelta as _td
        lunes_new   = ini_d - _td(days=ini_d.weekday())
        domingo_new = lunes_new + _td(days=6)
        existe = Agenda.objects.filter(grado=grado_obj).filter(
            _Q(semana_inicio__lte=fin_d, semana_fin__gte=ini_d)                 # fechas que se traslapan
            | _Q(semana_inicio__gte=lunes_new, semana_inicio__lte=domingo_new)  # misma semana (mismo lunes)
        ).exists()
        if existe:
            messages.error(
                request,
                f"Ya existe una agenda para {grado_obj.nombre} en esa semana. "
                "No se permite duplicar; búscala en tu Historial para editarla."
            )
            return render(request, 'agendas/form_agenda.html', {
                'grados': grados, 'grado_obj': grado_obj,
                'materias': materias, 'semana_ini': semana_ini, 'semana_fin': semana_fin,
                'clases_asociadas': _clases_asociadas_para_grado(grado_obj),
                **_rol_ctx(request),
            })

        Agenda.objects.create(
            usuario=request.user,
            grado=grado_obj,
            semana_inicio=semana_ini,
            semana_fin=semana_fin,
            materias_json=materias_list,
            nota_general=nota_general,
            creado_por=request.user,
        )
        messages.success(request, "¡Agenda registrada correctamente!")
        return redirect('agendas:form_agenda')

    return render(request, 'agendas/form_agenda.html', {
        'grados':     grados,
        'grado_obj':  grado_obj,
        'materias':   materias,
        'semana_ini': semana_ini,
        'semana_fin': semana_fin,
        'clases_asociadas': _clases_asociadas_para_grado(grado_obj),
        'countdown_to': _agenda_countdown(request),
        **_rol_ctx(request),
    })


# ── EDITAR AGENDA ─────────────────────────────────────────────────────────────
@login_required
def editar_agenda(request, pk):
    _msg_bloqueo = _agenda_bloqueado(request)
    if _msg_bloqueo:
        return render(request, 'agendas/bloqueado.html',
                      {'mensaje': _msg_bloqueo, **_rol_ctx(request)})
    agenda = get_object_or_404(Agenda, pk=pk)
    es_coord       = _es_coord_efectivo(request)
    usuario_actual = request.user.get_full_name() or request.user.username

    # URL de regreso para maestros: su Historial (según el área de la agenda)
    from django.urls import reverse
    _hist_name = 'historial_maestro_colegio' if agenda.grado.area == 'colegio' else 'historial_maestro_bilingue'
    volver_url = reverse(_hist_name)

    DIAS = ['lunes', 'martes', 'miercoles', 'jueves', 'viernes', 'nota']
    core = _clases_core(agenda.grado)

    # Imagen por materia (máx 1)
    imagenes_dict = {img.materia: img for img in agenda.imagenes.all() if img.materia}

    # Datos guardados indexados por nombre de clase
    saved = {m.get('materia', ''): m for m in (agenda.materias_json or [])}
    orden = _clases_orden(agenda.grado)
    # Conservar clases guardadas que no están en el catálogo (legado)
    legado = [n for n in saved if n and n not in orden]
    orden_total = orden + legado

    def _editable(prev):
        dg = (prev.get('docente', '').strip() if prev else '')
        prev_tiene = any(prev.get(d, '').strip() for d in DIAS) if prev else False
        return es_coord or (prev is None) or (dg == usuario_actual) or (not dg and not prev_tiene)

    if request.method == 'POST':
        nuevas_materias = []
        for nombre in orden_total:
            prev = saved.get(nombre)
            if not _editable(prev):
                # Preservar tal cual lo que llenó otro docente
                nuevas_materias.append({
                    'materia': nombre,
                    'lunes': prev.get('lunes', ''), 'martes': prev.get('martes', ''),
                    'miercoles': prev.get('miercoles', ''), 'jueves': prev.get('jueves', ''),
                    'viernes': prev.get('viernes', ''), 'nota': prev.get('nota', ''),
                    'docente': prev.get('docente', ''),
                    'es_asociada': nombre not in core,
                })
                continue
            vals = {d: request.POST.get(f'{d}_{nombre}', '').strip() for d in DIAS}
            tiene = any(vals.values())
            es_extra = nombre not in core
            if es_extra and not tiene:
                continue  # clase extra vacía → no se guarda
            dg = (prev.get('docente', '').strip() if prev else '')
            # Si la celda quedó vacía → se borra también el nombre de quien la llenó.
            if not tiene:
                docente = ''
            elif es_coord:
                docente = dg if dg else usuario_actual
            else:
                docente = usuario_actual
            nuevas_materias.append({
                'materia': nombre, **vals,
                'docente': docente, 'es_asociada': es_extra,
            })

        agenda.materias_json  = nuevas_materias
        agenda.nota_general   = request.POST.get('nota_general', '').strip()[:200]
        agenda.modificado_por = request.user
        agenda.save()
        # <--- hecho por claude code: autoguardado (AJAX) → responde JSON sin recargar.
        # El envío normal del formulario sigue redirigiendo como antes.
        if request.headers.get('X-Requested-With') == 'XMLHttpRequest' or request.POST.get('ajax'):
            return JsonResponse({'ok': True})
        messages.success(request, "Agenda actualizada.")
        if es_coord:
            return redirect('agendas:dashboard_coordinador')
        return redirect(volver_url)

    # ── GET: armar filas ordenadas (uniformes) ────────────────────────────────
    filas = []
    for nombre in orden_total:
        prev = saved.get(nombre)
        filas.append({
            'materia': nombre,
            'lunes': prev.get('lunes', '') if prev else '',
            'martes': prev.get('martes', '') if prev else '',
            'miercoles': prev.get('miercoles', '') if prev else '',
            'jueves': prev.get('jueves', '') if prev else '',
            'viernes': prev.get('viernes', '') if prev else '',
            'nota': prev.get('nota', '') if prev else '',
            'docente': (prev.get('docente', '') if prev else ''),
            'editable': _editable(prev),
            'imagen': imagenes_dict.get(nombre),
        })

    return render(request, 'agendas/editar_agenda.html', {
        'agenda': agenda,
        'filas':  filas,
        'volver_url': volver_url,
        **_rol_ctx(request),
    })


# ── HISTORIAL MAESTRO ─────────────────────────────────────────────────────────
_ORD_PARCIAL = {1: '1er', 2: '2do', 3: '3er', 4: '4to'}


def _agrupar_agendas_por_parcial(agendas_qs):
    """Agrupa las agendas en Parcial (del calendario escolar) → Semana → agendas.
    El parcial se deduce automáticamente por la fecha de inicio de la semana."""
    from conducta.models import PeriodoEscolarConducta
    from datetime import timedelta
    periodos = list(PeriodoEscolarConducta.objects.all().order_by('-anio', '-parcial'))

    def parcial_de(sem):
        for p in periodos:
            if p.fecha_inicio <= sem <= p.fecha_fin:
                return p
        return None

    parciales = {}
    for a in agendas_qs:
        p = parcial_de(a.semana_inicio)
        if p:
            pkey = (p.anio, p.parcial)
            label = f"{_ORD_PARCIAL.get(p.parcial, p.parcial)} Parcial {p.anio}"
            activo = p.activo
        else:
            pkey = (0, 0)
            label = "Sin parcial"
            activo = False
        pg = parciales.setdefault(pkey, {
            'key': f"p{pkey[0]}_{pkey[1]}", 'label': label, 'activo': activo, 'semanas': {}})
        # Agrupar por SEMANA (lunes), aunque la fecha fin difiera entre agendas
        lunes = a.semana_inicio - timedelta(days=a.semana_inicio.weekday())
        pg['semanas'].setdefault(lunes, []).append(a)

    out = []
    for pkey in sorted(parciales, reverse=True):
        pg = parciales[pkey]
        semanas = []
        for lunes in sorted(pg['semanas'], reverse=True):
            semanas.append({
                'inicio': lunes, 'fin': lunes + timedelta(days=4),  # Lun–Vie
                'agendas': pg['semanas'][lunes],
                'id': f"{pg['key']}_{lunes.isoformat().replace('-', '')}",
            })
        out.append({
            'key': pg['key'], 'label': pg['label'], 'activo': pg['activo'],
            'semanas': semanas, 'n_agendas': sum(len(s['agendas']) for s in semanas),
        })
    return out


@login_required
def historial_maestro(request):
    # El listado completo de agendas solo lo ve el coordinador.
    # Un maestro normal crea desde el formulario y revisa las suyas en su Historial.
    if not _es_coordinador(request.user):
        return redirect('agendas:form_agenda')
    if _es_coord_efectivo(request):
        return redirect('agendas:dashboard_coordinador')
    areas = _areas_para_usuario(request.user, request)
    agendas_qs = Agenda.objects.select_related('grado', 'usuario').order_by('-semana_inicio', 'grado__nombre')
    if areas is not None:
        agendas_qs = agendas_qs.filter(grado__area__in=areas)
    return render(request, 'agendas/historial_maestro.html', {
        'agendas': agendas_qs,
        'parciales': _agrupar_agendas_por_parcial(agendas_qs),
        **_rol_ctx(request),
    })


# ── DASHBOARD COORDINADOR ─────────────────────────────────────────────────────
@login_required
def dashboard_coordinador(request):
    if not _es_coordinador(request.user):
        return HttpResponseForbidden()
    areas = _areas_para_usuario(request.user, request)
    agendas_qs = Agenda.objects.select_related('grado', 'usuario').order_by('-semana_inicio', 'grado__nombre')
    if areas is not None:
        agendas_qs = agendas_qs.filter(grado__area__in=areas)

    # Config + maestros bilingües (excluye coordinadores: ellos llenan del lado del coord)
    _COORD_GROUPS = ['coordinador_bilingue', 'coord_progress_bl']
    bcfg = AgendaBloqueoConfig.get()
    bloqueo_ids = set(bcfg.maestros.values_list('id', flat=True))
    maestros_bl = (User.objects.filter(groups__name='maestros_bilingue')
                   .exclude(groups__name__in=_COORD_GROUPS)
                   .order_by('first_name', 'last_name').distinct())

    return render(request, 'agendas/dashboard_coordinador.html', {
        'agendas': agendas_qs,
        'parciales': _agrupar_agendas_por_parcial(agendas_qs),
        'today':   timezone.now().strftime('%Y-%m-%d'),
        'bloqueo_cfg':  bcfg,
        'bloqueo_ids':  bloqueo_ids,
        'maestros_bl':  maestros_bl,
        **_rol_ctx(request),
    })


@login_required
def bloqueo_config_guardar(request):
    """Guarda el horario de bloqueo de agendas (solo coordinador)."""
    if not _es_coordinador(request.user):
        return HttpResponseForbidden()
    if request.method != 'POST':
        return JsonResponse({'error': 'Método no permitido'}, status=405)
    cfg = AgendaBloqueoConfig.get()
    cfg.activo = request.POST.get('activo') in ('1', 'true', 'on', 'True')
    cfg.mensaje         = (request.POST.get('mensaje') or AgendaBloqueoConfig.MSG_LUNES_DEFAULT).strip()
    cfg.mensaje_jueves  = (request.POST.get('mensaje_jueves') or AgendaBloqueoConfig.MSG_JUEVES_DEFAULT).strip()
    cfg.mensaje_viernes = (request.POST.get('mensaje_viernes') or AgendaBloqueoConfig.MSG_VIERNES_DEFAULT).strip()
    hora = (request.POST.get('jueves_limite') or '').strip()
    if hora:
        cfg.jueves_limite = hora
    cfg.save()
    ids = request.POST.getlist('maestros')
    cfg.maestros.set(User.objects.filter(pk__in=ids))
    return JsonResponse({'ok': True})


# ── SUBIR IMAGEN ──────────────────────────────────────────────────────────────
@login_required
def subir_imagen(request):
    if request.method != 'POST':
        return JsonResponse({'error': 'Método no permitido'}, status=405)
    agenda_id   = request.POST.get('agenda_id')
    imagen      = request.FILES.get('imagen')
    descripcion = request.POST.get('descripcion', '')
    materia     = request.POST.get('materia', '')

    if not agenda_id or not imagen:
        return JsonResponse({'error': 'Datos incompletos'}, status=400)

    agenda = get_object_or_404(Agenda, pk=agenda_id)

    # Si ya existe imagen para esta materia, reemplazarla
    es_coord = _es_coordinador(request.user)
    if materia:
        ImagenAgenda.objects.filter(agenda=agenda, materia=materia).delete()

    img = ImagenAgenda.objects.create(
        agenda=agenda,
        imagen=imagen,
        descripcion=descripcion,
        materia=materia,
        subida_por=request.user,
    )
    return JsonResponse({'ok': True, 'url': img.imagen.url, 'id': img.pk})


# ── ELIMINAR IMAGEN ───────────────────────────────────────────────────────────
@login_required
def eliminar_imagen(request, pk):
    img = get_object_or_404(ImagenAgenda, pk=pk)
    es_coord = _es_coordinador(request.user)
    if img.subida_por != request.user and not es_coord:
        return HttpResponseForbidden()
    img.imagen.delete(save=False)
    img.delete()
    return JsonResponse({'ok': True})


# ── DOCX COLEGIO (Word, A4 horizontal) ──────────────────────────────────────
def _descargar_docx_colegio(agenda):
    from docx import Document
    from docx.shared import Cm, Pt, RGBColor as DocxRGB
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.enum.table import WD_TABLE_ALIGNMENT, WD_ALIGN_VERTICAL
    from docx.oxml.ns import qn as dqn
    from docx.oxml import OxmlElement

    _MESES_ES = ['enero','febrero','marzo','abril','mayo','junio',
                 'julio','agosto','septiembre','octubre','noviembre','diciembre']

    def _fecha_es(d):
        return f"{d.day} de {_MESES_ES[d.month - 1]}"

    materias = agenda.materias_json or []

    doc = Document()

    section = doc.sections[0]
    section.orientation = 1       # landscape
    section.page_width   = Cm(29.7)
    section.page_height  = Cm(21.0)
    section.left_margin  = Cm(1.5)
    section.right_margin = Cm(1.5)
    section.top_margin   = Cm(1.5)
    section.bottom_margin = Cm(1.0)

    # Estilo de párrafo base
    style = doc.styles['Normal']
    style.font.name = 'Calibri'
    style.font.size = Pt(11)

    def _hdr_p(text, size=12, bold=False):
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.paragraph_format.space_before = Pt(0)
        p.paragraph_format.space_after  = Pt(2)
        run = p.add_run(text)
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = 'Calibri'
        return p

    _hdr_p('C. E. M. N. G. NUEVO AMANECER', size=14, bold=True)
    _hdr_p('AGENDA SEMANAL DE TAREAS', size=12, bold=True)

    # Grado y semana — tabla sin bordes de 1 fila / 2 celdas
    sem = (f"Semana del {_fecha_es(agenda.semana_inicio)} "
           f"al {_fecha_es(agenda.semana_fin)} de {agenda.semana_fin.year}")

    info = doc.add_table(rows=1, cols=2)
    info.alignment = WD_TABLE_ALIGNMENT.CENTER
    for cell in info.rows[0].cells:
        tcPr = cell._tc.get_or_add_tcPr()
        tcBorders = OxmlElement('w:tcBorders')
        for side in ('top','left','bottom','right','insideH','insideV'):
            el = OxmlElement(f'w:{side}')
            el.set(dqn('w:val'), 'nil')
            tcBorders.append(el)
        tcPr.append(tcBorders)

    c0, c1 = info.rows[0].cells
    r0 = c0.paragraphs[0].add_run(f'Grado: {agenda.grado.nombre}')
    r0.font.bold = True; r0.font.size = Pt(10); r0.font.name = 'Calibri'
    c1.paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.RIGHT
    r1 = c1.paragraphs[0].add_run(sem)
    r1.font.bold = True; r1.font.size = Pt(10); r1.font.name = 'Calibri'

    doc.add_paragraph().paragraph_format.space_after = Pt(2)

    # Tabla de asignaturas
    headers = ['ASIGNATURA', 'LUNES', 'MARTES', 'MIÉRCOLES', 'JUEVES', 'VIERNES']
    tbl = doc.add_table(rows=len(materias) + 1, cols=6)
    tbl.style = 'Table Grid'
    tbl.alignment = WD_TABLE_ALIGNMENT.CENTER

    def _set_width(cell, twips):
        tcPr = cell._tc.get_or_add_tcPr()
        tcW  = OxmlElement('w:tcW')
        tcW.set(dqn('w:w'), str(int(twips)))
        tcW.set(dqn('w:type'), 'dxa')
        tcPr.append(tcW)

    def _shd(cell, fill_hex):
        tcPr = cell._tc.get_or_add_tcPr()
        shd  = OxmlElement('w:shd')
        shd.set(dqn('w:val'),   'clear')
        shd.set(dqn('w:color'), 'auto')
        shd.set(dqn('w:fill'),  fill_hex)
        tcPr.append(shd)

    # 1 cm ≈ 567 twips; A4 landscape usable ≈ 26.7 cm
    W_ASIG = int(4.2 * 567)
    W_DIA  = int((26.7 - 4.2) / 5 * 567)

    for ci, hdr in enumerate(headers):
        cell = tbl.cell(0, ci)
        _set_width(cell, W_ASIG if ci == 0 else W_DIA)
        _shd(cell, 'ADC6E0')
        cell.vertical_alignment = WD_ALIGN_VERTICAL.CENTER
        p = cell.paragraphs[0]
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = p.add_run(hdr)
        run.font.bold = True; run.font.size = Pt(10); run.font.name = 'Calibri'

    for ri, mat in enumerate(materias, start=1):
        vals = [mat.get('materia',''), mat.get('lunes',''), mat.get('martes',''),
                mat.get('miercoles',''), mat.get('jueves',''), mat.get('viernes','')]
        for ci, val in enumerate(vals):
            cell = tbl.cell(ri, ci)
            _set_width(cell, W_ASIG if ci == 0 else W_DIA)
            cell.vertical_alignment = WD_ALIGN_VERTICAL.CENTER
            p = cell.paragraphs[0]
            run = p.add_run(val or '')
            run.font.bold = (ci == 0)
            run.font.size = Pt(9); run.font.name = 'Calibri'

    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    grado_slug = agenda.grado.nombre.replace(' ', '_')
    filename = f'agenda_col_{grado_slug}_{agenda.semana_inicio}.docx'
    response = HttpResponse(
        buf.getvalue(),
        content_type='application/vnd.openxmlformats-officedocument.wordprocessingml.document'
    )
    response['Content-Disposition'] = f'attachment; filename="{filename}"'
    return response


# ── DESCARGA PPTX ─────────────────────────────────────────────────────────────
@login_required
def descargar_pptx_agenda(request, pk):
    if not _es_coord_efectivo(request):
        return HttpResponseForbidden()
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor
    from pptx.oxml.ns import qn
    from lxml import etree

    agenda   = get_object_or_404(Agenda, pk=pk)
    es_coord = _es_coordinador(request.user)

    # ── Colegio: genera DOCX (Word) ──────────────────────────────────────────
    if agenda.grado.area == 'colegio':
        return _descargar_docx_colegio(agenda)


    materias = agenda.materias_json or []

    # Buscar imagen de plantilla subida por el usuario
    import glob as _glob
    from pptx.util import Cm
    _base_img  = '/home/admin2/techcare_project/system_proyect/agendas/static/agendas/img/'
    fondo_path = ''
    for _ext in ('jpg', 'jpeg', 'png'):
        _p = f'{_base_img}plantilla.{_ext}'
        if os.path.exists(_p):
            fondo_path = _p
            break
    if not fondo_path:
        _any = [f for f in _glob.glob(f'{_base_img}*')
                if f.lower().endswith(('.jpg', '.jpeg', '.png'))]
        if _any:
            fondo_path = _any[0]

    # ── Helpers ───────────────────────────────────────────────────────────────
    def _shadow(run):
        rPr = run._r.get_or_add_rPr()
        for ex in rPr.findall(qn('a:effectLst')): rPr.remove(ex)
        ef  = etree.SubElement(rPr, qn('a:effectLst'))
        os_ = etree.SubElement(ef,  qn('a:outerShdw'),
                               attrib={'blurRad':'40000','dist':'23000','dir':'5400000','rotWithShape':'0'})
        cl  = etree.SubElement(os_, qn('a:srgbClr'), val='000000')
        etree.SubElement(cl, qn('a:alpha'), val='40000')

    def _run(para, txt, size, bold=False, italic=False, underline=False,
             shadow=False, color=None, font='Book Antiqua'):
        run = para.add_run()
        run.text = txt
        run.font.name      = font
        run.font.size      = Pt(size)
        run.font.bold      = bold
        run.font.italic    = italic
        run.font.underline = underline
        if color:
            run.font.color.rgb = color
        if shadow:
            _shadow(run)
        return run

    def _borders(cell, w=19050):
        tc   = cell._tc
        tcPr = tc.get_or_add_tcPr()
        for tag in ('lnL','lnR','lnT','lnB'):
            for ex in tcPr.findall(qn(f'a:{tag}')): tcPr.remove(ex)
        for i, tag in enumerate(('lnL','lnR','lnT','lnB')):
            ln = etree.Element(qn(f'a:{tag}'), attrib={'w': str(w), 'cap':'flat','cmpd':'sng'})
            sf = etree.SubElement(ln, qn('a:solidFill'))
            etree.SubElement(sf, qn('a:srgbClr'), val='000000')
            tcPr.insert(i, ln)

    def _add_line(slide, x1, y1, x2, y2, hex_color='3D3620', w_pt=3):
        """Agrega una línea recta al slide vía XML (cxnSp)."""
        sp_tree = slide.shapes._spTree
        ids = [int(e.get('id', 0)) for e in sp_tree.iter() if e.get('id')]
        sp_id = (max(ids) if ids else 100) + 1
        x1, y1, x2, y2 = int(x1), int(y1), int(x2), int(y2)
        off_x  = min(x1, x2);  off_y  = min(y1, y2)
        ext_cx = abs(x2 - x1); ext_cy = abs(y2 - y1)
        fh = 'true' if x2 < x1 else 'false'
        fv = 'true' if y2 < y1 else 'false'
        xml = (
            f'<p:cxnSp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
            f' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            f'<p:nvCxnSpPr><p:cNvPr id="{sp_id}" name="l{sp_id}"/>'
            f'<p:cNvCxnSpPr/><p:nvPr/></p:nvCxnSpPr>'
            f'<p:spPr><a:xfrm flipH="{fh}" flipV="{fv}">'
            f'<a:off x="{off_x}" y="{off_y}"/>'
            f'<a:ext cx="{ext_cx}" cy="{ext_cy}"/></a:xfrm>'
            f'<a:prstGeom prst="line"><a:avLst/></a:prstGeom>'
            f'<a:noFill/>'
            f'<a:ln w="{int(w_pt*12700)}">'
            f'<a:solidFill><a:srgbClr val="{hex_color}"/></a:solidFill>'
            f'</a:ln></p:spPr></p:cxnSp>'
        )
        sp_tree.append(etree.fromstring(xml))

    _MESES_ES = ['enero','febrero','marzo','abril','mayo','junio',
                 'julio','agosto','septiembre','octubre','noviembre','diciembre']

    def _fecha_es(d):
        return f"{d.day} de {_MESES_ES[d.month - 1]}"

    # ── Slide con fondo Espiral (PNG como imagen de fondo) ────────────────────
    LEFT_X  = Inches(0.12)
    LEFT_W  = Inches(2.1)
    TAB_X   = Inches(2.35)
    TAB_W   = Inches(10.8)

    prs = Presentation()
    prs.slide_width  = Emu(12192000)
    prs.slide_height = Emu(6858000)
    slide  = prs.slides.add_slide(prs.slide_layouts[6])
    SLIDE_W = prs.slide_width
    SLIDE_H = prs.slide_height

    if fondo_path:
        slide.shapes.add_picture(fondo_path, 0, 0, width=SLIDE_W, height=SLIDE_H)
    else:
        bg = slide.background.fill; bg.solid()
        bg.fore_color.rgb = RGBColor(245, 245, 228)
        pan = slide.shapes.add_shape(1, 0, 0, TAB_X, SLIDE_H)
        pan.fill.solid(); pan.fill.fore_color.rgb = RGBColor(50, 44, 24)
        pan.line.fill.background()
        for _i in range(-3, 20):
            _y0 = Inches(_i * 0.52)
            _add_line(slide, x1=0, y1=_y0, x2=TAB_X, y2=_y0 + Inches(2.4),
                      hex_color='3A3318', w_pt=3)
        acc = slide.shapes.add_shape(1, TAB_X - Inches(0.15), 0, Inches(0.15), SLIDE_H)
        acc.fill.solid(); acc.fill.fore_color.rgb = RGBColor(130, 100, 45)
        acc.line.fill.background()

    # ── Texto panel izquierdo ─────────────────────────────────────────────────
    CLR_DARK  = RGBColor(30, 30, 30)    # negro para título y semana
    CLR_WHITE = RGBColor(255, 255, 255) # blanco para texto dentro de la flecha

    tb = slide.shapes.add_textbox(LEFT_X, Inches(0.1), LEFT_W, Inches(2.2))
    tb.fill.background(); tb.line.fill.background()
    tf = tb.text_frame; tf.clear(); tf.word_wrap = True
    p  = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    # <--- hecho por claude code: centrado, Book Antiqua, tamaño 20
    _run(p, 'Agenda\nSemanal', 20, bold=True, italic=True, shadow=True, color=CLR_DARK, font='Book Antiqua')

    sem = f"{_fecha_es(agenda.semana_inicio)} al {_fecha_es(agenda.semana_fin)}"
    tb = slide.shapes.add_textbox(LEFT_X, Inches(2.5), LEFT_W, Inches(2.0))
    tb.fill.background(); tb.line.fill.background()
    tf = tb.text_frame; tf.clear(); tf.word_wrap = True
    # <--- hecho por claude code: Semana centrado, Book Antiqua, tamaño 20
    p  = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    _run(p, 'Semana:', 20, bold=True, italic=True, underline=True, shadow=True, color=CLR_DARK, font='Book Antiqua')
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    _run(p2, sem, 20, italic=True, shadow=True, color=CLR_DARK, font='Book Antiqua')

    # Textbox transparente sobre la flecha roja de la plantilla
    # <--- hecho por claude code: posición/tamaño exactos de la flecha (medidos en PowerPoint)
    tb = slide.shapes.add_textbox(Cm(-0.7), Cm(12.11), Cm(5.72), Cm(1.91))
    tb.fill.background(); tb.line.fill.background()
    tf = tb.text_frame; tf.clear(); tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE   # centrado vertical dentro de la flecha
    for _m in ('margin_left', 'margin_right', 'margin_top', 'margin_bottom'):
        setattr(tf, _m, 0)
    p  = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    _run(p, f'Grado: {agenda.grado.nombre}', 20, bold=True, italic=True, color=CLR_WHITE, font='Book Antiqua')

    # ── Tabla ─────────────────────────────────────────────────────────────────
    num_rows = len(materias) + 1
    tbl_shape = slide.shapes.add_table(num_rows, 7, TAB_X, Inches(0.05), TAB_W, SLIDE_H - Inches(0.1))
    tbl = tbl_shape.table

    # Sin bandas de estilo
    tblPr = tbl._tbl.find(qn('a:tblPr'))
    if tblPr is None:
        tblPr = etree.SubElement(tbl._tbl, qn('a:tblPr'))
    for _at in ('bandRow','bandCol','firstRow','firstCol','lastRow','lastCol'):
        tblPr.set(_at, '0')
    for _el in tblPr.findall(qn('a:tableStyleId')): tblPr.remove(_el)
    etree.SubElement(tblPr, qn('a:tableStyleId')).text = '{2D5ABB26-0587-4C30-8999-92F81FD0307C}'

    # Anchos columnas → total 10.8"
    for _i, _w in enumerate([Inches(1.3),
                              Inches(1.55), Inches(1.55), Inches(1.55),
                              Inches(1.55), Inches(1.55), Inches(1.75)]):
        tbl.columns[_i].width = _w

    def _cell(cell, txt, bold=False, align=PP_ALIGN.CENTER, header=False, size=9):
        # Limpiar retornos de carro Windows (\r) que generan "_x000D_" en PPTX
        txt = (txt or '').replace('\r\n', '\n').replace('\r', '')
        cell.text = txt
        _borders(cell)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(235, 238, 215) if header else RGBColor(255, 255, 255)
        tf   = cell.text_frame
        tf.word_wrap = True
        # <--- hecho por claude code: aplicar Century + tamaño a TODAS las líneas y runs
        # (antes solo la 1ª línea → el resto quedaba en 18pt y se veía disparejo)
        for para in tf.paragraphs:
            para.alignment      = align
            para.font.name      = 'Century Gothic'
            para.font.size      = Pt(size)
            para.font.bold      = bold
            para.font.color.rgb = RGBColor(0, 0, 0)
            for run in para.runs:
                run.font.name      = 'Century Gothic'
                run.font.size      = Pt(size)
                run.font.bold      = bold
                run.font.color.rgb = RGBColor(0, 0, 0)
        # Encoger texto para que quepa en la altura fija de 1.8 cm
        txBody = cell._tc.find(qn('a:txBody'))
        if txBody is not None:
            bodyPr = txBody.find(qn('a:bodyPr'))
            if bodyPr is not None:
                for _t in (qn('a:noAutofit'), qn('a:normAutofit'), qn('a:spAutoFit')):
                    for _e in bodyPr.findall(_t): bodyPr.remove(_e)
                bodyPr.set('anchor', 'ctr')     # centrado vertical
                bodyPr.set('anchorCtr', '1')
                # <--- hecho por claude code: sin encoger — el texto queda al tamaño fijo (uniforme)
                etree.SubElement(bodyPr, qn('a:noAutofit'))

    headers = ['Clase/Día', 'Lunes', 'Martes', 'Miércoles', 'Jueves', 'Viernes', 'Nota']
    for _i, _h in enumerate(headers):
        _cell(tbl.cell(0, _i), _h, bold=True, header=True, size=16)

    for _row, mat in enumerate(materias, start=1):
        _cell(tbl.cell(_row, 0), mat.get('materia',   ''), bold=True, size=16)  # <--- hecho por claude code: columna Clase/Día = 16
        _cell(tbl.cell(_row, 1), mat.get('lunes',     ''))
        _cell(tbl.cell(_row, 2), mat.get('martes',    ''))
        _cell(tbl.cell(_row, 3), mat.get('miercoles', ''))
        _cell(tbl.cell(_row, 4), mat.get('jueves',    ''))
        _cell(tbl.cell(_row, 5), mat.get('viernes',   ''))

    # Columna Nota = una sola celda combinada (nota general de la semana)
    if materias:
        if len(materias) > 1:
            tbl.cell(1, 6).merge(tbl.cell(len(materias), 6))
        _cell(tbl.cell(1, 6), agenda.nota_general or '')

    # Altura fija 1.8 cm por fila (el texto se encoge si no cabe)
    for _i in range(num_rows):
        tbl.rows[_i].height = Cm(1.8)

    # ── Descarga ──────────────────────────────────────────────────────────────
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    grado_slug = agenda.grado.nombre.replace(' ', '_')
    filename = f'agenda_{grado_slug}_{agenda.semana_inicio}.pptx'
    response = HttpResponse(buf.getvalue(),
                            content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
    response['Content-Disposition'] = f'attachment; filename="{filename}"'
    return response


# ── ELIMINAR AGENDA ───────────────────────────────────────────────────────────
@login_required
def eliminar_agenda(request, pk):
    if not _es_coordinador(request.user):
        return JsonResponse({'error': 'Sin permiso'}, status=403)
    if request.method != 'POST':
        return JsonResponse({'error': 'Método no permitido'}, status=405)
    agenda = get_object_or_404(Agenda, pk=pk)
    agenda.delete()
    return JsonResponse({'ok': True})
